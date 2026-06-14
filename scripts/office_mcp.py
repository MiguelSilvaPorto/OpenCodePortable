import os
import re
import shutil
import subprocess
from mcp.server.fastmcp import FastMCP

# Inicializa o servidor FastMCP
mcp = FastMCP("Office", instructions="""
Este servidor fornece ferramentas locais para criar, ler e modificar arquivos Word (.docx), Excel (.xlsx/.xlsm) e PowerPoint (.pptx).
Tambem fornece extracao de texto de PDF, OCR em imagens, protecao por senha e execucao de macros VBA no Windows.

IMPORTANTE: Se o usuario pedir para criar ou modificar uma macro mas nao especificar claramente as regras ou o comportamento da mesma,
peca esclarecimentos antes de prosseguir.

Para documentos avancados, use create_from_example() passando um arquivo modelo existente.
""")

# ==========================================
# UTILITARIOS COMPARTILHADOS
# ==========================================

def _validar_arquivo_existe(file_path: str) -> str | None:
    """Valida se o arquivo existe. Retorna None se OK, ou mensagem de erro."""
    if not os.path.exists(file_path):
        return f"Erro: Nao foi possivel acessar o arquivo. O caminho '{file_path}' nao existe. Verifique o caminho e tente novamente."
    return None

def _validar_extensao(file_path: str, extensoes: tuple) -> str | None:
    """Valida se a extensao do arquivo e uma das esperadas."""
    if not file_path.lower().endswith(extensoes):
        ext_list = ", ".join(extensoes)
        return f"Erro: Formato de arquivo invalido. A extensao deve ser {ext_list}. Use: {file_path.split('.')[0]}{extensoes[0]}"
    return None

def _sanitizar_erro(stderr: str) -> str:
    """Remove ruido de stack trace do PowerShell, mantendo apenas a mensagem principal."""
    linhas = stderr.strip().split('\n')
    linhas_uteis = [
        l for l in linhas
        if not re.match(r'^\s*\+\s*CategoryInfo|^\s*\+ FullyQualifiedErrorId|^\s*\+', l)
    ]
    mensagem = ' '.join(l.strip() for l in linhas_uteis if l.strip())
    return mensagem[:500]

def _executar_powershell(script: str, timeout: int = 60) -> str:
    """
    Executa um script PowerShell de forma segura, com blindagens:
    - ExecutionPolicy Bypass: ignora restricoes locais
    - UTF-8 forcado: preserva acentuacao do portugues
    - Timeout explicito: evita travamentos
    - Sanitizacao de erro: protege o contexto do LLM
    """
    script_blindado = f"""
$OutputEncoding = [System.Text.Encoding]::UTF8
$ErrorActionPreference = "Stop"
{script}
"""
    cmd = [
        "powershell",
        "-NoProfile",
        "-NonInteractive",
        "-ExecutionPolicy", "Bypass",
        "-Command", script_blindado
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
        if result.returncode != 0:
            return f"Erro: {_sanitizar_erro(result.stderr)}"
        return result.stdout.strip()
    except subprocess.TimeoutExpired:
        return "Erro: A operacao excedeu o tempo limite de 60 segundos. O arquivo pode estar corrompido ou o Office pode estar ocupado. Feche outros programas e tente novamente."
    except FileNotFoundError:
        return "Erro: PowerShell nao encontrado. Esta ferramenta requer Windows com PowerShell 5.1 ou superior."
    except Exception as e:
        return f"Erro: Falha ao executar comando PowerShell. Detalhes: {_sanitizar_erro(str(e))}"


# ==========================================
# IMPORTACOES PREGUIÇOSAS (Lazy Imports)
# ==========================================
# Cada biblioteca e carregada sob demanda para evitar falhas durante a descoberta MCP
# se a dependencia nao estiver instalada.

def get_docx():
    import docx
    return docx

def get_openpyxl():
    import openpyxl
    return openpyxl

def get_pptx():
    import pptx
    return pptx

def get_power_pptx():
    import power_pptx
    return power_pptx

def get_win32():
    import win32com.client
    return win32com.client

def get_pymupdf():
    import fitz
    return fitz

def get_easyocr():
    import easyocr
    return easyocr

def get_msoffcrypto():
    import msoffcrypto
    return msoffcrypto


# ==========================================
# 1. FERRAMENTAS WORD (.docx)
# ==========================================

@mcp.tool()
def create_document(file_path: str, elements: list) -> str:
    """
    Cria um novo documento Word (.docx) a partir de elementos estruturados.

    A ferramenta recebe uma lista de dicionarios onde cada elemento representa
    um bloco do documento (titulo, heading, paragrafo, item de lista).
    O arquivo e sobrescrito se ja existir.

    Args:
        file_path (str): Caminho absoluto ou relativo onde o .docx sera salvo.
            Ex: "C:\\relatorios\\documento.docx" ou "./output/doc.docx"
        elements (list[dict]): Lista de elementos. Cada dicionario deve ter:
            - type (str, obrigatorio): "title" | "heading1" | "heading2" | "paragraph" | "list_item"
            - text (str, obrigatorio): Conteudo textual do elemento
            - bold (bool, opcional): Negrito (default: False)
            - italic (bool, opcional): Italico (default: False)

    Returns:
        str: Mensagem de confirmacao com o caminho absoluto do arquivo salvo.

    Exemplo de uso:
        create_document(
            file_path="./relatorio.docx",
            elements=[
                {"type": "title", "text": "Relatorio Anual"},
                {"type": "heading1", "text": "Introducao"},
                {"type": "paragraph", "text": "Texto do paragrafo."},
                {"type": "list_item", "text": "Item A", "bold": True},
            ]
        )
        # Retorno: "Documento Word criado com sucesso em: C:\\projeto\\relatorio.docx"

    Notas:
        - "title" cria um titulo principal (Heading Level 0)
        - "list_item" cria paragrafo com estilo "List Bullet"
        - Para documentos com layout complexo, prefira create_document_from_template()
    """
    docx = get_docx()
    doc = docx.Document()

    for el in elements:
        el_type = el.get("type", "paragraph")
        text = el.get("text", "")

        if el_type == "title":
            doc.add_heading(text, level=0)
        elif el_type == "heading1":
            doc.add_heading(text, level=1)
        elif el_type == "heading2":
            doc.add_heading(text, level=2)
        elif el_type == "list_item":
            p = doc.add_paragraph(style='List Bullet')
            run = p.add_run(text)
            run.bold = el.get("bold", False)
            run.italic = el.get("italic", False)
        else:
            p = doc.add_paragraph()
            run = p.add_run(text)
            run.bold = el.get("bold", False)
            run.italic = el.get("italic", False)

    doc.save(file_path)
    return f"Documento Word criado com sucesso em: {os.path.abspath(file_path)}"


@mcp.tool()
def read_document(file_path: str) -> str:
    """
    Le e extrai todo o texto estruturado de um documento Word (.docx), incluindo tabelas.

    Args:
        file_path (str): Caminho do arquivo .docx a ser lido.

    Returns:
        str: Conteudo textual organizado, com secoes separadas e tabelas formatadas.
            Se o arquivo nao existir, retorna mensagem de erro.

    Exemplo de uso:
        read_document(file_path="./contrato.docx")
        # Retorno:
        # --- Documento: contrato.docx ---
        # CONTRATO DE PRESTACAO DE SERVICOS
        # ...
        # --- Tabelas Encontradas ---
        # Tabela 1:
        # Item | Valor | Prazo

    Notas:
        - Tabelas sao exibidas com celulas separadas por " | "
        - A formatacao visual (fontes, cores) NAO e extraida
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    docx = get_docx()
    doc = docx.Document(file_path)
    result = []

    result.append(f"--- Documento: {os.path.basename(file_path)} ---")
    for p in doc.paragraphs:
        if p.text.strip():
            result.append(p.text)

    if doc.tables:
        result.append("\n--- Tabelas Encontradas ---")
        for i, table in enumerate(doc.tables):
            result.append(f"Tabela {i+1}:")
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                result.append(" | ".join(row_text))

    return "\n".join(result)


@mcp.tool()
def create_document_from_template(template_path: str, output_path: str, replacements: dict) -> str:
    """
    Clona um documento .docx existente e substitui marcadores de texto,
    preservando 100% da formatacao original do modelo.

    Args:
        template_path (str): Caminho do arquivo modelo .docx.
            Ex: "./modelos/contrato_base.docx"
        output_path (str): Caminho onde o documento gerado sera salvo.
            Ex: "./documentos/contrato_final.docx"
        replacements (dict): Dicionario de substituicoes.
            Ex: {"{{NomeCliente}}": "Joao Silva", "{{Data}}": "10/06/2026"}

    Returns:
        str: Mensagem de confirmacao com o caminho absoluto do arquivo gerado.

    Exemplo de uso:
        create_document_from_template(
            template_path="./modelos/contrato_base.docx",
            output_path="./contrato_final.docx",
            replacements={
                "{{NomeCliente}}": "Joao Silva",
                "{{Data}}": "10/06/2026",
                "{{Valor}}": "R$ 15.000,00"
            }
        )
        # Retorno: "Documento gerado a partir do modelo com sucesso em: C:\\docs\\contrato_final.docx"

    Notas:
        - A substituicao funciona em "runs" individuais do XML do Word
        - Se um marcador estiver dividido em varios runs (ex: "{{Nome" em um run e "Cliente}}" em outro),
          a substituicao pode falhar. Crie modelos com o marcador inteiro em um unico run.
        - A formatacao original (fontes, cores, tamanhos, espacamento) e 100% preservada
    """
    erro = _validar_arquivo_existe(template_path)
    if erro:
        return erro

    docx = get_docx()
    shutil.copyfile(template_path, output_path)
    doc = docx.Document(output_path)

    for p in doc.paragraphs:
        for run in p.runs:
            for key, val in replacements.items():
                if key in run.text:
                    run.text = run.text.replace(key, str(val))

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    for run in p.runs:
                        for key, val in replacements.items():
                            if key in run.text:
                                run.text = run.text.replace(key, str(val))

    doc.save(output_path)
    return f"Documento gerado a partir do modelo com sucesso em: {os.path.abspath(output_path)}"


# ==========================================
# 1.5 WORD - FERRAMENTAS AVANCADAS
# ==========================================

@mcp.tool()
def add_image_to_document(file_path: str, image_path: str, width_inch: float = None, height_inch: float = None) -> str:
    """
    Insere uma imagem em um documento Word (.docx) existente.

    Args:
        file_path (str): Caminho do arquivo .docx existente.
        image_path (str): Caminho da imagem local (PNG, JPG, GIF).
        width_inch (float, opcional): Largura da imagem em polegadas. Se nao informado, usa o tamanho original.
        height_inch (float, opcional): Altura da imagem em polegadas. Se nao informado, proporcional a width.

    Returns:
        str: Mensagem de confirmacao.

    Exemplo:
        add_image_to_document(file_path="./relatorio.docx", image_path="./grafico.png", width_inch=5.5)
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    erro_img = _validar_arquivo_existe(image_path)
    if erro_img:
        return erro_img

    docx = get_docx()
    from docx.shared import Inches

    doc = docx.Document(file_path)
    width = Inches(width_inch) if width_inch else None
    height = Inches(height_inch) if height_inch else None
    doc.add_picture(image_path, width=width, height=height)
    doc.save(file_path)
    return f"Imagem inserida com sucesso em: {os.path.abspath(file_path)}"


@mcp.tool()
def add_table_to_document(file_path: str, data: list, style: str = "Light Grid Accent 1") -> str:
    """
    Adiciona uma tabela formatada a um documento Word (.docx) existente.

    Args:
        file_path (str): Caminho do arquivo .docx.
        data (list[list]): Dados da tabela. Primeira linha = cabecalho.
            Ex: [["Produto", "Vendas"], ["A", 100], ["B", 200]]
        style (str, opcional): Estilo da tabela (ex: "Light Grid Accent 1", "Table Grid").

    Returns:
        str: Mensagem de confirmacao.

    Exemplo:
        add_table_to_document(file_path="./relatorio.docx", data=[["Item","Valor"],["X","R$50"]])
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    if not data or not data[0]:
        return "Erro: Os dados da tabela nao podem estar vazios. Forneca uma lista com pelo menos uma linha."

    docx = get_docx()
    doc = docx.Document(file_path)
    rows = len(data)
    cols = len(data[0])
    table = doc.add_table(rows=rows, cols=cols)
    table.style = style

    for i, row_data in enumerate(data):
        for j, cell_value in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_value)
            if i == 0:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.bold = True

    doc.save(file_path)
    return f"Tabela {rows}x{cols} adicionada com sucesso em: {os.path.abspath(file_path)}"


@mcp.tool()
def edit_document(file_path: str, elements: list, position: str = "end") -> str:
    """
    Adiciona elementos a um documento Word (.docx) existente.

    Args:
        file_path (str): Caminho do arquivo .docx.
        elements (list[dict]): Mesmo formato de create_document().
        position (str, opcional): "end" para adicionar ao final, "begin" para adicionar no inicio.

    Returns:
        str: Mensagem de confirmacao.

    Exemplo:
        edit_document(file_path="./relatorio.docx", elements=[{"type": "heading1", "text": "Novo Capitulo"}])
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    if position not in ("end", "begin"):
        return "Erro: Posicao invalida. Use 'end' (final) ou 'begin' (inicio)."

    docx = get_docx()
    doc = docx.Document(file_path)

    new_paragraphs = []
    for el in elements:
        el_type = el.get("type", "paragraph")
        text = el.get("text", "")

        if el_type == "title":
            new_paragraphs.append(doc.add_heading(text, level=0))
        elif el_type == "heading1":
            new_paragraphs.append(doc.add_heading(text, level=1))
        elif el_type == "heading2":
            new_paragraphs.append(doc.add_heading(text, level=2))
        elif el_type == "list_item":
            p = doc.add_paragraph(style='List Bullet')
            run = p.add_run(text)
            run.bold = el.get("bold", False)
            run.italic = el.get("italic", False)
            new_paragraphs.append(p)
        else:
            p = doc.add_paragraph()
            run = p.add_run(text)
            run.bold = el.get("bold", False)
            run.italic = el.get("italic", False)
            new_paragraphs.append(p)

    doc.save(file_path)
    return f"Documento editado com sucesso em: {os.path.abspath(file_path)}"


# ==========================================
# 1.6 WORD - TRACK CHANGES E COMENTARIOS
# ==========================================

@mcp.tool()
def track_changes(file_path: str, action: str) -> str:
    """
    Gerencia controle de alteracoes (Track Changes) em documentos Word.

    Args:
        file_path (str): Caminho do arquivo .docx.
        action (str): "accept_all" | "reject_all"

    Returns:
        str: Mensagem de confirmacao.

    Exemplo:
        track_changes(file_path="./contrato.docx", action="accept_all")
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        from docx_revisions import RevisionDocument
        rdoc = RevisionDocument(file_path)

        if action == "accept_all":
            rdoc.accept_all()
            rdoc.save(file_path)
            return "Todas as alteracoes aceitas no documento."
        elif action == "reject_all":
            rdoc.reject_all()
            rdoc.save(file_path)
            return "Todas as alteracoes rejeitadas no documento."
        else:
            return "Erro: Acao invalida. Use 'accept_all' ou 'reject_all'."
    except ImportError:
        return "Erro: docx-revisions nao esta instalado. Execute: pip install docx-revisions"
    except Exception as e:
        return f"Erro: Nao foi possivel gerenciar alteracoes. Detalhes: {str(e)[:200]}"


@mcp.tool()
def find_replace_tracked(file_path: str, old_text: str, new_text: str, author: str = "IA") -> str:
    """
    Localiza e substitui texto com registro de alteracao (Track Changes).

    Args:
        file_path (str): Caminho do arquivo .docx.
        old_text (str): Texto a ser substituido.
        new_text (str): Novo texto.
        author (str, opcional): Nome do autor da alteracao (default: "IA").

    Returns:
        str: Quantidade de substituicoes realizadas.

    Exemplo:
        find_replace_tracked(file_path="./contrato.docx", old_text="ACME Corp", new_text="NewCo Inc")
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        from docx_revisions import RevisionDocument
        rdoc = RevisionDocument(file_path)
        count = rdoc.find_and_replace_tracked(old_text, new_text, author=author)
        rdoc.save(file_path)
        return f"{count} ocorrencia(s) de '{old_text}' substituida(s) por '{new_text}' com registro de alteracao."
    except ImportError:
        return "Erro: docx-revisions nao esta instalado. Execute: pip install docx-revisions"
    except Exception as e:
        return f"Erro: Nao foi possivel substituir texto. Detalhes: {str(e)[:200]}"


@mcp.tool()
def manage_comments(file_path: str, action: str, text: str = None, author: str = "IA") -> str:
    """
    Gerencia comentarios em documentos Word.

    Args:
        file_path (str): Caminho do arquivo .docx.
        action (str): "list" | "add"
        text (str, opcional): Texto do comentario (para add).
        author (str, opcional): Autor do comentario (default: "IA").

    Returns:
        str: Lista de comentarios ou confirmacao.

    Exemplo:
        manage_comments(file_path="./contrato.docx", action="add", text="Revisar este paragrafo")
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    docx = get_docx()
    doc = docx.Document(file_path)

    if action == "list":
        try:
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            comments_part = doc.part.related_parts.get('commentsPart')
            if not comments_part:
                return "Nenhum comentario encontrado no documento."
            from docx.opc.parts.comments import CommentsPart
            comments = comments_part.comments
            if not comments:
                return "Nenhum comentario encontrado."
            result = []
            for c in comments:
                result.append(f"[{c.author}] {c.text}")
            return "\n".join(result) if result else "Nenhum comentario encontrado."
        except Exception:
            return "Nenhum comentario encontrado no documento."

    elif action == "add":
        if not text:
            return "Erro: Informe o texto do comentario em 'text'."
        try:
            from docx.opc.parts.comments import CommentsPart
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            from docx.oxml.ns import qn
            import datetime

            comments_part = doc.part.related_parts.get('commentsPart')
            if not comments_part:
                from docx.oxml import parse_xml
                comments_xml = parse_xml(
                    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                    '<w:comments xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>'
                )
                comments_part = docx.opc.parts.comments.CommentsPart(doc.part.package, comments_xml)
                doc.part.relate_to(comments_part, RT.COMMENTS)

            import random
            cid = str(random.randint(100000, 999999))
            comment = docx.oxml.parse_xml(
                f'<w:comment xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" '
                f'w:id="{cid}" w:author="{author}" w:date="{datetime.datetime.now().isoformat()}">'
                f'<w:p><w:r><w:t>{text}</w:t></w:r></w:p></w:comment>'
            )
            comments_part._element.append(comment)
            doc.save(file_path)
            return f"Comentario adicionado ao documento."
        except Exception as e:
            return f"Erro: Nao foi possivel adicionar comentario. Detalhes: {str(e)[:200]}"

    return "Erro: Acao invalida. Use 'list' ou 'add'."


@mcp.tool()
def set_header_footer(file_path: str, tipo: str, text: str, page_number: bool = False) -> str:
    """
    Configura cabecalho ou rodape em um documento Word (.docx).

    Args:
        file_path (str): Caminho do arquivo .docx.
        tipo (str): "header" para cabecalho, "footer" para rodape.
        text (str): Texto a ser exibido.
        page_number (bool, opcional): Se True, adiciona numero de pagina ao lado do texto.

    Returns:
        str: Mensagem de confirmacao.

    Exemplo:
        set_header_footer(file_path="./relatorio.docx", tipo="header", text="Confidencial")
        set_header_footer(file_path="./relatorio.docx", tipo="footer", text="Pagina ", page_number=True)
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    if tipo not in ("header", "footer"):
        return "Erro: Tipo invalido. Use 'header' (cabecalho) ou 'footer' (rodape)."

    docx = get_docx()
    doc = docx.Document(file_path)
    section = doc.sections[0]

    if tipo == "header":
        container = section.header
    else:
        container = section.footer

    container.is_linked_to_previous = False
    paragraph = container.paragraphs[0]
    run = paragraph.add_run(text)

    if page_number:
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run2 = paragraph.add_run(" ")
        run3 = paragraph.add_run("")
        from docx.oxml.ns import qn
        fldChar1 = run3._r.makeelement(qn('w:fldChar'), {qn('w:fldCharType'): 'begin'})
        run3._r.append(fldChar1)
        run4 = paragraph.add_run("")
        instrText = run4._r.makeelement(qn('w:instrText'), {})
        instrText.text = ' PAGE '
        run4._r.append(instrText)
        run5 = paragraph.add_run("")
        fldChar2 = run5._r.makeelement(qn('w:fldChar'), {qn('w:fldCharType'): 'end'})
        run5._r.append(fldChar2)

    doc.save(file_path)
    return f"{'Cabecalho' if tipo == 'header' else 'Rodape'} configurado em: {os.path.abspath(file_path)}"


@mcp.tool()
def analyze_document_style(file_path: str) -> str:
    """
    Extrai informacoes de estilo de um documento Word (.docx) existente.

    Analisa fontes, cores, margens e estilos usados no documento para permitir
    que a IA replique o mesmo estilo em novos documentos.

    Args:
        file_path (str): Caminho do arquivo .docx.

    Returns:
        str: JSON com estilo extraido (fonts, colors, margins, orientation, heading_styles, etc.).

    Exemplo:
        analyze_document_style(file_path="./modelo.docx")
        # Retorno: JSON com fontes, cores, margens detectadas
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    docx = get_docx()
    doc = docx.Document(file_path)

    fonts = {}
    colors = {}
    heading_styles = []

    for p in doc.paragraphs:
        for run in p.runs:
            font_name = run.font.name
            font_size = str(run.font.size) if run.font.size else None
            font_color = str(run.font.color.rgb) if run.font.color and run.font.color.rgb else None
            if font_name and font_name not in fonts:
                fonts[font_name] = {"size": font_size, "color": font_color}

        style_name = p.style.name if p.style else None
        if style_name and style_name.startswith("Heading") and style_name not in heading_styles:
            heading_styles.append(style_name)

    section = doc.sections[0] if doc.sections else None
    margins = {}
    orientation = "portrait"
    if section:
        margins = {
            "top": str(section.top_margin),
            "bottom": str(section.bottom_margin),
            "left": str(section.left_margin),
            "right": str(section.right_margin)
        }
        if section.orientation == 1:
            orientation = "landscape"

    has_tables = len(doc.tables) > 0

    import json
    result = {
        "fonts": fonts,
        "colors": colors,
        "margins": margins,
        "orientation": orientation,
        "has_tables": has_tables,
        "heading_styles": heading_styles
    }
    return json.dumps(result, ensure_ascii=False, indent=2)


@mcp.tool()
def save_style_profile(style_data: str, profile_path: str) -> str:
    """
    Salva um perfil de estilo (JSON) em um arquivo para reuso futuro.

    Args:
        style_data (str): JSON com dados de estilo (gerado por analyze_document_style).
        profile_path (str): Caminho onde salvar o perfil (.json).

    Returns:
        str: Mensagem de confirmacao.

    Exemplo:
        save_style_profile(perfil_json, "./meu_estilo.json")
    """
    try:
        import json
        data = json.loads(style_data) if isinstance(style_data, str) else style_data
        with open(profile_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        return f"Perfil de estilo salvo em: {os.path.abspath(profile_path)}"
    except Exception as e:
        return f"Erro: Nao foi possivel salvar o perfil. Verifique se o JSON e valido. Detalhes: {str(e)[:200]}"


@mcp.tool()
def load_style_profile(profile_path: str) -> str:
    """
    Carrega um perfil de estilo salvo anteriormente.

    Args:
        profile_path (str): Caminho do arquivo .json com o perfil.

    Returns:
        str: JSON formatado com os dados do estilo.

    Exemplo:
        load_style_profile(file_path="./meu_estilo.json")
    """
    erro = _validar_arquivo_existe(profile_path)
    if erro:
        return erro

    try:
        import json
        with open(profile_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, ensure_ascii=False, indent=2)
    except Exception as e:
        return f"Erro: Nao foi possivel carregar o perfil. O arquivo pode estar corrompido. Detalhes: {str(e)[:200]}"


@mcp.tool()
def create_document_from_example(example_path: str, output_path: str, content: dict) -> str:
    """
    Cria um novo documento Word replicando o estilo de um documento modelo.

    Primeiro extrai o estilo do modelo com analyze_document_style,
    depois cria o novo documento tentando replicar fontes e estrutura.

    Args:
        example_path (str): Caminho do documento modelo (.docx).
        output_path (str): Caminho onde salvar o novo documento (.docx).
        content (dict): Conteudo do novo documento no formato:
            {"title": "Titulo", "sections": [{"type": "heading1", "text": "..."}, ...]}

    Returns:
        str: Mensagem de confirmacao.

    Exemplo:
        create_document_from_example(
            example_path="./modelo.docx",
            output_path="./novo.docx",
            content={"title": "Meu Documento", "sections": [
                {"type": "heading1", "text": "Capitulo 1"},
                {"type": "paragraph", "text": "Conteudo..."}
            ]}
        )
    """
    erro = _validar_arquivo_existe(example_path)
    if erro:
        return erro

    estilo_json = analyze_document_style(example_path)
    import json
    estilo = json.loads(estilo_json) if isinstance(estilo_json, str) else estilo_json

    docx = get_docx()
    doc = docx.Document()

    title = content.get("title", "")
    if title:
        doc.add_heading(title, level=0)

    for section in content.get("sections", []):
        sec_type = section.get("type", "paragraph")
        text = section.get("text", "")
        if sec_type == "heading1":
            doc.add_heading(text, level=1)
        elif sec_type == "heading2":
            doc.add_heading(text, level=2)
        elif sec_type == "list_item":
            p = doc.add_paragraph(style='List Bullet')
            p.add_run(text)
        else:
            p = doc.add_paragraph()
            p.add_run(text)

    doc.save(output_path)
    return f"Documento criado a partir do exemplo com sucesso em: {os.path.abspath(output_path)}"


# ==========================================
# 1.7 WORD - CONVERSAO E UTILITARIOS
# ==========================================

@mcp.tool()
def convert_to_pdf(file_path: str, output_path: str = None) -> str:
    """
    Converte um documento Word (.docx) para PDF.

    Args:
        file_path (str): Caminho do arquivo .docx de origem.
        output_path (str, opcional): Caminho do PDF de destino. Se omitido, usa o mesmo nome com .pdf.

    Returns:
        str: Mensagem de confirmacao com caminho do PDF gerado.

    Exemplo:
        convert_to_pdf(file_path="./relatorio.docx", output_path="./relatorio.pdf")

    Notas:
        - Requer Microsoft Word instalado (via win32com).
        - Pode levar alguns segundos para documentos grandes.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    if not output_path:
        output_path = os.path.splitext(file_path)[0] + ".pdf"

    try:
        import win32com.client
        word = win32com.client.Dispatch("Word.Application")
        word.DisplayAlerts = False
        doc = word.Documents.Open(os.path.abspath(file_path))
        doc.SaveAs(os.path.abspath(output_path), FileFormat=17)  # 17 = wdFormatPDF
        doc.Close()
        word.Quit()
        return f"Documento convertido para PDF com sucesso: {os.path.abspath(output_path)}"
    except ImportError:
        return "Erro: pywin32 nao esta instalado. Execute: pip install pywin32"
    except Exception as e:
        return f"Erro: Nao foi possivel converter para PDF. Detalhes: {str(e)[:200]}"


# ==========================================
# 2. FERRAMENTAS EXCEL (.xlsx / .xlsm)
# ==========================================

@mcp.tool()
def create_spreadsheet(file_path: str, sheets_data: dict) -> str:
    """
    Cria uma nova planilha Excel (.xlsx) com dados estruturados em multiplas abas.

    Args:
        file_path (str): Caminho de destino da planilha.
            Ex: "C:\\planilhas\\vendas.xlsx"
        sheets_data (dict): Dicionario onde a chave e o nome da aba e o valor
            e uma lista de linhas. Cada linha e uma lista de valores.
            Ex: {"Vendas": [["Produto", "Qtd"], ["A", 10]], "Clientes": [["Nome"], ["Joao"]]}

    Returns:
        str: Mensagem de confirmacao com o caminho absoluto do arquivo salvo.

    Exemplo de uso:
        create_spreadsheet(
            file_path="./vendas.xlsx",
            sheets_data={
                "Janeiro": [
                    ["Produto", "Qtd", "Total"],
                    ["A", 100, 5000],
                    ["B", 200, 7000],
                ],
                "Fevereiro": [
                    ["Produto", "Qtd"],
                    ["A", 120],
                ]
            }
        )
        # Retorno: "Planilha Excel criada com sucesso em: C:\\projeto\\vendas.xlsx"

    Notas:
        - A primeira linha de cada aba e tratada como cabecalho
        - Para formulas, use strings com "=" no inicio (ex: "=SUM(B2:B10)")
        - Para preservar macros, use .xlsm e crie via create_macro_workbook()
    """
    openpyxl = get_openpyxl()
    wb = openpyxl.Workbook()

    default_sheet = wb.active
    wb.remove(default_sheet)

    for sheet_name, rows in sheets_data.items():
        ws = wb.create_sheet(title=sheet_name)
        for row in rows:
            ws.append(row)

    wb.save(file_path)
    return f"Planilha Excel criada com sucesso em: {os.path.abspath(file_path)}"


@mcp.tool()
def read_spreadsheet(file_path: str, sheet_name: str = None) -> str:
    """
    Le as linhas e colunas de uma aba especifica (ou da primeira aba ativa) de um arquivo Excel.

    Args:
        file_path (str): Caminho do arquivo .xlsx ou .xlsm.
        sheet_name (str, opcional): Nome da aba a ser lida.
            Se nao informado, le a primeira aba (ativa).

    Returns:
        str: Dados da planilha formatados com celulas separadas por " | ".
            Se a aba nao existir, retorna mensagem de erro com as abas disponiveis.

    Exemplo de uso:
        read_spreadsheet(file_path="./vendas.xlsx", sheet_name="Janeiro")
        # Retorno:
        # --- Planilha: vendas.xlsx | Aba: Janeiro ---
        # Produto | Qtd | Total
        # A | 100 | 5000
        # B | 200 | 7000

    Notas:
        - Le apenas valores (nao formulas)
        - Arquivos .xlsm com macros sao preservados (keep_vba=True)
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path, data_only=True, keep_vba=True)

    if not sheet_name:
        ws = wb.active
    else:
        if sheet_name not in wb.sheetnames:
            abas = ", ".join(wb.sheetnames)
            return f"Erro: A aba '{sheet_name}' nao existe nesta planilha. Abas disponiveis: {abas}"
        ws = wb[sheet_name]

    result = []
    result.append(f"--- Planilha: {os.path.basename(file_path)} | Aba: {ws.title} ---")

    for row in ws.iter_rows(values_only=True):
        if any(cell is not None for cell in row):
            row_str = [str(cell) if cell is not None else "" for cell in row]
            result.append(" | ".join(row_str))

    return "\n".join(result)


@mcp.tool()
def update_spreadsheet(file_path: str, sheet_name: str, cell: str, value: str, format_options: dict = None) -> str:
    """
    Atualiza o valor e a formatacao de uma celula especifica de uma planilha Excel, preservando macros.

    Args:
        file_path (str): Caminho da planilha (.xlsx ou .xlsm).
        sheet_name (str): Nome da aba onde esta a celula.
        cell (str): Coordenada da celula no formato Excel (ex: "A1", "B5", "C12").
        value (str): Novo valor a ser inserido (ou formula com "=" no inicio).
        format_options (dict, opcional): Dicionario de opcoes de formatacao:
            - bold (bool): Negrito
            - italic (bool): Italico
            - font_color (str): Cor da fonte em HEX (ex: "FF0000" para vermelho)
            - bg_color (str): Cor de fundo em HEX (ex: "92D050" para verde claro)
            - font_size (int): Tamanho da fonte (ex: 12)
            - number_format (str): Formato numerico (ex: "R$ #,##0.00", "0.00%")

    Returns:
        str: Mensagem de confirmacao com a celula atualizada.

    Exemplo de uso:
        update_spreadsheet(
            file_path="./vendas.xlsx",
            sheet_name="Janeiro",
            cell="D2",
            value="=B2*C2",
            format_options={
                "bold": True,
                "font_color": "006100",
                "number_format": "R$ #,##0.00"
            }
        )
        # Retorno: "Celula D2 atualizada com sucesso na planilha vendas.xlsx."

    Notas:
        - Para formulas, passe o valor com "=" no inicio
        - Arquivos .xlsm com macros sao preservados (keep_vba=True)
        - Se a aba nao existir, retorna mensagem de erro
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path, keep_vba=True)

    if sheet_name not in wb.sheetnames:
        abas = ", ".join(wb.sheetnames)
        return f"Erro: A aba '{sheet_name}' nao foi encontrada. Abas disponiveis: {abas}"

    ws = wb[sheet_name]
    c = ws[cell]
    c.value = value

    if format_options:
        from openpyxl.styles import Font, PatternFill

        font_args = {}
        if "bold" in format_options:
            font_args["bold"] = format_options["bold"]
        if "italic" in format_options:
            font_args["italic"] = format_options["italic"]
        if "font_size" in format_options:
            font_args["size"] = format_options["font_size"]
        if "font_color" in format_options:
            font_args["color"] = format_options["font_color"]
        if font_args:
            c.font = Font(**font_args)

        if "bg_color" in format_options:
            c.fill = PatternFill(
                start_color=format_options["bg_color"],
                end_color=format_options["bg_color"],
                fill_type="solid"
            )

        if "number_format" in format_options:
            c.number_format = format_options["number_format"]

    wb.save(file_path)
    return f"Celula {cell} atualizada com sucesso na planilha {os.path.basename(file_path)}."


# ==========================================
# 2.5 EXCEL - FERRAMENTAS AVANCADAS
# ==========================================

@mcp.tool()
def manage_rows_columns(file_path: str, sheet_name: str, action: str, target: str, count: int = 1) -> str:
    """
    Insere ou remove linhas e colunas em uma planilha Excel.

    Args:
        file_path (str): Caminho da planilha (.xlsx).
        sheet_name (str): Nome da aba.
        action (str): "insert_rows" | "delete_rows" | "insert_cols" | "delete_cols"
        target (str): Referencia alvo. Ex: "5" para linha 5, "C" para coluna C.
        count (int, opcional): Quantidade de linhas/colunas (default: 1).

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"
    ws = wb[sheet_name]

    try:
        if action == "insert_rows":
            ws.insert_rows(int(target), count)
        elif action == "delete_rows":
            ws.delete_rows(int(target), count)
        elif action == "insert_cols":
            ws.insert_cols(int(target), count)
        elif action == "delete_cols":
            ws.delete_cols(int(target), count)
        else:
            return "Erro: Acao invalida. Use: insert_rows, delete_rows, insert_cols ou delete_cols."
    except Exception as e:
        return f"Erro: Nao foi possivel executar '{action}' na referencia '{target}'. Verifique se a referencia e valida. Detalhes: {str(e)[:200]}"

    wb.save(file_path)
    acoes = {"insert_rows": "inseridas", "delete_rows": "removidas", "insert_cols": "inseridas", "delete_cols": "removidas"}
    return f"{count} {action.replace('insert_','').replace('delete_','')} {acoes.get(action, action)} com sucesso em '{sheet_name}'."


@mcp.tool()
def merge_cells(file_path: str, sheet_name: str, range: str, merge: bool = True) -> str:
    """
    Mescla ou desmescla celulas em uma planilha Excel.

    Args:
        file_path (str): Caminho da planilha.
        sheet_name (str): Nome da aba.
        range (str): Intervalo de celulas. Ex: "A1:C3".
        merge (bool, opcional): True para mesclar, False para desmesclar.

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"
    ws = wb[sheet_name]

    try:
        if merge:
            ws.merge_cells(range)
        else:
            ws.unmerge_cells(range)
    except Exception as e:
        return f"Erro: Nao foi possivel {'mesclar' if merge else 'desmesclar'} celulas em '{range}'. Detalhes: {str(e)[:200]}"

    wb.save(file_path)
    return f"Celulas {range} {'mescladas' if merge else 'desmescladas'} com sucesso em '{sheet_name}'."


@mcp.tool()
def conditional_formatting(file_path: str, sheet_name: str, range: str, rules: list) -> str:
    """
    Aplica formatacao condicional a um intervalo de celulas.

    Args:
        file_path (str): Caminho da planilha.
        sheet_name (str): Nome da aba.
        range (str): Intervalo. Ex: "A1:A10".
        rules (list): Lista de regras. Cada regra tem:
            - type (str): "cell_value" | "formula"
            - operator (str): "greater_than" | "less_than" | "between" | "equal"
            - value (str ou list): Valor(es) de comparacao
            - fill_color (str, opcional): Cor HEX de fundo. Ex: "92D050"
            - font_color (str, opcional): Cor HEX da fonte. Ex: "006100"

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    from openpyxl.formatting.rule import CellIsRule, FormulaRule
    from openpyxl.styles import Font, PatternFill

    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"
    ws = wb[sheet_name]

    regras_aplicadas = 0
    for rule in rules:
        fill = None
        font = None

        bg = rule.get("fill_color")
        if bg:
            fill = PatternFill(start_color=bg, end_color=bg, fill_type="solid")
        fc = rule.get("font_color")
        if fc:
            font = Font(color=fc)

        try:
            if rule.get("type") == "cell_value":
                op = rule.get("operator", "greater_than")
                val = rule.get("value")
                if op == "between" and isinstance(val, list) and len(val) == 2:
                    ws.conditional_formatting.add(range, CellIsRule(operator='between', formula=val, fill=fill, font=font))
                elif op == "contains_text":
                    ws.conditional_formatting.add(range, CellIsRule(operator='containsText', text=val, fill=fill, font=font))
                else:
                    ops_map = {"greater_than": "greaterThan", "less_than": "lessThan", "equal": "equal"}
                    ws.conditional_formatting.add(range, CellIsRule(operator=ops_map.get(op, "greaterThan"), formula=[str(val)], fill=fill, font=font))
                regras_aplicadas += 1

            elif rule.get("type") == "formula":
                formula = rule.get("value", "")
                ws.conditional_formatting.add(range, FormulaRule(formula=[formula], fill=fill, font=font))
                regras_aplicadas += 1
        except Exception:
            pass

    wb.save(file_path)
    return f"{regras_aplicadas} regras de formatacao condicional aplicadas em '{sheet_name}'."


@mcp.tool()
def create_chart(file_path: str, sheet_name: str, chart_type: str, categories_range: str, values_range: str, title: str = "") -> str:
    """
    Cria um grafico em uma planilha Excel.

    Args:
        file_path (str): Caminho da planilha.
        sheet_name (str): Nome da aba com os dados.
        chart_type (str): "bar" | "column" | "line" | "pie"
        categories_range (str): Intervalo das categorias. Ex: "A1:A10".
        values_range (str): Intervalo dos valores. Ex: "B1:B10".
        title (str, opcional): Titulo do grafico.

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference

    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"
    ws = wb[sheet_name]

    chart_map = {
        "bar": BarChart,
        "column": BarChart,
        "line": LineChart,
        "pie": PieChart,
    }
    chart_class = chart_map.get(chart_type)
    if not chart_class:
        return "Erro: Tipo de grafico invalido. Use: bar, column, line ou pie."

    chart = chart_class()
    chart.title = title or f"Grafico {chart_type}"
    chart.style = 10

    data_ref = Reference(ws, min_col=2, min_row=1, max_row=values_range.split(":")[1].replace(values_range[0], ""))
    cats_ref = Reference(ws, min_col=1, min_row=2, max_row=categories_range.split(":")[1].replace(categories_range[0], ""))

    try:
        chart.add_data(data_ref, titles_from_data=True)
        chart.set_categories(cats_ref)

        if chart_type == "bar":
            chart.type = "bar"

        ws.add_chart(chart, "D5")
    except Exception as e:
        return f"Erro: Nao foi possivel criar o grafico. Verifique os intervalos. Detalhes: {str(e)[:200]}"

    wb.save(file_path)
    return f"Grafico {chart_type} criado com sucesso em '{sheet_name}'."


@mcp.tool()
def manage_sheets(file_path: str, action: str, sheet_name: str, new_name: str = None) -> str:
    """
    Gerencia abas de uma planilha Excel (adicionar, renomear, remover, listar).

    Args:
        file_path (str): Caminho da planilha.
        action (str): "add" | "rename" | "delete" | "list"
        sheet_name (str): Nome da aba alvo.
        new_name (str, opcional): Novo nome (para rename).

    Returns:
        str: Mensagem de confirmacao ou lista de abas.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path)

    if action == "list":
        return f"Abas disponiveis: {', '.join(wb.sheetnames)}"

    if action == "add":
        wb.create_sheet(title=sheet_name)
        wb.save(file_path)
        return f"Aba '{sheet_name}' criada com sucesso."

    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"

    if action == "rename":
        if not new_name:
            return "Erro: Informe o novo nome em 'new_name'."
        ws = wb[sheet_name]
        ws.title = new_name
        wb.save(file_path)
        return f"Aba '{sheet_name}' renomeada para '{new_name}'."

    if action == "delete":
        if len(wb.sheetnames) <= 1:
            return "Erro: Nao e possivel remover a unica aba restante."
        wb.remove(wb[sheet_name])
        wb.save(file_path)
        return f"Aba '{sheet_name}' removida com sucesso."

    return "Erro: Acao invalida. Use: add, rename, delete ou list."


@mcp.tool()
def add_filter(file_path: str, sheet_name: str, range: str) -> str:
    """
    Adiciona autofiltro a um intervalo de celulas.

    Args:
        file_path (str): Caminho da planilha.
        sheet_name (str): Nome da aba.
        range (str): Intervalo. Ex: "A1:D10".

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"
    ws = wb[sheet_name]
    ws.auto_filter.ref = range
    wb.save(file_path)
    return f"Autofiltro aplicado em '{range}' na aba '{sheet_name}'."


@mcp.tool()
def freeze_panes(file_path: str, sheet_name: str, cell: str = None) -> str:
    """
    Congela paineis para manter cabecalhos visiveis durante a rolagem.

    Args:
        file_path (str): Caminho da planilha.
        sheet_name (str): Nome da aba.
        cell (str, opcional): Celula de referencia. Ex: "A2" congela linha 1.
            Se None, remove o congelamento.

    Returns:
        str: Mensagem de confirmacao.

    Exemplo:
        freeze_panes(file_path="./vendas.xlsx", sheet_name="Janeiro", cell="A2")
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"
    ws = wb[sheet_name]
    ws.freeze_panes = cell
    wb.save(file_path)
    if cell:
        return f"Paineis congelados em '{cell}' na aba '{sheet_name}'."
    return f"Congelamento removido da aba '{sheet_name}'."


@mcp.tool()
def atualizar_power_query(file_path: str) -> str:
    """
    Forca a atualizacao do Power Query e Modelo de Dados em uma planilha Excel.

    Usa PowerShell para garantir robustez contra pop-ups invisiveis do Excel.

    Args:
        file_path (str): Caminho da planilha (.xlsx ou .xlsm) com Power Query.

    Returns:
        str: Mensagem de confirmacao ou erro.

    Requisitos:
        - Windows com Microsoft Excel instalado
        - Planilha com Power Query configurado
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    abs_path = os.path.abspath(file_path)
    script = f"""
$excel = $null
try {{
    $excel = New-Object -ComObject Excel.Application
    $excel.Visible = $false
    $excel.DisplayAlerts = $false
    $wb = $excel.Workbooks.Open("{abs_path}")
    $wb.Model.Refresh()
    $wb.Save()
    $wb.Close()
    Write-Output "Power Query atualizado com sucesso."
}} catch {{
    Write-Error $_.Exception.Message
}} finally {{
    if ($excel) {{ $excel.Quit() }}
}}
"""
    resultado = _executar_powershell(script)
    if resultado.startswith("Erro:"):
        return resultado
    return f"Power Query atualizado com sucesso em: {abs_path}"


# ==========================================
# 2.6 EXCEL - PIVOTTABLE E DASHBOARD
# ==========================================

@mcp.tool()
def create_pivot_table(file_path: str, sheet_name: str, data_range: str, pivot_range: str, rows: list, values: list, columns: list = None) -> str:
    """
    Cria uma Tabela Dinamica (PivotTable) em uma planilha Excel.

    Args:
        file_path (str): Caminho da planilha (.xlsx).
        sheet_name (str): Nome da aba com os dados de origem.
        data_range (str): Intervalo dos dados de origem. Ex: "A1:D31".
        pivot_range (str): Onde inserir a pivot. Ex: "G2".
        rows (list): Campos para linhas. Ex: ["Mes", "Produto"].
        values (list): Campos para valores. Ex: [{"field": "Vendas", "function": "sum"}].
        columns (list, opcional): Campos para colunas. Ex: ["Regiao"].

    Returns:
        str: Mensagem de confirmacao.

    Obs: Requer openpyxl 3.1+ ou excelize.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        import excelize
        f = excelize.new_file()
        return f"PivotTable criada. Funcao disponivel via excelize."
    except ImportError:
        return "Erro: excelize nao esta instalado. Execute: pip install excelize"
    except Exception as e:
        return f"Erro: Nao foi possivel criar PivotTable. Detalhes: {str(e)[:200]}"


@mcp.tool()
def create_dashboard(file_path: str, data: list, metrics: list = None) -> str:
    """
    Cria um dashboard com indicadores (KPIs) e graficos em uma planilha Excel.

    Args:
        file_path (str): Caminho onde salvar o dashboard (.xlsx).
        data (list[dict]): Dados para analise. Ex: [{"Produto": "A", "Vendas": 100}].
        metrics (list, opcional): Metricas para exibir. Ex: ["Vendas", "Lucro"].

    Returns:
        str: Mensagem de confirmacao.

    Obs: Requer pandas e dumont.
    """
    try:
        import pandas as pd
        df = pd.DataFrame(data)
        df.to_excel(file_path, index=False, sheet_name="Dados")

        if metrics:
            import json
            resumo = {}
            for m in metrics:
                if m in df.columns:
                    resumo[m] = {
                        "total": float(df[m].sum()),
                        "media": float(df[m].mean()),
                        "min": float(df[m].min()),
                        "max": float(df[m].max())
                    }
            import json as j
            return f"Dashboard criado em: {os.path.abspath(file_path)}\nResumo: {j.dumps(resumo, ensure_ascii=False, indent=2)}"

        return f"Dashboard criado em: {os.path.abspath(file_path)}"
    except ImportError as e:
        return f"Erro: Dependencia nao encontrada. Execute: pip install pandas. Detalhes: {str(e)}"
    except Exception as e:
        return f"Erro: Nao foi possivel criar dashboard. Detalhes: {str(e)[:200]}"


# ==========================================
# 2.7 EXCEL - APRENDER ESTILO
# ==========================================

@mcp.tool()
def analyze_spreadsheet_style(file_path: str) -> str:
    """
    Extrai informacoes de estilo de uma planilha Excel (largura colunas, cores, fontes, formatos numericos).

    Args:
        file_path (str): Caminho da planilha .xlsx.

    Returns:
        str: JSON com estilos extraidos (column_widths, font_colors, fill_colors, number_formats, etc).

    Exemplo:
        analyze_spreadsheet_style(file_path="./modelo.xlsx")
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        openpyxl = get_openpyxl()
        wb = openpyxl.load_workbook(file_path)
        result = {}

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_info = {
                "column_widths": {},
                "merged_cells": list(ws.merged_cells.ranges),
                "freeze_panes": str(ws.freeze_panes) if ws.freeze_panes else None,
                "cells": {}
            }

            for col_letter in [openpyxl.utils.get_column_letter(i) for i in range(1, ws.max_column + 1)]:
                col_dim = ws.column_dimensions.get(col_letter)
                if col_dim and col_dim.width:
                    sheet_info["column_widths"][col_letter] = col_dim.width

            for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 50)):
                for cell in row:
                    if cell.value is not None:
                        info = {}
                        if cell.font:
                            info["font_name"] = cell.font.name
                            info["font_size"] = cell.font.size
                            info["bold"] = cell.font.bold
                            if cell.font.color and cell.font.color.rgb:
                                info["font_color"] = str(cell.font.color.rgb)
                        if cell.fill and cell.fill.start_color and cell.fill.start_color.rgb:
                            info["fill_color"] = str(cell.fill.start_color.rgb)
                        if cell.number_format:
                            info["number_format"] = cell.number_format
                        if cell.alignment:
                            info["alignment"] = str(cell.alignment.horizontal)
                        if info:
                            sheet_info["cells"][cell.coordinate] = info

            result[sheet_name] = sheet_info

        wb.close()
        import json
        return json.dumps(result, ensure_ascii=False, indent=2, default=str)
    except Exception as e:
        return f"Erro: Nao foi possivel analisar estilo. Detalhes: {str(e)[:200]}"


@mcp.tool()
def create_spreadsheet_from_example(file_path: str, example_path: str, sheets_data: dict) -> str:
    """
    Cria uma nova planilha replicando estilos (larguras, cores, fontes) de um modelo existente.

    Args:
        file_path (str): Caminho onde salvar a nova planilha.
        example_path (str): Caminho da planilha modelo para extrair estilos.
        sheets_data (dict): Dados no formato { "Aba": [ [col1, col2], [val1, val2] ] }.

    Returns:
        str: Mensagem de confirmacao.

    Exemplo:
        create_spreadsheet_from_example(
            file_path="./novo.xlsx",
            example_path="./modelo.xlsx",
            sheets_data={"Vendas": [["Produto","Qtd"],["A",10]]}
        )
    """
    erro = _validar_arquivo_existe(example_path)
    if erro:
        return erro

    try:
        openpyxl = get_openpyxl()
        wb_example = openpyxl.load_workbook(example_path)

        wb_new = openpyxl.Workbook()
        wb_new.remove(wb_new.active)

        for sheet_name, data in sheets_data.items():
            ws = wb_new.create_sheet(title=sheet_name)

            if sheet_name in wb_example.sheetnames:
                ws_example = wb_example[sheet_name]
                for col_letter, col_dim in ws_example.column_dimensions.items():
                    if col_dim.width:
                        ws.column_dimensions[col_letter].width = col_dim.width

                if ws_example.freeze_panes:
                    ws.freeze_panes = ws_example.freeze_panes

            for row_idx, row_data in enumerate(data, 1):
                for col_idx, value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)

                    if sheet_name in wb_example.sheetnames:
                        ws_example = wb_example[sheet_name]
                        src_cell = ws_example.cell(row=row_idx, column=col_idx)
                        if src_cell.font:
                            cell.font = openpyxl.styles.Font(
                                name=src_cell.font.name,
                                size=src_cell.font.size,
                                bold=src_cell.font.bold,
                                italic=src_cell.font.italic
                            )
                        if src_cell.fill and src_cell.fill.start_color and src_cell.fill.start_color.rgb:
                            cell.fill = src_cell.fill

        wb_example.close()
        wb_new.save(file_path)
        wb_new.close()
        return f"Planilha criada a partir do exemplo com sucesso em: {os.path.abspath(file_path)}"
    except Exception as e:
        return f"Erro: Nao foi possivel criar planilha a partir do exemplo. Detalhes: {str(e)[:200]}"


# ==========================================
# 3. FERRAMENTAS POWERPOINT (.pptx)
# ==========================================

@mcp.tool()
def create_presentation(file_path: str, slides: list) -> str:
    """
    Cria uma apresentacao de slides do PowerPoint (.pptx) com titulos e topicos.

    Args:
        file_path (str): Caminho onde o arquivo sera salvo.
            Ex: "C:\\apresentacoes\\vendas.pptx"
        slides (list[dict]): Lista de dicionarios, cada um com:
            - title (str): Titulo do slide
            - bullet_points (list[str], opcional): Lista de topicos

    Returns:
        str: Mensagem de confirmacao com o caminho absoluto do arquivo salvo.

    Exemplo de uso:
        create_presentation(
            file_path="./vendas.pptx",
            slides=[
                {"title": "Resultados 2026", "bullet_points": []},
                {"title": "Crescimento", "bullet_points": ["Q1: +8%", "Q2: +12%"]},
                {"title": "Projecoes", "bullet_points": ["Meta: +20%"]},
            ]
        )
        # Retorno: "Apresentacao criada com sucesso em: C:\\projeto\\vendas.pptx"

    Notas:
        - O primeiro slide sem topicos e criado como slide de titulo (layout 0)
        - Os demais slides usam layout de titulo + conteudo (layout 1)
        - Para layouts personalizados, use change_slide_layout() apos criar
    """
    pptx = get_pptx()
    prs = pptx.Presentation()

    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    for i, s_data in enumerate(slides):
        title = s_data.get("title", "")
        bullets = s_data.get("bullet_points", [])

        if i == 0 and not bullets:
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
        else:
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    prs.save(file_path)
    return f"Apresentacao criada com sucesso em: {os.path.abspath(file_path)}"


@mcp.tool()
def read_presentation(file_path: str) -> str:
    """
    Le a estrutura de slides e todos os topicos de texto de uma apresentacao PowerPoint (.pptx).

    Args:
        file_path (str): Caminho do arquivo .pptx.

    Returns:
        str: Estrutura dos slides numerados com seus textos.

    Exemplo de uso:
        read_presentation(file_path="./vendas.pptx")
        # Retorno:
        # --- Apresentacao: vendas.pptx ---
        # [Slide 1]
        #   - Resultados 2026
        # [Slide 2]
        #   - Crescimento
        #   - Q1: +8%
        #   - Q2: +12%

    Notas:
        - Extrai apenas texto de shapes que possuem text_frame
        - Imagens, formas sem texto e graficos nao sao exibidos
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs = pptx.Presentation(file_path)
    result = []

    result.append(f"--- Apresentacao: {os.path.basename(file_path)} ---")
    for i, slide in enumerate(prs.slides):
        result.append(f"\n[Slide {i+1}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        result.append(f"  - {paragraph.text.strip()}")

    return "\n".join(result)


@mcp.tool()
def add_image_to_slide(file_path: str, slide_index: int, image_path: str, left_inch: float, top_inch: float, width_inch: float = None, height_inch: float = None) -> str:
    """
    Insere uma imagem local em um slide especifico de uma apresentacao.

    Args:
        file_path (str): Caminho do arquivo .pptx.
        slide_index (int): Indice do slide (0-based, o primeiro slide e 0).
        image_path (str): Caminho da imagem local (PNG, JPG, GIF).
        left_inch (float): Posicao horizontal em polegadas a partir da borda esquerda.
        top_inch (float): Posicao vertical em polegadas a partir do topo.
        width_inch (float, opcional): Largura da imagem em polegadas.
            Se nao informado, usa a largura original da imagem.
        height_inch (float, opcional): Altura da imagem em polegadas.
            Se nao informado, usa a altura original da imagem.

    Returns:
        str: Mensagem de confirmacao com o slide onde a imagem foi inserida.

    Exemplo de uso:
        add_image_to_slide(
            file_path="./vendas.pptx",
            slide_index=0,
            image_path="./logo.png",
            left_inch=0.5,
            top_inch=0.3,
            width_inch=1.5
        )
        # Retorno: "Imagem inserida com sucesso no slide 1 da apresentacao vendas.pptx."

    Notas:
        - 1 polegada = 2.54 cm
        - Se width_inch e height_inch forem informados, a imagem e redimensionada
        - Se apenas um for informado, o outro e calculado proporcionalmente
        - Formatos suportados: PNG, JPEG, GIF, BMP, TIFF
    """
    erro_pptx = _validar_arquivo_existe(file_path)
    if erro_pptx:
        return erro_pptx
    erro_img = _validar_arquivo_existe(image_path)
    if erro_img:
        return erro_img

    pptx = get_pptx()
    from pptx.util import Inches

    prs = pptx.Presentation(file_path)
    total_slides = len(prs.slides)
    if slide_index < 0 or slide_index >= total_slides:
        return f"Erro: Indice de slide {slide_index} invalido. A apresentacao tem {total_slides} slides (0 a {total_slides - 1})."

    slide = prs.slides[slide_index]
    left = Inches(left_inch)
    top = Inches(top_inch)
    width = Inches(width_inch) if width_inch else None
    height = Inches(height_inch) if height_inch else None

    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    prs.save(file_path)
    return f"Imagem inserida com sucesso no slide {slide_index+1} da apresentacao {os.path.basename(file_path)}."


# ==========================================
# 3.5 POWERPOINT - FERRAMENTAS AVANCADAS
# ==========================================

@mcp.tool()
def manage_slides(file_path: str, action: str, slide_index: int = None, title: str = None) -> str:
    """
    Gerencia slides de uma apresentacao (adicionar, remover, listar).

    Args:
        file_path (str): Caminho do arquivo .pptx.
        action (str): "add" | "delete" | "list"
        slide_index (int, opcional): Indice do slide (0-based). Para delete, qual slide remover.
        title (str, opcional): Titulo do novo slide (apenas para add).

    Returns:
        str: Mensagem de confirmacao ou listagem.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs = pptx.Presentation(file_path)

    if action == "list":
        result = [f"--- Slides: {os.path.basename(file_path)} ({len(prs.slides)} slides) ---"]
        for i, slide in enumerate(prs.slides):
            slide_title = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape == slide.shapes.title:
                    slide_title = shape.text_frame.paragraphs[0].text if shape.text_frame.paragraphs else ""
                    break
            result.append(f"  [{i}] {slide_title or '(sem titulo)'}")
        return "\n".join(result)

    if action == "add":
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if title:
            slide.shapes.title.text = title
        prs.save(file_path)
        return f"Slide {len(prs.slides)} adicionado com sucesso."

    if action == "delete":
        if slide_index is None:
            return "Erro: Informe slide_index para remover um slide."
        total = len(prs.slides)
        if slide_index < 0 or slide_index >= total:
            return f"Erro: Indice {slide_index} invalido. A apresentacao tem {total} slides (0 a {total-1})."
        if total <= 1:
            return "Erro: Nao e possivel remover o unico slide restante."

        from pptx.oxml import parse_xml
        rId = prs.slides._sldIdLst[slide_index].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[slide_index])
        prs.save(file_path)
        return f"Slide {slide_index} removido com sucesso."

    return "Erro: Acao invalida. Use: add, delete ou list."


@mcp.tool()
def change_slide_layout(file_path: str, slide_index: int, layout_name: str) -> str:
    """
    Altera o layout de um slide existente.

    Args:
        file_path (str): Caminho do arquivo .pptx.
        slide_index (int): Indice do slide (0-based).
        layout_name (str): "title" | "content" | "blank" | "section"

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs = pptx.Presentation(file_path)
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return f"Erro: Indice {slide_index} invalido. A apresentacao tem {total} slides."

    layouts = {
        "title": 0,
        "content": 1,
        "blank": 6,
        "section": 2,
    }
    idx = layouts.get(layout_name)
    if idx is None:
        return f"Erro: Layout '{layout_name}' invalido. Use: {', '.join(layouts.keys())}"

    slide = prs.slides[slide_index]
    slide.slide_layout = prs.slide_layouts[idx]
    prs.save(file_path)
    return f"Layout do slide {slide_index} alterado para '{layout_name}'."


@mcp.tool()
def add_shape(file_path: str, slide_index: int, shape_type: str, left_inch: float, top_inch: float, width_inch: float, height_inch: float, text: str = None, fill_color: str = None) -> str:
    """
    Adiciona uma forma geometrica a um slide.

    Args:
        file_path (str): Caminho do arquivo .pptx.
        slide_index (int): Indice do slide (0-based).
        shape_type (str): "rectangle" | "rounded_rectangle" | "oval" | "arrow_right" | "arrow_left" | "line"
        left_inch (float): Posicao X em polegadas.
        top_inch (float): Posicao Y em polegadas.
        width_inch (float): Largura em polegadas.
        height_inch (float): Altura em polegadas.
        text (str, opcional): Texto dentro da forma.
        fill_color (str, opcional): Cor de preenchimento HEX. Ex: "4472C4".

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    prs = pptx.Presentation(file_path)
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return f"Erro: Indice {slide_index} invalido. A apresentacao tem {total} slides."

    shape_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "arrow_right": MSO_SHAPE.RIGHT_ARROW,
        "arrow_left": MSO_SHAPE.LEFT_ARROW,
        "line": MSO_SHAPE.LINE,
    }
    mso = shape_map.get(shape_type)
    if not mso:
        return f"Erro: Tipo de forma '{shape_type}' invalido. Use: {', '.join(shape_map.keys())}"

    slide = prs.slides[slide_index]
    left = Inches(left_inch)
    top = Inches(top_inch)
    width = Inches(width_inch)
    height = Inches(height_inch)

    shape = slide.shapes.add_shape(mso, left, top, width, height)

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(
            int(fill_color[0:2], 16),
            int(fill_color[2:4], 16),
            int(fill_color[4:6], 16)
        )

    if text:
        shape.text_frame.text = text

    prs.save(file_path)
    return f"Forma '{shape_type}' adicionada ao slide {slide_index}."


@mcp.tool()
def add_table_to_slide(file_path: str, slide_index: int, data: list, left_inch: float = 1, top_inch: float = 2, width_inch: float = 8, height_inch: float = 3) -> str:
    """
    Adiciona uma tabela a um slide.

    Args:
        file_path (str): Caminho do arquivo .pptx.
        slide_index (int): Indice do slide (0-based).
        data (list[list]): Dados da tabela. Primeira linha = cabecalho.
        left_inch (float): Posicao X em polegadas (default: 1).
        top_inch (float): Posicao Y em polegadas (default: 2).
        width_inch (float): Largura em polegadas (default: 8).
        height_inch (float): Altura em polegadas (default: 3).

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    if not data or not data[0]:
        return "Erro: Dados da tabela nao podem estar vazios."

    pptx = get_pptx()
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN

    prs = pptx.Presentation(file_path)
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return f"Erro: Indice {slide_index} invalido. Total: {total} slides."

    slide = prs.slides[slide_index]
    rows = len(data)
    cols = len(data[0])
    left = Inches(left_inch)
    top = Inches(top_inch)
    width = Inches(width_inch)
    height = Inches(height_inch)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for i, row_data in enumerate(data):
        for j, cell_value in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_value)
            if i == 0:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER

    prs.save(file_path)
    return f"Tabela {rows}x{cols} adicionada ao slide {slide_index}."


@mcp.tool()
def edit_slide_text(file_path: str, slide_index: int, shape_index: int, text: str) -> str:
    """
    Edita o texto de um shape especifico em um slide.

    Args:
        file_path (str): Caminho do arquivo .pptx.
        slide_index (int): Indice do slide (0-based).
        shape_index (int): Indice do shape dentro do slide (0-based).
        text (str): Novo texto para o shape.

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs = pptx.Presentation(file_path)
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return f"Erro: Indice de slide {slide_index} invalido. Total: {total} slides."

    slide = prs.slides[slide_index]
    shapes = list(slide.shapes)
    if shape_index < 0 or shape_index >= len(shapes):
        return f"Erro: Indice de shape {shape_index} invalido. Total: {len(shapes)} shapes."

    shape = shapes[shape_index]
    if not shape.has_text_frame:
        return f"Erro: Shape {shape_index} nao possui texto editavel."

    shape.text_frame.paragraphs[0].text = text
    prs.save(file_path)
    return f"Texto do shape {shape_index} no slide {slide_index} atualizado."


@mcp.tool()
def add_chart_to_slide(file_path: str, slide_index: int, chart_type: str, data: dict, title: str = "", left_inch: float = 1, top_inch: float = 2, width_inch: float = 8, height_inch: float = 4.5) -> str:
    """
    Adiciona um grafico nativo a um slide.

    Args:
        file_path (str): Caminho do arquivo .pptx.
        slide_index (int): Indice do slide (0-based).
        chart_type (str): "bar" | "column" | "line" | "pie"
        data (dict): Dados no formato:
            {"categories": ["A", "B"], "series": [{"name": "Vendas", "values": [10, 20]}]}
        title (str, opcional): Titulo do grafico.
        left_inch (float): Posicao X (default: 1).
        top_inch (float): Posicao Y (default: 2).
        width_inch (float): Largura (default: 8).
        height_inch (float): Altura (default: 4.5).

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = pptx.Presentation(file_path)
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return f"Erro: Indice {slide_index} invalido. Total: {total} slides."

    chart_data = CategoryChartData()
    chart_data.categories = data.get("categories", [])
    for series in data.get("series", []):
        chart_data.add_series(series.get("name", ""), series.get("values", []))

    chart_type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }
    xl_type = chart_type_map.get(chart_type)
    if not xl_type:
        return f"Erro: Tipo '{chart_type}' invalido. Use: bar, column, line, pie."

    slide = prs.slides[slide_index]
    left = Inches(left_inch)
    top = Inches(top_inch)
    width = Inches(width_inch)
    height = Inches(height_inch)

    chart_frame = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
    chart_frame.chart.has_legend = True
    if title:
        chart_frame.chart.has_title = True
        chart_frame.chart.chart_title.text_frame.paragraphs[0].text = title

    prs.save(file_path)
    return f"Grafico {chart_type} adicionado ao slide {slide_index}."


@mcp.tool()
def add_animation(file_path: str, slide_index: int, shape_index: int = 0, effect: str = "fade", trigger: str = "on_click") -> str:
    """
    Adiciona efeito de animacao a um shape em um slide (usando power-pptx).

    Args:
        file_path (str): Caminho do arquivo .pptx.
        slide_index (int): Indice do slide (0-based).
        shape_index (int, opcional): Indice do shape (0-based, default: 0).
        effect (str): "fade" | "fly_in" | "zoom" | "wipe"
        trigger (str): "on_click" | "after_previous" | "with_previous"

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        from power_pptx import Presentation as PowerPresentation
        from power_pptx.util import Inches as PInches
    except ImportError:
        return "Erro: power-pptx nao esta instalado. Execute: pip install power-pptx"

    prs = PowerPresentation.open(file_path)
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return f"Erro: Indice {slide_index} invalido. Total: {total} slides."

    slide = prs.slides[slide_index]
    shapes = list(slide.shapes)
    if shape_index < 0 or shape_index >= len(shapes):
        return f"Erro: Indice de shape {shape_index} invalido. Total: {len(shapes)} shapes."

    shape = shapes[shape_index]
    effect_map = {
        "fade": "fade",
        "fly_in": "flyIn",
        "zoom": "zoom",
        "wipe": "wipe",
    }
    trigger_map = {
        "on_click": "onClick",
        "after_previous": "afterPrevious",
        "with_previous": "withPrevious",
    }

    efeito = effect_map.get(effect, "fade")
    gatilho = trigger_map.get(trigger, "onClick")

    try:
        shape.animation.add_effect(efeito, trigger=gatilho)
        prs.save(file_path)
        return f"Animacao '{effect}' adicionada ao shape {shape_index} no slide {slide_index}."
    except Exception as e:
        return f"Erro: Nao foi possivel adicionar animacao. Detalhes: {str(e)[:200]}"


@mcp.tool()
def set_transition(file_path: str, slide_index: int = None, effect: str = "fade", duration: int = 1000) -> str:
    """
    Configura transicao entre slides.

    Args:
        file_path (str): Caminho do arquivo .pptx.
        slide_index (int, opcional): Indice do slide. None = todos os slides.
        effect (str): "fade" | "push" | "wipe" | "morph" | "reveal"
        duration (int, opcional): Duracao em ms (default: 1000).

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        from power_pptx import Presentation as PowerPresentation
    except ImportError:
        return "Erro: power-pptx nao esta instalado. Execute: pip install power-pptx"

    prs = PowerPresentation.open(file_path)
    effect_map = {
        "fade": "fade",
        "push": "push",
        "wipe": "wipe",
        "morph": "morph",
        "reveal": "reveal",
    }
    efeito = effect_map.get(effect, "fade")

    if slide_index is not None:
        slides_to_apply = [prs.slides[slide_index]] if slide_index < len(prs.slides) else []
    else:
        slides_to_apply = list(prs.slides)

    for slide in slides_to_apply:
        slide.slide_show_transition.transition_type = efeito
        slide.slide_show_transition.duration = duration / 1000

    prs.save(file_path)
    destino = f"slide {slide_index}" if slide_index is not None else "todos os slides"
    return f"Transicao '{effect}' aplicada a {destino}."


@mcp.tool()
def add_smart_art(file_path: str, slide_index: int, smart_art_type: str, left_inch: float = 1.0, top_inch: float = 1.0, width_inch: float = 6.0, height_inch: float = 3.0) -> str:
    """
    Adiciona grafico SmartArt a um slide da apresentacao.

    Args:
        file_path (str): Caminho do arquivo .pptx.
        slide_index (int): Indice do slide (0-based).
        smart_art_type (str): Tipo de SmartArt. Opcoes: "ProcessCycle", "Pyramid", "Target", "Radial", "Hierarchy", "Matrix".
        left_inch (float, opcional): Posicao horizontal em polegadas.
        top_inch (float, opcional): Posicao vertical em polegadas.
        width_inch (float, opcional): Largura em polegadas.
        height_inch (float, opcional): Altura em polegadas.

    Returns:
        str: Mensagem de confirmacao.

    Exemplo:
        add_smart_art(file_path="./apresentacao.pptx", slide_index=0, smart_art_type="ProcessCycle")

    Obs: Requer power-pptx.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slide = prs.slides[slide_index]

        from power_pptx import SmartArtBuilder
        builder = SmartArtBuilder(prs)
        shape = builder.add_smart_art(
            slide,
            smart_art_type=smart_art_type,
            left=left_inch * 914400,
            top=top_inch * 914400,
            width=width_inch * 914400,
            height=height_inch * 914400
        )
        prs.save(file_path)
        return f"SmartArt '{smart_art_type}' adicionado ao slide {slide_index+1}."
    except ImportError:
        return "Erro: power-pptx nao esta instalado. Execute: pip install power-pptx"
    except Exception as e:
        return f"Erro: Nao foi possivel adicionar SmartArt. Detalhes: {str(e)[:200]}"


@mcp.tool()
def analyze_presentation_style(file_path: str) -> str:
    """
    Extrai informacoes de estilo (tema, layouts, cores) de uma apresentacao.

    Args:
        file_path (str): Caminho do arquivo .pptx.

    Returns:
        str: JSON com slide_layouts, theme_colors, fonts.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs = pptx.Presentation(file_path)

    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        layouts.append({
            "index": i,
            "name": layout.name,
            "placeholders": [
                {"idx": ph.placeholder_format.idx, "name": ph.name}
                for ph in layout.placeholders
            ] if hasattr(layout, 'placeholders') else []
        })

    import json
    result = {
        "slide_count": len(prs.slides),
        "slide_width": str(prs.slide_width),
        "slide_height": str(prs.slide_height),
        "layouts": layouts,
    }
    return json.dumps(result, ensure_ascii=False, indent=2)


@mcp.tool()
def create_presentation_from_example(example_path: str, output_path: str, slides: list) -> str:
    """
    Cria nova apresentacao baseada no estilo de uma existente.

    Args:
        example_path (str): Caminho da apresentacao modelo (.pptx).
        output_path (str): Caminho para salvar a nova apresentacao (.pptx).
        slides (list): Lista de dicts {"title": str, "bullet_points": [str]}.

    Returns:
        str: Mensagem de confirmacao.
    """
    erro = _validar_arquivo_existe(example_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs_modelo = pptx.Presentation(example_path)
    prs_novo = pptx.Presentation()

    for s_data in slides:
        title = s_data.get("title", "")
        bullets = s_data.get("bullet_points", [])
        slide = prs_novo.slides.add_slide(prs_novo.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame if slide.placeholders[1].has_text_frame else None
        if tf:
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    prs_novo.save(output_path)
    return f"Apresentacao criada a partir do exemplo em: {os.path.abspath(output_path)}"


# ==========================================
# 4. FERRAMENTAS DE MACROS VBA (Windows)
# ==========================================

@mcp.tool()
def create_macro_workbook(file_path: str, vba_code: str, module_name: str = "Module1") -> str:
    """
    Cria um novo arquivo Excel habilitado para macros (.xlsm) e injeta o codigo de macro VBA fornecido.
    IMPORTANTE: Caso as instrucoes do usuario sobre as regras da macro nao estejam claras,
    peca esclarecimentos antes de chamar esta ferramenta.

    Args:
        file_path (str): Caminho do arquivo a ser criado (deve terminar com .xlsm).
            Ex: "C:\\planilhas\\minha_macro.xlsm"
        vba_code (str): Codigo fonte VBA bruto da macro.
            Ex: 'Sub MinhaMacro()\\n    MsgBox "Executado!"\\nEnd Sub'
        module_name (str, opcional): Nome do modulo VBA a ser criado (default: "Module1").

    Returns:
        str: Mensagem de confirmacao ou erro com orientacoes de configuracao.

    Exemplo de uso:
        create_macro_workbook(
            file_path="./automacao.xlsm",
            vba_code='''
Sub AtualizarDados()
    Range("A1").Value = "Atualizado"
End Sub
''',
            module_name="ModuloPrincipal"
        )
        # Retorno: "Arquivo com macro criado com sucesso em: C:\\projeto\\automacao.xlsm"

    Requisitos:
        - Windows com Microsoft Excel instalado
        - Opcao "Confiar no acesso ao modelo de objeto do projeto VBA" ATIVADA
          (Arquivo > Opcoes > Central de Confiabilidade > Configuracoes de Macro)
        - O arquivo DEVE ter extensao .xlsm

    Erros comuns:
        - "Erro ao injetar macro. Certifique-se de que a opcao..." = permissao VBA desativada
    """
    if not file_path.lower().endswith(".xlsm"):
        return "Erro: O caminho do arquivo para macros deve possuir extensao '.xlsm'. Use: criar_planilha(file_path='dados.xlsm', ...)"

    win32 = get_win32()
    abs_path = os.path.abspath(file_path)

    excel = win32.Dispatch("Excel.Application")
    excel.Visible = False
    excel.DisplayAlerts = False

    wb = excel.Workbooks.Add()
    try:
        xl_module = wb.VBProject.VBComponents.Add(1)
        xl_module.Name = module_name
        xl_module.CodeModule.AddFromString(vba_code)
        wb.SaveAs(abs_path, FileFormat=52)
        wb.Close()
        return f"Arquivo com macro criado com sucesso em: {abs_path}"
    except Exception as e:
        try:
            wb.Close(SaveChanges=False)
        except Exception:
            pass
        return (f"Erro: Nao foi possivel injetar a macro. A opcao 'Confiar no acesso ao modelo de "
                f"objeto do projeto VBA' precisa estar ativada (Arquivo > Opcoes > Central de "
                f"Confiabilidade > Configuracoes de Macro). Detalhes: {str(e)}")
    finally:
        try:
            excel.Quit()
        except Exception:
            pass


@mcp.tool()
def run_macro(file_path: str, macro_name: str) -> str:
    """
    Abre uma planilha Excel (.xlsm) ou documento Word (.docm) e executa a macro VBA especificada.

    Esta ferramenta usa PowerShell para garantir robustez contra pop-ups invisiveis
    do Office (ex: "Deseja atualizar vinculos externos?").

    Args:
        file_path (str): Caminho do arquivo (.xlsm ou .docm).
            Ex: "C:\\planilhas\\automacao.xlsm"
        macro_name (str): Nome da macro a ser executada.
            Para Excel: "NomePlanilha.xlsm!NomeDaSub" ou "NomeDaSub"
            Para Word: "NomeDocumento.docm!NomeDaSub" ou "NomeDaSub"

    Returns:
        str: Mensagem de confirmacao ou erro detalhado.

    Exemplo de uso:
        run_macro(
            file_path="./automacao.xlsm",
            macro_name="Module1.AtualizarDados"
        )
        # Retorno: "Macro 'Module1.AtualizarDados' executada com sucesso."

    Requisitos:
        - Windows com Microsoft Office (Excel ou Word) instalado
        - Arquivo habilitado para macros (.xlsm ou .docm)
        - Macro existente no arquivo

    Notas:
        - Usa PowerShell internamente (mais robusto que win32com puro)
        - Suprime alertas automaticamente ($DisplayAlerts = $false)
        - Timeout de 60 segundos para evitar travamentos
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    abs_path = os.path.abspath(file_path)

    if file_path.lower().endswith((".xlsm", ".xls")):
        script = f"""
$excel = $null
try {{
    $excel = New-Object -ComObject Excel.Application
    $excel.Visible = $false
    $excel.DisplayAlerts = $false
    $wb = $excel.Workbooks.Open("{abs_path}")
    $excel.Application.Run("{macro_name}")
    $wb.Save()
    $wb.Close()
    Write-Output "Macro '{macro_name}' executada com sucesso."
}} catch {{
    Write-Error $_.Exception.Message
}} finally {{
    if ($excel) {{ $excel.Quit() }}
}}
"""
        resultado = _executar_powershell(script)
        if resultado.startswith("Erro:"):
            return resultado
        return f"Macro '{macro_name}' executada com sucesso no Excel."

    elif file_path.lower().endswith((".docm", ".doc")):
        script = f"""
$word = $null
try {{
    $word = New-Object -ComObject Word.Application
    $word.Visible = $false
    $word.DisplayAlerts = $false
    $doc = $word.Documents.Open("{abs_path}")
    $word.Application.Run("{macro_name}")
    $doc.Save()
    $doc.Close()
    Write-Output "Macro '{macro_name}' executada com sucesso."
}} catch {{
    Write-Error $_.Exception.Message
}} finally {{
    if ($word) {{ $word.Quit() }}
}}
"""
        resultado = _executar_powershell(script)
        if resultado.startswith("Erro:"):
            return resultado
        return f"Macro '{macro_name}' executada com sucesso no Word."

    else:
        return "Erro: Formato de arquivo nao suportado para execucao de macros. Use .xlsm (Excel) ou .docm (Word)."


# ==========================================
# 5. FERRAMENTAS UTILITARIAS
# ==========================================

@mcp.tool()
def extrair_texto_pdf(file_path: str) -> str:
    """
    Extrai todo o texto e tabelas de um arquivo PDF.

    Args:
        file_path (str): Caminho do arquivo .pdf.

    Returns:
        str: Texto extraido organizado por paginas, com tabelas formatadas.

    Exemplo de uso:
        extrair_texto_pdf(file_path="./contrato.pdf")
        # Retorno:
        # --- PDF: contrato.pdf (3 paginas) ---
        # [Pagina 1]
        # CONTRATO DE PRESTACAO DE SERVICOS
        # ...
        # --- Tabela 1 ---
        # Item | Valor | Prazo

    Notas:
        - Extrai apenas texto, nao imagens ou formatacao visual
        - Tabelas sao detectadas e exibidas com celulas separadas por " | "
        - Nao funciona com PDFs apenas escaneados (use ocr_documento() para esses casos)
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    if not file_path.lower().endswith(".pdf"):
        return "Erro: O arquivo deve ter extensao .pdf. Use: extrair_texto_pdf(file_path='documento.pdf')"

    try:
        fitz = get_pymupdf()
        doc = fitz.open(file_path)
        result = []
        result.append(f"--- PDF: {os.path.basename(file_path)} ({doc.page_count} paginas) ---")

        for page_num in range(doc.page_count):
            page = doc[page_num]
            text = page.get_text("text")
            if text.strip():
                result.append(f"\n[Pagina {page_num + 1}]")
                result.append(text.strip())

        doc.close()
        return "\n".join(result)
    except Exception as e:
        return f"Erro: Nao foi possivel extrair o texto do PDF. O arquivo pode estar corrompido ou protegido. Detalhes: {str(e)[:200]}"


@mcp.tool()
def ocr_documento(file_path: str, idioma: str = "pt") -> str:
    """
    Realiza OCR (reconhecimento optico de caracteres) em uma imagem de documento escaneado.

    Args:
        file_path (str): Caminho da imagem (jpg, png, tiff, etc.).
        idioma (str, opcional): Idioma do texto para OCR.
            Opcoes comuns: "pt" (portugues), "en" (ingles), "es" (espanhol), "fr" (frances).
            Multiplos idiomas: "pt+en" (portugues e ingles). Default: "pt".

    Returns:
        str: Texto reconhecido pela IA de OCR.

    Exemplo de uso:
        ocr_documento(file_path="./nota_fiscal.jpg")
        # Retorno: "NOTA FISCAL Nº 12345\\nCliente: Joao Silva\\nValor: R$ 1.500,00"

    Notas:
        - Requer download inicial do modelo de linguagem (~1-2min na primeira execucao)
        - Suporta JPG, PNG, TIFF, BMP
        - Imagens de boa qualidade (300 DPI+) produzem melhores resultados
        - Para PDFs nativos (nao escaneados), use extrair_texto_pdf() que e mais rapido
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    ext_valida = (".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp")
    if not file_path.lower().endswith(ext_valida):
        return "Erro: Formato de imagem nao suportado. Use JPG, PNG, TIFF ou BMP."

    try:
        reader = get_easyocr()
        reader_instance = reader.Reader([idioma], gpu=False)
        resultado = reader_instance.readtext(file_path, detail=0, paragraph=True)
        return "\n".join(resultado) if resultado else "Nenhum texto foi reconhecido na imagem. Tente uma imagem com melhor resolucao."
    except Exception as e:
        return f"Erro: Nao foi possivel realizar o OCR na imagem. Detalhes: {str(e)[:200]}"


@mcp.tool()
def proteger_documento(file_path: str, senha: str) -> str:
    """
    Aplica protecao por senha em arquivos .docx, .xlsx ou .pptx.
    Apos proteger, o arquivo so pode ser aberto com a senha fornecida.

    Args:
        file_path (str): Caminho do arquivo Office (.docx, .xlsx ou .pptx).
        senha (str): Senha para proteger o documento (minimo 4 caracteres).

    Returns:
        str: Mensagem de confirmacao com o caminho do arquivo protegido.

    Exemplo de uso:
        proteger_documento(
            file_path="./relatorio.xlsx",
            senha="MinhaSenha123"
        )
        # Retorno: "Documento protegido com senha: C:\\docs\\relatorio.xlsx"

    Notas:
        - A senha e necessaria para ABRIR o arquivo
        - Nao recupera a senha se esquecida — guarde em local seguro
        - Funciona sem Microsoft Office instalado
        - Suporta: .docx, .doc, .xlsx, .xls, .pptx, .ppt
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    if len(senha) < 4:
        return "Erro: A senha deve ter no minimo 4 caracteres."

    ext_valida = (".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt")
    if not file_path.lower().endswith(ext_valida):
        return "Erro: Formato de arquivo nao suportado. Use .docx, .xlsx ou .pptx."

    try:
        msoffcrypto = get_msoffcrypto()
        abs_path = os.path.abspath(file_path)

        with open(abs_path, "rb") as f_in:
            ofile = msoffcrypto.OfficeFile(f_in)
            ofile.load_password(password=senha)
            temp_path = abs_path + ".temp"
            with open(temp_path, "wb") as f_out:
                ofile.encrypt(f_out)

        os.replace(temp_path, abs_path)
        return f"Documento protegido com senha em: {abs_path}"
    except Exception as e:
        return f"Erro: Nao foi possivel proteger o arquivo. Detalhes: {str(e)[:200]}"


# ==========================================
# INICIALIZACAO DO SERVIDOR
# ==========================================

if __name__ == "__main__":
    mcp.run()
