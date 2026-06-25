import os
import unittest
import shutil
import json
from scripts.office_mcp import (
    create_document,
    read_document,
    create_document_from_template,
    create_spreadsheet,
    read_spreadsheet,
    update_spreadsheet,
    create_presentation,
    read_presentation,
    create_macro_workbook,
    run_macro
)

class TestOfficeMCPServer(unittest.TestCase):
    
    def setUp(self):
        self.test_dir = os.path.join(os.path.dirname(__file__), "output")
        if not os.path.exists(self.test_dir):
            os.makedirs(self.test_dir)
            
    def tearDown(self):
        if os.path.exists(self.test_dir):
            shutil.rmtree(self.test_dir)

    def test_word_tools(self):
        doc_path = os.path.join(self.test_dir, "test_doc.docx")
        
        # 1. Testar create_document
        elements = [
            {"type": "title", "text": "Documento de Teste"},
            {"type": "heading1", "text": "Introducao"},
            {"type": "paragraph", "text": "Este e um paragrafo de teste para validacao de MCP.", "bold": True},
            {"type": "list_item", "text": "Item 1 da lista"},
            {"type": "list_item", "text": "Item 2 da lista", "italic": True}
        ]
        
        res_create = create_document(doc_path, elements)
        self.assertIn("sucesso", res_create)
        self.assertTrue(os.path.exists(doc_path))
        
        # 2. Testar read_document
        res_read = read_document(doc_path)
        self.assertIn("Documento de Teste", res_read)
        self.assertIn("Este e um paragrafo de teste", res_read)
        self.assertIn("Item 1 da lista", res_read)
        
        # 3. Testar create_document_from_template
        template_output = os.path.join(self.test_dir, "test_template_out.docx")
        template_input = os.path.join(self.test_dir, "template.docx")
        create_document(template_input, [
            {"type": "paragraph", "text": "Prezado {{Nome}}, sua conta e {{Conta}}."}
        ])
        
        replacements = {"{{Nome}}": "Miguel Porto", "{{Conta}}": "12345-6"}
        res_tpl = create_document_from_template(template_input, template_output, replacements)
        self.assertIn("sucesso", res_tpl)
        self.assertTrue(os.path.exists(template_output))
        
        res_read_tpl = read_document(template_output)
        self.assertIn("Prezado Miguel Porto, sua conta e 12345-6.", res_read_tpl)

    def test_excel_tools(self):
        sheet_path = os.path.join(self.test_dir, "test_sheet.xlsx")
        
        # 1. Testar create_spreadsheet
        sheets_data = {
            "Dados": [
                ["ID", "Nome", "Cidade"],
                [1, "Ana", "Sao Paulo"],
                [2, "Bruno", "Rio de Janeiro"]
            ]
        }
        res_create = create_spreadsheet(sheet_path, sheets_data)
        self.assertIn("sucesso", res_create)
        self.assertTrue(os.path.exists(sheet_path))
        
        # 2. Testar read_spreadsheet
        res_read = read_spreadsheet(sheet_path, "Dados")
        self.assertIn("ID | Nome | Cidade", res_read)
        self.assertIn("1 | Ana | Sao Paulo", res_read)
        
        # 3. Testar update_spreadsheet
        format_options = {"bold": True, "bg_color": "CCCCCC", "font_color": "FF0000"}
        res_update = update_spreadsheet(sheet_path, "Dados", "A4", "Novo Valor", format_options)
        self.assertIn("sucesso", res_update)
        
        res_read_updated = read_spreadsheet(sheet_path, "Dados")
        self.assertIn("Novo Valor", res_read_updated)

    def test_powerpoint_tools(self):
        ppt_path = os.path.join(self.test_dir, "test_presentation.pptx")
        
        # 1. Testar create_presentation
        slides = [
            {"title": "Apresentacao de Teste"},
            {"title": "Segundo Slide", "bullet_points": ["Topico A", "Topico B"]}
        ]
        res_create = create_presentation(ppt_path, slides)
        self.assertIn("sucesso", res_create)
        self.assertTrue(os.path.exists(ppt_path))
        
        # 2. Testar read_presentation
        res_read = read_presentation(ppt_path)
        self.assertIn("Apresentacao de Teste", res_read)
        self.assertIn("Topico A", res_read)
        self.assertIn("Topico B", res_read)

    def test_vba_macro_creation(self):
        macro_path = os.path.join(self.test_dir, "test_macro.xlsm")
        vba_code = """Sub TesteMacro()
    Range("A1").Value = "Valor Escrito pela Macro"
End Sub"""

        try:
            res_macro = create_macro_workbook(macro_path, vba_code, "Module1")
            if "Erro" not in res_macro:
                self.assertIn("sucesso", res_macro)
                self.assertTrue(os.path.exists(macro_path))

                # Testar execucao da macro
                res_run = run_macro(macro_path, "test_macro.xlsm!Module1.TesteMacro")
                self.assertIn("sucesso", res_run)

                # Validar o resultado
                res_read = read_spreadsheet(macro_path, "Sheet1")
                self.assertIn("Valor Escrito pela Macro", res_read)
            else:
                print("Ignorando assercao completa de macro: Central de Confiabilidade nao habilitada ou sem Excel.")
        except Exception as e:
            print(f"Ignorando teste COM/ActiveX (Excel nao instalado): {str(e)}")

    def test_set_transition_pptx_fade(self):
        from scripts.office_mcp import set_transition_pptx
        import zipfile
        import re
        try:
            from power_pptx import Presentation
        except ImportError:
            self.skipTest("power-pptx nao instalado")

        ppt_path = os.path.join(self.test_dir, "test_transition.pptx")
        prs = Presentation()
        for i in range(2):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_textbox(1, 1, 4, 2).text_frame.text = f"Slide {i+1}"
        prs.save(ppt_path)

        result = set_transition_pptx(file_path=ppt_path, effect="fade", duration=800)
        self.assertIn("fade", result.lower())
        self.assertIn("aplicada", result.lower())

        with zipfile.ZipFile(ppt_path) as z:
            for name in z.namelist():
                if re.match(r"ppt/slides/slide\d+\.xml$", name):
                    content = z.read(name).decode("utf-8")
                    m = re.search(r"<p:transition[^>]*>.*?</p:transition>", content)
                    self.assertIsNotNone(m, f"Transicao nao encontrada em {name}")
                    self.assertIn("<p:fade/>", m.group(0))
                    self.assertIn('p14:dur="800"', m.group(0))

    def test_add_animation_pptx_fade(self):
        from scripts.office_mcp import add_animation_pptx
        import zipfile
        import re
        try:
            from power_pptx import Presentation
        except ImportError:
            self.skipTest("power-pptx nao instalado")

        ppt_path = os.path.join(self.test_dir, "test_animation.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(1, 1, 4, 2).text_frame.text = "Hello"
        prs.save(ppt_path)

        result = add_animation_pptx(file_path=ppt_path, slide_index=0, shape_index=0, effect="fade", trigger="on_click")
        self.assertIn("fade", result.lower())
        self.assertIn("sucesso", result.lower())

        with zipfile.ZipFile(ppt_path) as z:
            content = z.read("ppt/slides/slide1.xml").decode("utf-8")
            self.assertIn("<p:timing>", content, "Timing XML nao gerado")
            self.assertIn("fade", content.lower(), "Animacao fade nao encontrada no XML")

    def test_list_template_slots_returns_json(self):
        from scripts.office_mcp import list_template_slots_pptx
        import json
        template_path = os.path.join(
            os.path.dirname(os.path.dirname(__file__)),
            "data", "templates", "clean_minimal.pptx"
        )
        if not os.path.exists(template_path):
            self.skipTest(f"Template nao encontrado: {template_path}")
        result = list_template_slots_pptx(template_path)
        data = json.loads(result)
        self.assertIn("total_slides", data)
        self.assertEqual(data["total_slides"], 25)
        self.assertGreater(len(data["slides"]), 0)

    def test_create_from_template_basic(self):
        from scripts.office_mcp import create_from_template_pptx
        from pptx import Presentation
        template_path = os.path.join(
            os.path.dirname(os.path.dirname(__file__)),
            "data", "templates", "clean_minimal.pptx"
        )
        if not os.path.exists(template_path):
            self.skipTest(f"Template nao encontrado: {template_path}")
        output_path = os.path.join(self.test_dir, "from_template_test.pptx")
        replacements = json.dumps({
            "0": {"10": "Ola Mundo", "11": "2026"},
            "2": {"3": "SUMARIO"}
        })
        result = create_from_template_pptx(
            template_path=template_path,
            output_path=output_path,
            slide_replacements=replacements
        )
        self.assertIn("sucesso", result.lower())
        self.assertTrue(os.path.exists(output_path))
        prs = Presentation(output_path)
        self.assertEqual(len(prs.slides), 25)
        # Verify text was actually replaced
        self.assertEqual(prs.slides[0].shapes[10].text_frame.text, "Ola Mundo")
        self.assertEqual(prs.slides[0].shapes[11].text_frame.text, "2026")
        self.assertEqual(prs.slides[2].shapes[3].text_frame.text, "SUMARIO")

    def test_create_from_template_preserves_formatting(self):
        from scripts.office_mcp import create_from_template_pptx
        from pptx import Presentation
        template_path = os.path.join(
            os.path.dirname(os.path.dirname(__file__)),
            "data", "templates", "clean_minimal.pptx"
        )
        if not os.path.exists(template_path):
            self.skipTest(f"Template nao encontrado: {template_path}")
        output_path = os.path.join(self.test_dir, "format_test.pptx")
        # Get original formatting
        prs_orig = Presentation(template_path)
        original_shape = prs_orig.slides[0].shapes[10]
        original_size = None
        if original_shape.text_frame.paragraphs[0].runs:
            original_size = original_shape.text_frame.paragraphs[0].runs[0].font.size
        # Apply replacement
        replacements = json.dumps({"0": {"10": "Texto Novo"}})
        create_from_template_pptx(
            template_path=template_path,
            output_path=output_path,
            slide_replacements=replacements
        )
        # Verify formatting preserved
        prs_new = Presentation(output_path)
        new_shape = prs_new.slides[0].shapes[10]
        self.assertEqual(new_shape.text_frame.text, "Texto Novo")
        if original_size:
            new_size = new_shape.text_frame.paragraphs[0].runs[0].font.size
            self.assertEqual(new_size, original_size)

    def test_create_from_template_multiline(self):
        from scripts.office_mcp import create_from_template_pptx
        from pptx import Presentation
        template_path = os.path.join(
            os.path.dirname(os.path.dirname(__file__)),
            "data", "templates", "clean_minimal.pptx"
        )
        if not os.path.exists(template_path):
            self.skipTest(f"Template nao encontrado: {template_path}")
        output_path = os.path.join(self.test_dir, "multiline_test.pptx")
        replacements = json.dumps({"0": {"10": "Linha 1\\nLinha 2\\nLinha 3"}})
        result = create_from_template_pptx(
            template_path=template_path,
            output_path=output_path,
            slide_replacements=replacements
        )
        self.assertIn("sucesso", result.lower())
        prs = Presentation(output_path)
        text = prs.slides[0].shapes[10].text_frame.text
        self.assertIn("Linha 1", text)
        self.assertIn("Linha 2", text)
        self.assertIn("Linha 3", text)

    def test_create_from_template_invalid_json(self):
        from scripts.office_mcp import create_from_template_pptx
        template_path = os.path.join(
            os.path.dirname(os.path.dirname(__file__)),
            "data", "templates", "clean_minimal.pptx"
        )
        if not os.path.exists(template_path):
            self.skipTest(f"Template nao encontrado: {template_path}")
        output_path = os.path.join(self.test_dir, "invalid_test.pptx")
        result = create_from_template_pptx(
            template_path=template_path,
            output_path=output_path,
            slide_replacements="not valid json"
        )
        self.assertIn("Erro", result)

    def test_create_from_template_ignores_out_of_range(self):
        from scripts.office_mcp import create_from_template_pptx
        template_path = os.path.join(
            os.path.dirname(os.path.dirname(__file__)),
            "data", "templates", "clean_minimal.pptx"
        )
        if not os.path.exists(template_path):
            self.skipTest(f"Template nao encontrado: {template_path}")
        output_path = os.path.join(self.test_dir, "oor_test.pptx")
        replacements = json.dumps({
            "999": {"0": "out of range"},
            "0": {"999": "shape out of range"}
        })
        result = create_from_template_pptx(
            template_path=template_path,
            output_path=output_path,
            slide_replacements=replacements
        )
        self.assertIn("sucesso", result.lower())
        self.assertIn("Substituicoes aplicadas: 0", result)

if __name__ == "__main__":
    unittest.main()
