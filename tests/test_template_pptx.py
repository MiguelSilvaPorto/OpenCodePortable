"""Tests for template-based presentation generation (create_from_template_pptx, list_template_slots_pptx)."""
import os
import sys
import unittest
import json
import shutil

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))

from scripts.office_mcp import list_template_slots_pptx, create_from_template_pptx

TEMPLATE_PATH = os.path.join(
    os.path.dirname(os.path.dirname(__file__)),
    "data", "templates", "clean_minimal.pptx"
)


class TestTemplateFunctions(unittest.TestCase):
    """Tests for template listing and creation functions."""

    def setUp(self):
        self.test_dir = os.path.join(
            os.path.dirname(__file__), "output_template"
        )
        if not os.path.exists(self.test_dir):
            os.makedirs(self.test_dir)

    def tearDown(self):
        # Clean up generated files
        for f in os.listdir(self.test_dir):
            if f.endswith('.pptx'):
                os.remove(os.path.join(self.test_dir, f))

    def test_list_template_slots_returns_json(self):
        if not os.path.exists(TEMPLATE_PATH):
            self.skipTest(f"Template nao encontrado: {TEMPLATE_PATH}")
        result = list_template_slots_pptx(TEMPLATE_PATH)
        data = json.loads(result)
        self.assertIn("total_slides", data)
        self.assertIn("slides", data)
        self.assertEqual(data["total_slides"], 25)
        self.assertGreater(len(data["slides"]), 0)
        first_slide = data["slides"][0]
        self.assertIn("shapes", first_slide)
        self.assertGreater(len(first_slide["shapes"]), 0)
        for shape in first_slide["shapes"]:
            self.assertIn("shape_id", shape)
            self.assertIn("current_text", shape)
            self.assertIn("shape_index", shape)

    def test_list_template_slots_invalid_path(self):
        result = list_template_slots_pptx(r"C:\nonexistent\file.pptx")
        self.assertIn("Erro", result)

    def test_create_from_template_basic(self):
        if not os.path.exists(TEMPLATE_PATH):
            self.skipTest(f"Template nao encontrado: {TEMPLATE_PATH}")
        output_path = os.path.join(self.test_dir, "test1.pptx")
        replacements = json.dumps({
            "0": {"10": "Ola Mundo", "11": "2026"},
            "2": {"3": "SUMARIO"}
        })
        result = create_from_template_pptx(
            template_path=TEMPLATE_PATH,
            output_path=output_path,
            slide_replacements=replacements
        )
        self.assertIn("sucesso", result.lower())
        self.assertTrue(os.path.exists(output_path))
        from pptx import Presentation
        prs = Presentation(output_path)
        self.assertEqual(len(prs.slides), 25)
        self.assertEqual(prs.slides[0].shapes[10].text_frame.text, "Ola Mundo")
        self.assertEqual(prs.slides[0].shapes[11].text_frame.text, "2026")
        self.assertEqual(prs.slides[2].shapes[3].text_frame.text, "SUMARIO")

    def test_create_from_template_preserves_formatting(self):
        if not os.path.exists(TEMPLATE_PATH):
            self.skipTest(f"Template nao encontrado: {TEMPLATE_PATH}")
        output_path = os.path.join(self.test_dir, "test_format.pptx")
        from pptx import Presentation
        prs_orig = Presentation(TEMPLATE_PATH)
        original_shape = prs_orig.slides[0].shapes[10]
        original_size = None
        if original_shape.text_frame.paragraphs[0].runs:
            original_size = original_shape.text_frame.paragraphs[0].runs[0].font.size
        # Use texto curto para nao acionar auto-shrink
        replacements = json.dumps({"0": {"10": "Oi"}})
        create_from_template_pptx(
            template_path=TEMPLATE_PATH,
            output_path=output_path,
            slide_replacements=replacements
        )
        prs_new = Presentation(output_path)
        new_shape = prs_new.slides[0].shapes[10]
        self.assertEqual(new_shape.text_frame.text, "Oi")
        if original_size:
            new_size = new_shape.text_frame.paragraphs[0].runs[0].font.size
            # Texto curto: tamanho deve ser exatamente o mesmo
            self.assertEqual(new_size, original_size)
        # Tambem verificar que outras formatacoes foram preservadas
        # (bold, font family, color)
        if original_shape.text_frame.paragraphs[0].runs:
            orig_bold = original_shape.text_frame.paragraphs[0].runs[0].font.bold
            new_bold = new_shape.text_frame.paragraphs[0].runs[0].font.bold
            self.assertEqual(orig_bold, new_bold)

    def test_create_from_template_multiline(self):
        if not os.path.exists(TEMPLATE_PATH):
            self.skipTest(f"Template nao encontrado: {TEMPLATE_PATH}")
        output_path = os.path.join(self.test_dir, "test_multiline.pptx")
        replacements = json.dumps({"0": {"10": "Linha 1\\nLinha 2\\nLinha 3"}})
        result = create_from_template_pptx(
            template_path=TEMPLATE_PATH,
            output_path=output_path,
            slide_replacements=replacements
        )
        self.assertIn("sucesso", result.lower())
        from pptx import Presentation
        prs = Presentation(output_path)
        text = prs.slides[0].shapes[10].text_frame.text
        self.assertIn("Linha 1", text)
        self.assertIn("Linha 2", text)
        self.assertIn("Linha 3", text)

    def test_create_from_template_invalid_json(self):
        if not os.path.exists(TEMPLATE_PATH):
            self.skipTest(f"Template nao encontrado: {TEMPLATE_PATH}")
        output_path = os.path.join(self.test_dir, "test_invalid.pptx")
        result = create_from_template_pptx(
            template_path=TEMPLATE_PATH,
            output_path=output_path,
            slide_replacements="not valid json"
        )
        self.assertIn("Erro", result)
        self.assertIn("JSON", result)

    def test_create_from_template_ignores_out_of_range(self):
        if not os.path.exists(TEMPLATE_PATH):
            self.skipTest(f"Template nao encontrado: {TEMPLATE_PATH}")
        output_path = os.path.join(self.test_dir, "test_oor.pptx")
        replacements = json.dumps({
            "999": {"0": "out of range"},
            "0": {"999": "shape out of range"}
        })
        result = create_from_template_pptx(
            template_path=TEMPLATE_PATH,
            output_path=output_path,
            slide_replacements=replacements
        )
        self.assertIn("sucesso", result.lower())
        self.assertIn("Substituicoes aplicadas: 0", result)

    def test_create_from_template_autoshrinks_long_text(self):
        """Testa que texto longo aciona auto-shrink e fica menor."""
        if not os.path.exists(TEMPLATE_PATH):
            self.skipTest(f"Template nao encontrado: {TEMPLATE_PATH}")
        output_path = os.path.join(self.test_dir, "test_autoshrink.pptx")
        # Texto longo que nao cabe no shape 10 do slide 0 (Author, 3.61 in)
        replacements = json.dumps({"0": {"10": "Texto Muito Longo Para Caber"}})
        result = create_from_template_pptx(
            template_path=TEMPLATE_PATH,
            output_path=output_path,
            slide_replacements=replacements
        )
        self.assertIn("sucesso", result.lower())
        from pptx import Presentation
        prs = Presentation(output_path)
        new_shape = prs.slides[0].shapes[10]
        self.assertEqual(new_shape.text_frame.text, "Texto Muito Longo Para Caber")
        # Auto-shrink deve ter reduzido a fonte (original era 32pt)
        if new_shape.text_frame.paragraphs[0].runs:
            new_size = new_shape.text_frame.paragraphs[0].runs[0].font.size.pt
            self.assertLess(new_size, 32.0, f"Auto-shrink nao reduziu a fonte: {new_size}pt")


if __name__ == "__main__":
    unittest.main(verbosity=2)
