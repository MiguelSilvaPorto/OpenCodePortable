"""build_cars_pptx.py — Gera PPTX de Indústria Automobilística Global

Aplicando DESIGN.md:
- Tipografia: Inter (título bold, corpo regular)
- Paleta: branco + slate + vermelho/blue (corporate analytical)
- Layout: 5 canônicos, sem decoração supérflua
- Espaçamento: margins 0.6in, grid de terços
- 12 slides, ~20min

Uso: python build_cars_pptx.py
Saída: ./output/cars_industry.pptx + previews HTML
"""
import tomllib
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).parent
THEME_PATH = ROOT / "themes" / "cars_industry.toml"
OUT_DIR = ROOT / "output"
PREVIEW_DIR = ROOT / "preview" / "cars"
OUT_DIR.mkdir(exist_ok=True)
PREVIEW_DIR.mkdir(parents=True, exist_ok=True)


def load_theme() -> dict:
    with open(THEME_PATH, "rb") as f:
        return tomllib.load(f)


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_textbox(
    slide,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    text: str,
    *,
    font_size: int,
    font_weight: str = "400",
    color_hex: str = "#0F172A",
    align: str = "left",
    anchor: str = "top",
    font_family: str = "Inter",
    line_spacing: float = 1.15,
):
    """Adiciona textbox com tipografia consistente."""
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)

    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.line_spacing = line_spacing
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]

    run = p.add_run()
    run.text = text
    run.font.name = font_family
    run.font.size = Pt(font_size)
    run.font.bold = font_weight in ("600", "700", "bold")
    run.font.color.rgb = hex_to_rgb(color_hex)
    return tb


def add_footer(slide, n: int, total: int, theme: dict):
    """Rodapé consistente: tema à esquerda, paginação à direita."""
    f = theme["footer"]
    fs = theme["font_sizes"]["footer"]
    color = theme["colors"]["text_muted"]
    add_textbox(
        slide, 0.6, 7.15, 8, 0.3,
        f["left"], font_size=fs, color_hex=color, align="left"
    )
    add_textbox(
        slide, 11.5, 7.15, 1.2, 0.3,
        f["right_format"].format(n=n, total=total),
        font_size=fs, color_hex=color, align="right"
    )


def slide_blank(prs, theme: dict):
    """Slide em branco 16:9 com fundo da paleta."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme["colors"]["background"])
    return slide


def slide_01_cover(prs, theme: dict, n: int, total: int):
    """S1 — Capa corporativa.

    Layout:
    - Acento vertical vermelho (1cm) à esquerda
    - Eyebrow text pequeno no topo
    - Título grande bold
    - Subtítulo médio
    - Rodapé com autor/ano
    """
    s = slide_blank(prs, theme)
    colors = theme["colors"]
    fs = theme["font_sizes"]

    accent_bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0),
        Inches(0.25), Inches(7.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(colors["accent"])
    accent_bar.line.fill.background()

    add_textbox(
        s, 0.9, 1.0, 11, 0.4,
        "RELATÓRIO SETORIAL  ·  2024-2025",
        font_size=fs["caption"] + 1, font_weight="600",
        color_hex=colors["accent"], align="left",
    )

    add_textbox(
        s, 0.9, 1.6, 11.5, 2.4,
        "Indústria\nAutomobilística\nGlobal",
        font_size=fs["cover_title"], font_weight="700",
        color_hex=colors["text_primary"], align="left",
        line_spacing=1.05,
    )

    add_textbox(
        s, 0.9, 5.0, 11.5, 0.8,
        "Panorama de mercado, revolução elétrica e geopolítica",
        font_size=fs["cover_subtitle"], font_weight="400",
        color_hex=colors["text_secondary"], align="left",
    )

    div = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(6.0),
        Inches(1.5), Emu(20000)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = hex_to_rgb(colors["text_primary"])
    div.line.fill.background()

    add_textbox(
        s, 0.9, 6.2, 11.5, 0.4,
        "OpenCode Studio  ·  Análise Setorial",
        font_size=fs["caption"] + 1, font_weight="500",
        color_hex=colors["text_muted"], align="left",
    )

    add_footer(s, n, total, theme)


def slide_02_panorama(prs, theme: dict, n: int, total: int):
    """S2 — Panorama 2024: 3 big numbers.

    Layout:
    - Título + subtítulo no topo
    - Grid 3 colunas com big number + label + contexto
    - Fonte dos dados no rodapé
    """
    s = slide_blank(prs, theme)
    colors = theme["colors"]
    fs = theme["font_sizes"]

    add_textbox(
        s, 0.6, 0.45, 12, 0.5,
        "PANORAMA  ·  NÚMEROS-CHAVE",
        font_size=fs["caption"] + 1, font_weight="600",
        color_hex=colors["accent"], align="left",
    )

    add_textbox(
        s, 0.6, 0.95, 12, 0.8,
        "A indústria em três números",
        font_size=fs["slide_title"], font_weight="700",
        color_hex=colors["text_primary"], align="left",
    )

    add_textbox(
        s, 0.6, 1.85, 12, 0.5,
        "Produção, frota ativa e valor de mercado em 2024",
        font_size=fs["body"], font_weight="400",
        color_hex=colors["text_secondary"], align="left",
    )

    metrics = [
        ("90M", "veículos produzidos", "por ano no mundo", "OICA / ACEA"),
        ("1.5bi", "carros em circulação", "1 para cada 5 pessoas", "Hedges & Company"),
        ("3T", "em receita (US$)", "setor + cadeia", "BloombergNEF"),
    ]

    col_w = 4.0
    gap = 0.13
    total_w = col_w * 3 + gap * 2
    start_x = (13.333 - total_w) / 2
    top_y = 3.0

    for i, (number, label, context, source) in enumerate(metrics):
        x = start_x + i * (col_w + gap)
        add_textbox(
            s, x, top_y, col_w, 1.6,
            number, font_size=fs["big_number"], font_weight="700",
            color_hex=colors["accent"], align="left", line_spacing=1.0,
        )
        add_textbox(
            s, x, top_y + 1.7, col_w, 0.5,
            label, font_size=fs["body"], font_weight="600",
            color_hex=colors["text_primary"], align="left",
        )
        add_textbox(
            s, x, top_y + 2.2, col_w, 0.4,
            context, font_size=fs["body_small"], font_weight="400",
            color_hex=colors["text_secondary"], align="left",
        )
        add_textbox(
            s, x, top_y + 2.6, col_w, 0.3,
            f"Fonte: {source}", font_size=fs["caption"], font_weight="400",
            color_hex=colors["text_muted"], align="left",
        )

    add_footer(s, n, total, theme)


def slide_03_top_oem(prs, theme: dict, n: int, total: int):
    """S3 — Top fabricantes 2024: tabela limpa.

    Layout:
    - Título + subtítulo
    - Tabela: rank, fabricante, vendas, share
    - Linha de destaque (Toyota) em cor de accent
    - Fonte no rodapé
    """
    s = slide_blank(prs, theme)
    colors = theme["colors"]
    fs = theme["font_sizes"]

    add_textbox(
        s, 0.6, 0.45, 12, 0.5,
        "FABRICANTES  ·  RANKING 2024",
        font_size=fs["caption"] + 1, font_weight="600",
        color_hex=colors["accent"], align="left",
    )

    add_textbox(
        s, 0.6, 0.95, 12, 0.7,
        "Os 8 maiores fabricantes globais",
        font_size=fs["slide_title"], font_weight="700",
        color_hex=colors["text_primary"], align="left",
    )

    add_textbox(
        s, 0.6, 1.75, 12, 0.5,
        "Toyota lidera; BYD cresce 41% e ultrapassa Honda em vendas",
        font_size=fs["body"], font_weight="400",
        color_hex=colors["text_secondary"], align="left",
    )

    rows = [
        ("1", "Toyota",       "10.5",  "11.7%"),
        ("2", "Volkswagen",   " 9.0",  "10.0%"),
        ("3", "Hyundai-Kia",  " 7.2",  " 8.0%"),
        ("4", "Stellantis",   " 6.0",  " 6.7%"),
        ("5", "GM",           " 5.9",  " 6.6%"),
        ("6", "BYD",          " 4.3",  " 4.8%"),
        ("7", "Honda",        " 3.8",  " 4.2%"),
        ("8", "Ford",         " 3.5",  " 3.9%"),
    ]

    table_x = 0.6
    table_y = 2.55
    col_widths = [0.6, 4.5, 2.0, 2.0]
    col_x = [table_x]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)

    header_h = 0.4
    row_h = 0.38

    for j, header in enumerate(["#", "Fabricante", "Vendas (M)", "Share"]):
        add_textbox(
            s, col_x[j], table_y, col_widths[j], header_h,
            header, font_size=fs["body_small"], font_weight="700",
            color_hex=colors["text_secondary"], align="left" if j == 1 else "right",
        )

    underline = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(table_x), Inches(table_y + header_h),
        Inches(sum(col_widths)), Emu(15000)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = hex_to_rgb(colors["text_primary"])
    underline.line.fill.background()

    for i, row in enumerate(rows):
        y = table_y + header_h + 0.05 + i * row_h
        is_highlight = row[1] == "BYD"
        for j, val in enumerate(row):
            color = colors["accent"] if is_highlight and j > 0 else colors["text_primary"]
            weight = "700" if is_highlight and j > 0 else ("400")
            add_textbox(
                s, col_x[j], y, col_widths[j], row_h,
                val, font_size=16, font_weight=weight,
                color_hex=color,
                align="left" if j == 1 else "right",
            )

    add_textbox(
        s, 0.6, 6.85, 12, 0.3,
        "Fonte: Focus2Move / OICA 2024 (estimativa anual)",
        font_size=fs["caption"], font_weight="400",
        color_hex=colors["text_muted"], align="left",
    )

    add_footer(s, n, total, theme)


def slide_04_regiao(prs, theme: dict, n: int, total: int):
    """S4 — Produção por região: barras horizontais simples (desenhadas).

    Layout:
    - Título + subtítulo
    - Lista de regiões com barra de share proporcional + número
    - 1 região em destaque (China)
    """
    s = slide_blank(prs, theme)
    colors = theme["colors"]
    fs = theme["font_sizes"]

    add_textbox(
        s, 0.6, 0.45, 12, 0.5,
        "GEOGRAFIA  ·  PRODUÇÃO GLOBAL",
        font_size=fs["caption"] + 1, font_weight="600",
        color_hex=colors["accent"], align="left",
    )

    add_textbox(
        s, 0.6, 0.95, 12, 0.8,
        "China fabrica 1 em cada 3 carros do mundo",
        font_size=fs["slide_title"], font_weight="700",
        color_hex=colors["text_primary"], align="left",
    )

    add_textbox(
        s, 0.6, 1.85, 12, 0.5,
        "Produção de veículos por região, 2024 (milhões de unidades / share global)",
        font_size=fs["body"], font_weight="400",
        color_hex=colors["text_secondary"], align="left",
    )

    regions = [
        ("China",         "29.0M",  32,  True),
        ("Europa",        "17.0M",  19,  False),
        ("América Norte", "15.5M",  17,  False),
        ("Japão + Coreia","12.0M",  13,  False),
        ("Índia",          "5.5M",   6,  False),
        ("Brasil + Mercosul", "4.0M",  4, False),
        ("Outros",         "7.0M",   9,  False),
    ]

    bar_max_w = 7.0
    bar_x = 4.0
    bar_y = 2.9
    row_h = 0.5

    for i, (name, value, share, highlight) in enumerate(regions):
        y = bar_y + i * row_h
        add_textbox(
            s, 0.6, y, 3.3, 0.4,
            name, font_size=fs["body"], font_weight="600" if highlight else "400",
            color_hex=colors["text_primary"], align="left",
        )
        bar_w = bar_max_w * (share / 32.0)
        bar = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(bar_x), Inches(y + 0.05),
            Inches(bar_w), Inches(0.30)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(colors["accent"] if highlight else colors["text_muted"])
        bar.line.fill.background()
        add_textbox(
            s, bar_x + bar_w + 0.1, y, 1.5, 0.4,
            f"{value}  ({share}%)", font_size=fs["body_small"],
            font_weight="700" if highlight else "400",
            color_hex=colors["text_primary"], align="left",
        )

    add_textbox(
        s, 0.6, 6.85, 12, 0.3,
        "Fonte: OICA 2024 — produção total estimada em 90M unidades",
        font_size=fs["caption"], font_weight="400",
        color_hex=colors["text_muted"], align="left",
    )

    add_footer(s, n, total, theme)


def slide_05_ev_growth(prs, theme: dict, n: int, total: int):
    """S5 — % EV nas vendas: linha do tempo 2018-2024.

    Layout:
    - Título + subtítulo
    - 6 anos como cards: 2019, 2020, 2021, 2022, 2023, 2024
    - Cada card: ano + % EV + variação
    - 2024 destacado
    """
    s = slide_blank(prs, theme)
    colors = theme["colors"]
    fs = theme["font_sizes"]

    add_textbox(
        s, 0.6, 0.45, 12, 0.5,
        "REVOLUÇÃO ELÉTRICA  ·  EV NAS VENDAS GLOBAIS",
        font_size=fs["caption"] + 1, font_weight="600",
        color_hex=colors["accent"], align="left",
    )

    add_textbox(
        s, 0.6, 0.95, 12, 0.8,
        "EVs saltam de 2% para 22% em 6 anos",
        font_size=fs["slide_title"], font_weight="700",
        color_hex=colors["text_primary"], align="left",
    )

    add_textbox(
        s, 0.6, 1.85, 12, 0.5,
        "Participação de EVs (BEV + PHEV) nas vendas globais de carros novos",
        font_size=fs["body"], font_weight="400",
        color_hex=colors["text_secondary"], align="left",
    )

    data = [
        (2019, "2.5%"),
        (2020, "4.2%"),
        (2021, "8.9%"),
        (2022, "13.7%"),
        (2023, "17.5%"),
        (2024, "22.0%"),
    ]

    card_w = 1.85
    card_h = 3.2
    gap = 0.1
    total_w = card_w * 6 + gap * 5
    start_x = (13.333 - total_w) / 2
    top_y = 3.0

    for i, (year, pct) in enumerate(data):
        x = start_x + i * (card_w + gap)
        is_last = i == len(data) - 1
        fill_color = colors["accent"] if is_last else colors["surface"]
        text_color = "#FFFFFF" if is_last else colors["text_primary"]
        sub_color = "#FFFFFF" if is_last else colors["text_secondary"]

        card = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(top_y),
            Inches(card_w), Inches(card_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(fill_color)
        card.line.color.rgb = hex_to_rgb(colors["divider"])
        card.line.width = Pt(0.5)

        add_textbox(
            s, x, top_y + 0.3, card_w, 0.5,
            str(year), font_size=fs["body"], font_weight="600",
            color_hex=sub_color, align="center",
        )
        add_textbox(
            s, x, top_y + 1.1, card_w, 1.4,
            pct, font_size=44, font_weight="700",
            color_hex=text_color, align="center", line_spacing=1.0,
        )
        if i > 0:
            prev_pct = float(data[i - 1][1].rstrip("%"))
            curr_pct = float(pct.rstrip("%"))
            growth = ((curr_pct - prev_pct) / prev_pct) * 100
            add_textbox(
                s, x, top_y + 2.55, card_w, 0.4,
                f"+{growth:.0f}% a.a.", font_size=fs["caption"],
                font_weight="600", color_hex=sub_color, align="center",
            )

    add_textbox(
        s, 0.6, 6.85, 12, 0.3,
        "Fonte: IEA Global EV Outlook 2024 / BloombergNEF",
        font_size=fs["caption"], font_weight="400",
        color_hex=colors["text_muted"], align="left",
    )

    add_footer(s, n, total, theme)


def main():
    theme = load_theme()
    prs = Presentation()
    prs.slide_width = Inches(theme["layout"]["slide_width_in"])
    prs.slide_height = Inches(theme["layout"]["slide_height_in"])

    total = 12
    slide_01_cover(prs, theme, 1, total)
    slide_02_panorama(prs, theme, 2, total)
    slide_03_top_oem(prs, theme, 3, total)
    slide_04_regiao(prs, theme, 4, total)
    slide_05_ev_growth(prs, theme, 5, total)

    out_path = OUT_DIR / "cars_industry.pptx"
    prs.save(out_path)
    print(f"OK: {out_path} ({out_path.stat().st_size / 1024:.1f} KB)")


if __name__ == "__main__":
    main()
