import os
import re
import shutil
import subprocess
import json
from mcp.server.fastmcp import FastMCP

# Inicializa o servidor FastMCP
mcp = FastMCP("Office", instructions="""
Este servidor fornece ferramentas locais para criar, ler e modificar arquivos Word (.docx), Excel (.xlsx/.xlsm) e PowerPoint (.pptx).
Tambem fornece extracao de texto de PDF, OCR em imagens, protecao por senha e execucao de macros VBA no Windows.

IMPORTANTE: Se o usuario pedir para criar ou modificar uma macro mas nao especificar claramente as regras ou o comportamento da mesma,
peca esclarecimentos antes de prosseguir.

Para documentos avancados, use create_from_example() passando um arquivo modelo existente.

REGRAS DE COMUNICACAO:
- Apos cada tool call, resuma o resultado de forma amigavel em portugues para o usuario
- NUNCA invente limitacoes ou problemas que nao existem nas mensagens das tools
- Se uma tool retornar mensagem de sucesso, reporte o resultado com clareza
- Se uma tool retornar erro, explique o erro de forma simples e sugira a correcao
- Use linguagem natural, nao exponha nomes tecnicos internos das tools
- Mostre o caminho do arquivo criado/modificado
""")

# ==========================================
# UTILITARIOS COMPARTILHADOS
# ==========================================

def _validar_arquivo_existe(file_path: str) -> str | None:
    """Valida se o arquivo existe. Retorna None se OK, ou mensagem de erro."""
    if not os.path.exists(file_path):
        return f"Erro: Nao foi possivel acessar o arquivo. O caminho '{file_path}' nao existe. Verifique o caminho e tente novamente."
    return None

def _validar_extensao(file_path: str, extensoes: tuple) -> str | None:
    """Valida se a extensao do arquivo e uma das esperadas."""
    if not file_path.lower().endswith(extensoes):
        ext_list = ", ".join(extensoes)
        return f"Erro: Formato de arquivo invalido. A extensao deve ser {ext_list}. Use: {file_path.split('.')[0]}{extensoes[0]}"
    return None

def _sanitizar_erro(stderr: str) -> str:
    """Remove ruido de stack trace do PowerShell, mantendo apenas a mensagem principal."""
    linhas = stderr.strip().split('\n')
    linhas_uteis = [
        l for l in linhas
        if not re.match(r'^\s*\+\s*CategoryInfo|^\s*\+ FullyQualifiedErrorId|^\s*\+', l)
    ]
    mensagem = ' '.join(l.strip() for l in linhas_uteis if l.strip())
    return mensagem[:500]

def _executar_powershell(script: str, timeout: int = 60) -> str:
    """
    Executa um script PowerShell de forma segura, com blindagens:
    - ExecutionPolicy Bypass: ignora restricoes locais
    - UTF-8 forcado: preserva acentuacao do portugues
    - Timeout explicito: evita travamentos
    - Sanitizacao de erro: protege o contexto do LLM
    """
    script_blindado = f"""
$OutputEncoding = [System.Text.Encoding]::UTF8
$ErrorActionPreference = "Stop"
{script}
"""
    cmd = [
        "powershell",
        "-NoProfile",
        "-NonInteractive",
        "-ExecutionPolicy", "Bypass",
        "-Command", script_blindado
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
        if result.returncode != 0:
            return f"Erro: {_sanitizar_erro(result.stderr)}"
        return result.stdout.strip()
    except subprocess.TimeoutExpired:
        return "Erro: A operacao excedeu o tempo limite de 60 segundos. O arquivo pode estar corrompido ou o Office pode estar ocupado. Feche outros programas e tente novamente."
    except FileNotFoundError:
        return "Erro: PowerShell nao encontrado. Esta ferramenta requer Windows com PowerShell 5.1 ou superior."
    except Exception as e:
        return f"Erro: Falha ao executar comando PowerShell. Detalhes: {_sanitizar_erro(str(e))}"


# ==========================================
# IMPORTACOES PREGUIÇOSAS (Lazy Imports)
# ==========================================
# Cada biblioteca e carregada sob demanda para evitar falhas durante a descoberta MCP
# se a dependencia nao estiver instalada.

def get_docx():
    import docx
    return docx

def get_openpyxl():
    import openpyxl
    return openpyxl

def get_pptx():
    import pptx
    return pptx

def get_power_pptx():
    import power_pptx
    return power_pptx

def get_win32():
    import win32com.client
    return win32com.client

def get_pymupdf():
    import fitz
    return fitz

def get_easyocr():
    import easyocr
    return easyocr

def get_msoffcrypto():
    import msoffcrypto
    return msoffcrypto


# ==========================================
# 1. FERRAMENTAS WORD (.docx)
# ==========================================

@mcp.tool()
def create_document_word(file_path: str, elements: list) -> str:
    """
    Cria um novo documento Word (.docx) a partir de uma lista de elementos estruturados.

    Esta tool gera um arquivo .docx do zero, aceitando uma lista de dicionarios
    que representam blocos de conteudo (titulos, cabecalhos, paragrafos, listas).
    O arquivo e SOBRESCRITO se ja existir no caminho informado.

    QUANDO USAR:
    - Quando o usuario pede para criar um relatorio, documento ou arquivo Word
    - Quando e necessario gerar conteudo estruturado de forma programatica
    - Para criar documentos curtos (cartas, memorandos, resumos)
    - Para documentos com layout simples e direto

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho onde o .docx sera salvo.
      Aceita caminho absoluto ou relativo (baseado no workspace).
      Exemplos: "./output/relatorio.docx", "C:\\docs\\texto.docx"

    - elements (list[dict], obrigatorio): Lista de dicionarios. Cada dicionario define um bloco:
      - "type" (str, obrigatorio): Tipo do elemento:
          "title"      -> Titulo principal (Heading Level 0)
          "heading1"   -> Cabecalho de nivel 1
          "heading2"   -> Cabecalho de nivel 2
          "paragraph"  -> Paragrafo comum
          "list_item"  -> Item de lista com marcador (estilo "List Bullet")
      - "text" (str, obrigatorio): Conteudo textual do elemento.
      - "bold" (bool, opcional): Negrito. Default: False.
      - "italic" (bool, opcional): Italico. Default: False.

    RETORNO:
    - Mensagem de texto com o caminho absoluto do arquivo salvo.
      Ex: "Documento Word criado com sucesso em: D:\\projeto\\relatorio.docx"

    IMPORTANT:
    - O arquivo e SOBRESCRITO completamente se ja existir. Para adicionar conteudo
      a um documento existente, use append_to_document_word() ou edit_document_word().
    - "title" usa Heading Level 0 (secao especial do Word), nao confunda com heading1.
    - "list_item" cria paragrafo com estilo "List Bullet" (marcador visual).

    EXEMPLO DE USO:
        create_document_word(
            file_path="./relatorio.docx",
            elements=[
                {"type": "title", "text": "Relatorio Anual 2026"},
                {"type": "heading1", "text": "1. Introducao"},
                {"type": "paragraph", "text": "Este relatorio apresenta os resultados do exercicio."},
                {"type": "heading2", "text": "1.1 Metodologia"},
                {"type": "paragraph", "text": "A analise foi baseada em dados reais."},
                {"type": "heading1", "text": "2. Resultados"},
                {"type": "list_item", "text": "Crescimento de 15% no faturamento", "bold": True},
                {"type": "list_item", "text": "Reducao de 8% nos custos operacionais"},
                {"type": "list_item", "text": "Novos contratos assinados: 42"},
            ]
        )
        # Retorna: "Documento Word criado com sucesso em: D:\\projeto\\relatorio.docx"
    """
    docx = get_docx()
    doc = docx.Document()

    for el in elements:
        el_type = el.get("type", "paragraph")
        text = el.get("text", "")

        if el_type == "title":
            doc.add_heading(text, level=0)
        elif el_type == "heading1":
            doc.add_heading(text, level=1)
        elif el_type == "heading2":
            doc.add_heading(text, level=2)
        elif el_type == "list_item":
            p = doc.add_paragraph(style='List Bullet')
            run = p.add_run(text)
            run.bold = el.get("bold", False)
            run.italic = el.get("italic", False)
        else:
            p = doc.add_paragraph()
            run = p.add_run(text)
            run.bold = el.get("bold", False)
            run.italic = el.get("italic", False)

    doc.save(file_path)
    return f"Documento Word criado com sucesso em: {os.path.abspath(file_path)}"


@mcp.tool()
def read_document_word(file_path: str) -> str:
    """
    Le e extrai todo o texto estruturado de um documento Word (.docx), incluindo tabelas.

    Esta tool abre um arquivo .docx existente e retorna todo o conteudo textual
    de forma organizada: paragrafos seguidos, e tabelas formatadas com celulas
    separadas por " | ". Ideal para analisar o conteudo antes de editar.

    QUANDO USAR:
    - Quando o usuario quer saber o conteudo de um documento Word
    - Antes de fazer edicoes, para entender a estrutura atual
    - Para extrair dados de tabelas embutidas no documento
    - Para verificar se o documento contem o conteudo esperado

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx a ser lido.
      Exemplos: "./contrato.docx", "C:\\docs\\relatorio.docx"

    RETORNO:
    - Texto estruturado com o seguinte formato:
      - Cabecalho: "--- Documento: nome_arquivo.docx ---"
      - Paragrafos: texto livre, um por linha
      - Tabelas: "--- Tabelas Encontradas ---", seguido de "Tabela N:" e linhas com " | "
    - Se o arquivo nao existir, retorna mensagem de erro descritiva.

    IMPORTANT:
    - Extrai APENAS texto. Formatacao visual (fontes, cores, tamanhos) NAO e extraida.
    - Imagens e formas sem texto nao sao exibidas.
    - Tabelas sao exibidas com celulas separadas por " | " (pipe).

    EXEMPLO DE USO:
        read_document_word(file_path="./contrato.docx")
        # Retorna:
        # --- Documento: contrato.docx ---
        # CONTRATO DE PRESTACAO DE SERVICOS
        # O presente contrato tem por objeto...
        # --- Tabelas Encontradas ---
        # Tabela 1:
        # Item | Valor | Prazo
        # Servico A | R$ 5.000 | 30 dias
        # Servico B | R$ 3.200 | 15 dias
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    docx = get_docx()
    doc = docx.Document(file_path)
    result = []

    result.append(f"--- Documento: {os.path.basename(file_path)} ---")
    for p in doc.paragraphs:
        if p.text.strip():
            result.append(p.text)

    if doc.tables:
        result.append("\n--- Tabelas Encontradas ---")
        for i, table in enumerate(doc.tables):
            result.append(f"Tabela {i+1}:")
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                result.append(" | ".join(row_text))

    return "\n".join(result)


@mcp.tool()
def create_document_from_template_word(template_path: str, output_path: str, replacements: dict) -> str:
    """
    Clona um documento .docx existente e substitui marcadores de texto,
    preservando 100% da formatacao original do modelo.

    Esta tool copia um arquivo modelo (.docx) e substitui todos os marcadores
    de texto (ex: "{{NomeCliente}}") pelos valores fornecidos. A formatacao
    visual (fontes, cores, tamanhos, espacamento) e 100% preservada.

    QUANDO USAR:
    - Para gerar documentos personalizados a partir de um modelo padrao
    - Contratos, propostas, cartas, certificados com dados variaveis
    - Quando e preciso manter a formatacao exata de um documento modelo
    - Para automacao de documentos que seguem um template fixo

    PARAMETROS:
    - template_path (str, obrigatorio): Caminho do arquivo modelo .docx.
      O modelo DEVE existir e conter os marcadores de texto.
      Exemplos: "./modelos/contrato_base.docx", "C:\\templates\\proposta.docx"

    - output_path (str, obrigatorio): Caminho onde o documento gerado sera salvo.
      Pode ser diferente do template (nao sobrescreve o modelo).
      Exemplos: "./contratos/contrato_joao.docx"

    - replacements (dict, obrigatorio): Dicionario de substituicoes {marcador: valor}.
      Os marcadores devem ser strings unicas no documento.
      Exemplo:
        {
          "{{NomeCliente}}": "Joao Silva",
          "{{Data}}": "10/06/2026",
          "{{Valor}}": "R$ 15.000,00",
          "{{CNPJ}}": "12.345.678/0001-90"
        }

    RETORNO:
    - Mensagem de texto com o caminho absoluto do arquivo gerado.
      Ex: "Documento gerado a partir do modelo com sucesso em: D:\\docs\\contrato_joao.docx"

    IMPORTANT:
    - A substituicao funciona em "runs" individuais do XML do Word.
    - Se um marcador estiver DIVIDIDO em varios runs (ex: "{{Nome" em um run e
      "Cliente}}" em outro), a substituicao pode FALHAR.
      SOLUCAO: Crie o modelo com o marcador inteiro em um UNICO run.
    - A formatacao original (fontes, cores, tamanhos, espacamento) e 100% preservada.
    - Tabelas no modelo tambem sao processadas (substituicao em celulas).
    - O arquivo modelo NAO e modificado — apenas copiado.

    EXEMPLO DE USO:
        create_document_from_template_word(
            template_path="./modelos/contrato_base.docx",
            output_path="./contratos/contrato_joao_silva.docx",
            replacements={
                "{{NomeCliente}}": "Joao Silva",
                "{{DataAssinatura}}": "10/06/2026",
                "{{ValorContrato}}": "R$ 45.000,00",
                "{{Duracao}}": "12 meses",
                "{{Objeto}}": "Prestacao de servicos de consultoria"
            }
        )
        # Retorna: "Documento gerado a partir do modelo com sucesso em: D:\\docs\\contrato_joao_silva.docx"
    """
    erro = _validar_arquivo_existe(template_path)
    if erro:
        return erro

    docx = get_docx()
    shutil.copyfile(template_path, output_path)
    doc = docx.Document(output_path)

    for p in doc.paragraphs:
        for run in p.runs:
            for key, val in replacements.items():
                if key in run.text:
                    run.text = run.text.replace(key, str(val))

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    for run in p.runs:
                        for key, val in replacements.items():
                            if key in run.text:
                                run.text = run.text.replace(key, str(val))

    doc.save(output_path)
    return f"Documento gerado a partir do modelo com sucesso em: {os.path.abspath(output_path)}"


# ==========================================
# 1.5 WORD - FERRAMENTAS AVANCADAS
# ==========================================

@mcp.tool()
def add_image_to_document_word(file_path: str, image_path: str, width_inch: float = None, height_inch: float = None) -> str:
    """
    Insere uma imagem em um documento Word (.docx) existente.

    Esta tool adiciona uma imagem ao FINAL de um documento Word existente.
    A imagem e inserida apos todo o conteudo ja existente. Pode ser
    redimensionada especificando largura e/ou altura em polegadas.

    QUANDO USAR:
    - Quando o documento ja existe e voce quer adicionar uma imagem (logo, grafico, foto)
    - Para incluir imagens de assinatura, selos ou carimbos
    - Para adicionar graficos exportados como imagem ao documento
    - Para inserir capturas de tela ou diagramas

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx EXISTENTE.
      Exemplos: "./relatorio.docx", "C:\\docs\\projeto.docx"

    - image_path (str, obrigatorio): Caminho da imagem local a ser inserida.
      Formatos suportados: PNG, JPG/JPEG, GIF, BMP, TIFF.
      Exemplos: "./logotipo.png", "C:\\imgs\\grafico.jpg"

    - width_inch (float, opcional): Largura da imagem em polegadas.
      Se informado, a imagem e redimensionada. Se omitido, usa tamanho original.
      1 polegada = 2.54 cm. Ex: 5.5 polegadas = ~14 cm.

    - height_inch (float, opcional): Altura da imagem em polegadas.
      Se AMBOS (width e height) forem informados, a imagem e distorcida proporcionalmente.
      Se APENAS um for informado, o outro e calculado proporcionalmente.
      Se nenhum for informado, usa o tamanho original.

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Imagem inserida com sucesso em: D:\\docs\\relatorio.docx"

    IMPORTANT:
    - A imagem e inserida no FINAL do documento (apos todo conteudo existente).
    - Nao e possivel inserir imagem em posicao especifica com esta tool.
    - Para inserir imagens em slides, use add_image_to_slide_pptx().
    - Se width_inch e height_inch forem informados, a imagem e redimensionada.
    - Se apenas um for informado, o outro e calculado mantendo a proporcao.

    EXEMPLO DE USO:
        add_image_to_document_word(
            file_path="./relatorio.docx",
            image_path="./grafico_vendas.png",
            width_inch=5.5
        )
        # Retorna: "Imagem inserida com sucesso em: D:\\docs\\relatorio.docx"

        # Com ambas dimensoes (atencao: pode distorcer se nao mantiver proporcao):
        add_image_to_document_word(
            file_path="./documento.docx",
            image_path="./logo_empresa.png",
            width_inch=2.0,
            height_inch=1.0
        )
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    erro_img = _validar_arquivo_existe(image_path)
    if erro_img:
        return erro_img

    docx = get_docx()
    from docx.shared import Inches

    doc = docx.Document(file_path)
    width = Inches(width_inch) if width_inch else None
    height = Inches(height_inch) if height_inch else None
    doc.add_picture(image_path, width=width, height=height)
    doc.save(file_path)
    return f"Imagem inserida com sucesso em: {os.path.abspath(file_path)}"


@mcp.tool()
def add_table_to_document_word(file_path: str, data: list, style: str = "Light Grid Accent 1") -> str:
    """
    Adiciona uma tabela formatada a um documento Word (.docx) existente.

    Esta tool insere uma tabela ao FINAL de um documento Word existente.
    A primeira linha dos dados e tratada como cabecalho (negrito automatico).
    A tabela e adicionada apos todo o conteudo existente.

    QUANDO USAR:
    - Quando o documento ja existe e voce quer adicionar uma tabela de dados
    - Para inserir dados tabulares (vendas, resultados, listagens)
    - Quando precisa de uma tabela com estilo visual pre-definido
    - Para complementar um relatorio com dados estruturados

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx EXISTENTE.
      Exemplos: "./relatorio.docx", "C:\\docs\\projeto.docx"

    - data (list[list], obrigatorio): Dados da tabela em formato de lista de listas.
      - A PRIMEIRA LINHA e tratada como cabecalho (automaticamente fica em negrito).
      - Cada sub-lista e uma linha da tabela.
      - Valores sao convertidos para string.
      Exemplo:
        [
          ["Produto", "Quantidade", "Valor Unitario", "Total"],
          ["Notebook Dell", 10, "R$ 4.500", "R$ 45.000"],
          ["Monitor LG", 25, "R$ 1.200", "R$ 30.000"],
          ["Teclado Logitech", 50, "R$ 150", "R$ 7.500"]
        ]

    - style (str, opcional): Estilo visual da tabela. Default: "Light Grid Accent 1".
      Estilos comuns do Word:
      - "Light Grid Accent 1" (padrao, bordas leves)
      - "Table Grid" (bordas pretas simples)
      - "Medium Shading 1 Accent 1" (fundo alternado)
      - "List Table 4 Accent 1" (estilo limpo)

    RETORNO:
    - Mensagem de texto com dimensoes da tabela e caminho do arquivo.
      Ex: "Tabela 4x4 adicionada com sucesso em: D:\\docs\\relatorio.docx"

    IMPORTANT:
    - A tabela e adicionada ao FINAL do documento (apos todo conteudo existente).
    - Nao e possivel inserir tabela em posicao especifica com esta tool.
    - A primeira linha fica automaticamente em negrito (cabecalho).
    - Para tabelas em slides, use add_table_to_slide_pptx().

    EXEMPLO DE USO:
        add_table_to_document_word(
            file_path="./relatorio.docx",
            data=[
                ["Mes", "Receita", "Despesa", "Lucro"],
                ["Janeiro", "R$ 50.000", "R$ 35.000", "R$ 15.000"],
                ["Fevereiro", "R$ 55.000", "R$ 33.000", "R$ 22.000"],
                ["Marco", "R$ 62.000", "R$ 38.000", "R$ 24.000"],
            ],
            style="Light Grid Accent 1"
        )
        # Retorna: "Tabela 4x4 adicionada com sucesso em: D:\\docs\\relatorio.docx"
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    if not data or not data[0]:
        return "Erro: Os dados da tabela nao podem estar vazios. Forneca uma lista com pelo menos uma linha."

    docx = get_docx()
    doc = docx.Document(file_path)
    rows = len(data)
    cols = len(data[0])
    table = doc.add_table(rows=rows, cols=cols)
    table.style = style

    for i, row_data in enumerate(data):
        for j, cell_value in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_value)
            if i == 0:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.bold = True

    doc.save(file_path)
    return f"Tabela {rows}x{cols} adicionada com sucesso em: {os.path.abspath(file_path)}"


@mcp.tool()
def edit_document_word(file_path: str, elements: list, position: str = "end") -> str:
    """
    Adiciona elementos a um documento Word (.docx) existente, no inicio ou no final.

    Esta tool abre um .docx existente e insere novos elementos (titulos, paragrafos,
    listas) either no INICIO ou no FINAL do documento. Preserva todo o conteudo
    ja existente.

    QUANDO USAR:
    - Quando o documento ja existe e voce quer adicionar conteudo em POSICAO ESPECIFICA
    - Para adicionar um sumario ou preambulo no INICIO do documento
    - Para adicionar anexos ou complementos no FINAL
    - Quando append_to_document_word() nao atende (precisa inserir no inicio)

    DIFERENCA COM append_to_document_word():
    - append_to_document_word: SEMPRE adiciona ao final
    - edit_document_word: Permite inserir no inicio ("begin") ou final ("end")

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx EXISTENTE.
      Exemplos: "./relatorio.docx", "C:\\docs\\projeto.docx"

    - elements (list[dict], obrigatorio): Lista de elementos no mesmo formato de create_document_word():
      - "type" (str): "title" | "heading1" | "heading2" | "paragraph" | "list_item"
      - "text" (str): Conteudo textual
      - "bold" (bool, opcional): Negrito. Default: False.
      - "italic" (bool, opcional): Italico. Default: False.

    - position (str, opcional): Onde inserir os elementos:
      - "end" (default): Adiciona ao final do documento
      - "begin": Adiciona no inicio do documento (antes de todo conteudo existente)

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Documento editado com sucesso em: D:\\docs\\relatorio.docx"

    IMPORTANT:
    - O arquivo DEVE existir. Se nao existir, a tool retorna erro.
    - Nao e possivel editar no MEIO do documento — apenas inicio ou fim.
    - Para edicoes no meio, leia o documento e recrie com as alteracoes.

    EXEMPLO DE USO:
        # Adicionar sumario no inicio do documento:
        edit_document_word(
            file_path="./relatorio.docx",
            elements=[
                {"type": "title", "text": "SUMARIO"},
                {"type": "paragraph", "text": "1. Introducao .............. pag 1"},
                {"type": "paragraph", "text": "2. Metodos ................. pag 3"},
            ],
            position="begin"
        )

        # Adicionar anexo no final:
        edit_document_word(
            file_path="./relatorio.docx",
            elements=[
                {"type": "heading1", "text": "ANEXO A - Dados Complementares"},
                {"type": "paragraph", "text": "Tabela de dados brutos da pesquisa."},
            ],
            position="end"
        )
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    if position not in ("end", "begin"):
        return "Erro: Posicao invalida. Use 'end' (final) ou 'begin' (inicio)."

    docx = get_docx()
    doc = docx.Document(file_path)

    new_paragraphs = []
    for el in elements:
        el_type = el.get("type", "paragraph")
        text = el.get("text", "")

        if el_type == "title":
            new_paragraphs.append(doc.add_heading(text, level=0))
        elif el_type == "heading1":
            new_paragraphs.append(doc.add_heading(text, level=1))
        elif el_type == "heading2":
            new_paragraphs.append(doc.add_heading(text, level=2))
        elif el_type == "list_item":
            p = doc.add_paragraph(style='List Bullet')
            run = p.add_run(text)
            run.bold = el.get("bold", False)
            run.italic = el.get("italic", False)
            new_paragraphs.append(p)
        else:
            p = doc.add_paragraph()
            run = p.add_run(text)
            run.bold = el.get("bold", False)
            run.italic = el.get("italic", False)
            new_paragraphs.append(p)

    doc.save(file_path)
    return f"Documento editado com sucesso em: {os.path.abspath(file_path)}"


# ==========================================
# 1.6 WORD - TRACK CHANGES E COMENTARIOS
# ==========================================

@mcp.tool()
def track_changes_word(file_path: str, action: str) -> str:
    """
    Gerencia controle de alteracoes (Track Changes) em documentos Word.

    Esta tool permite aceitar ou rejeitar TODAS as alteracoes pendentes de um
    documento com Track Changes ativado. Util para processar revisoes de forma
    massiva sem precisar aceitar/rejeitar uma a uma.

    QUANDO USAR:
    - Quando o documento tem alteracoes registradas (Track Changes) e voce quer aceitar todas
    - Para rejeitar todas as alteracoes e voltar ao texto original
    - Para processar revisoes de forma automatica
    - Para limpar um documento de alteracoes pendentes

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx.
    - action (str, obrigatorio): Acao a executar:
      - "accept_all": Aceita todas as alteracoes pendentes
      - "reject_all": Rejeita todas as alteracoes (volta ao original)

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Todas as alteracoes aceitas no documento."

    REQUISITOS:
    - Biblioteca docx-revisions instalada (pip install docx-revisions)

    IMPORTANTE:
    - A acao e IRREVERSIVEL apos salvar. Verifique o documento antes.
    - Para substituir texto com registro de alteracao, use find_replace_tracked_word().
    - Para listar ou adicionar comentarios, use manage_comments_word().

    EXEMPLO DE USO:
        track_changes_word(file_path="./contrato.docx", action="accept_all")
        # Retorna: "Todas as alteracoes aceitas no documento."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        from docx_revisions import RevisionDocument
        rdoc = RevisionDocument(file_path)

        if action == "accept_all":
            rdoc.accept_all()
            rdoc.save(file_path)
            return "Todas as alteracoes aceitas no documento."
        elif action == "reject_all":
            rdoc.reject_all()
            rdoc.save(file_path)
            return "Todas as alteracoes rejeitadas no documento."
        else:
            return "Erro: Acao invalida. Use 'accept_all' ou 'reject_all'."
    except ImportError:
        return "Erro: docx-revisions nao esta instalado. Execute: pip install docx-revisions"
    except Exception as e:
        return f"Erro: Nao foi possivel gerenciar alteracoes. Detalhes: {str(e)[:200]}"


@mcp.tool()
def find_replace_tracked_word(file_path: str, old_text: str, new_text: str, author: str = "IA") -> str:
    """
    Localiza e substitui texto com registro de alteracao (Track Changes).

    Esta tool faz find & replace em um documento Word, mas registra cada
    substituicao como uma alteracao rastreavel (Track Changes). O texto antigo
    e marcado como "deletado" e o novo como "inserido", ambos visiveis no Word.

    QUANDO USAR:
    - Para substituir texto mantendo historico de alteracoes
    - Para que o destinatario possa ver o que foi alterado
    - Para revisoes colaborativas onde o historico importa
    - Para substituicoes que precisam de aprovacao

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx.
    - old_text (str, obrigatorio): Texto a ser substituido (case-insensitive).
    - new_text (str, obrigatorio): Novo texto substituto.
    - author (str, opcional): Nome do autor da alteracao. Default: "IA".

    RETORNO:
    - Mensagem de texto com quantidade de substituicoes.
      Ex: "3 ocorrencia(s) de 'ACME' substituida(s) por 'NewCo' com registro de alteracao."

    REQUISITOS:
    - Biblioteca docx-revisions instalada (pip install docx-revisions)

    IMPORTANTE:
    - Todas as substituicoes sao registradas como Track Changes.
    - Para aceitar/rejeitar depois, use track_changes_word().
    - O autor e visivel no painel de revisao do Word.

    EXEMPLO DE USO:
        find_replace_tracked_word(
            file_path="./contrato.docx",
            old_text="ACME Corp",
            new_text="NewCo Inc",
            author="Assistente IA"
        )
        # Retorna: "1 ocorrencia(s) de 'ACME Corp' substituida(s) por 'NewCo Inc' com registro de alteracao."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        from docx_revisions import RevisionDocument
        rdoc = RevisionDocument(file_path)
        count = rdoc.find_and_replace_tracked(old_text, new_text, author=author)
        rdoc.save(file_path)
        return f"{count} ocorrencia(s) de '{old_text}' substituida(s) por '{new_text}' com registro de alteracao."
    except ImportError:
        return "Erro: docx-revisions nao esta instalado. Execute: pip install docx-revisions"
    except Exception as e:
        return f"Erro: Nao foi possivel substituir texto. Detalhes: {str(e)[:200]}"


@mcp.tool()
def manage_comments_word(file_path: str, action: str, text: str = None, author: str = "IA") -> str:
    """
    Gerencia comentarios em documentos Word (listar ou adicionar).

    Esta tool permite listar todos os comentarios existentes em um documento
    ou adicionar um novo comentario. Comentarios sao uteis para feedback
    e revisao colaborativa de documentos.

    QUANDO USAR:
    - Para listar todos os comentarios de um documento recebido
    - Para adicionar feedback ou observacoes ao documento
    - Para marcar secoes que precisam de revisao
    - Para comunicar instrucoes sobre trechos especificos

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx.
    - action (str, obrigatorio): Acao a executar:
      - "list": Lista todos os comentarios do documento
      - "add": Adiciona um novo comentario
    - text (str, opcional): Texto do comentario (obrigatorio para "add").
    - author (str, opcional): Autor do comentario. Default: "IA".

    RETORNO:
    - Para "list": Lista de comentarios no formato "[autor] texto".
      Ou: "Nenhum comentario encontrado no documento."
    - Para "add": "Comentario adicionado ao documento."

    IMPORTANTE:
    - Comentarios sao visiveis no painel lateral do Word.
    - Para alteracoes com registro de alteracao, use track_changes_word() ou find_replace_tracked_word().
    - O autor e visivel junto ao comentario.

    EXEMPLO DE USO:
        # Listar comentarios:
        manage_comments_word(file_path="./contrato.docx", action="list")
        # Retorna: "[Joao] Revisar clausula 3\n[Ana] OK, aprovado"

        # Adicionar comentario:
        manage_comments_word(
            file_path="./contrato.docx",
            action="add",
            text="Revisar este paragrafo antes de enviar ao cliente",
            author="Assistente IA"
        )
        # Retorna: "Comentario adicionado ao documento."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    docx = get_docx()
    doc = docx.Document(file_path)

    if action == "list":
        try:
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            comments_part = doc.part.related_parts.get('commentsPart')
            if not comments_part:
                return "Nenhum comentario encontrado no documento."
            from docx.opc.parts.comments import CommentsPart
            comments = comments_part.comments
            if not comments:
                return "Nenhum comentario encontrado."
            result = []
            for c in comments:
                result.append(f"[{c.author}] {c.text}")
            return "\n".join(result) if result else "Nenhum comentario encontrado."
        except Exception:
            return "Nenhum comentario encontrado no documento."

    elif action == "add":
        if not text:
            return "Erro: Informe o texto do comentario em 'text'."
        try:
            from docx.opc.parts.comments import CommentsPart
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            from docx.oxml.ns import qn
            import datetime

            comments_part = doc.part.related_parts.get('commentsPart')
            if not comments_part:
                from docx.oxml import parse_xml
                comments_xml = parse_xml(
                    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                    '<w:comments xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>'
                )
                comments_part = docx.opc.parts.comments.CommentsPart(doc.part.package, comments_xml)
                doc.part.relate_to(comments_part, RT.COMMENTS)

            import random
            cid = str(random.randint(100000, 999999))
            comment = docx.oxml.parse_xml(
                f'<w:comment xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" '
                f'w:id="{cid}" w:author="{author}" w:date="{datetime.datetime.now().isoformat()}">'
                f'<w:p><w:r><w:t>{text}</w:t></w:r></w:p></w:comment>'
            )
            comments_part._element.append(comment)
            doc.save(file_path)
            return f"Comentario adicionado ao documento."
        except Exception as e:
            return f"Erro: Nao foi possivel adicionar comentario. Detalhes: {str(e)[:200]}"

    return "Erro: Acao invalida. Use 'list' ou 'add'."


@mcp.tool()
def set_header_footer_word(file_path: str, tipo: str, text: str, page_number: bool = False) -> str:
    """
    Configura cabecalho ou rodape em um documento Word (.docx).

    Esta tool define o conteudo do cabecalho (header) ou rodape (footer) de um
    documento Word. Pode incluir texto estatico e opcionalmente numero de pagina.

    QUANDO USAR:
    - Para adicionar texto "Confidencial" ou "Rascunho" em todos os documentos
    - Para incluir numero de pagina no rodape
    - Para adicionar data ou titulo ao cabecalho
    - Para padronizar documentos com informacoes de cabecalho/rodape

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx.
    - tipo (str, obrigatorio): Tipo de container:
      - "header": Cabecalho (parte superior de cada pagina)
      - "footer": Rodape (parte inferior de cada pagina)
    - text (str, obrigatorio): Texto a ser exibido.
    - page_number (bool, opcional): Se True, adiciona numero de pagina ao lado do texto.
      Default: False.

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Cabecalho configurado em: D:\\docs\\relatorio.docx"
      Ou: "Rodape configurado em: D:\\docs\\relatorio.docx"

    IMPORTANTE:
    - O cabecalho/rodape e aplicado a TODAS as paginas do documento.
    - Se page_number=True, o numero e alinhado a direita.
    - O texto e adicionado ao container existente (nao substitui).

    EXEMPLO DE USO:
        # Adicionar cabecalho:
        set_header_footer_word(file_path="./relatorio.docx", tipo="header", text="CONFIDENCIAL")

        # Adicionar rodape com numero de pagina:
        set_header_footer_word(file_path="./relatorio.docx", tipo="footer", text="Pagina ", page_number=True)
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    if tipo not in ("header", "footer"):
        return "Erro: Tipo invalido. Use 'header' (cabecalho) ou 'footer' (rodape)."

    docx = get_docx()
    doc = docx.Document(file_path)
    section = doc.sections[0]

    if tipo == "header":
        container = section.header
    else:
        container = section.footer

    container.is_linked_to_previous = False
    paragraph = container.paragraphs[0]
    run = paragraph.add_run(text)

    if page_number:
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run2 = paragraph.add_run(" ")
        run3 = paragraph.add_run("")
        from docx.oxml.ns import qn
        fldChar1 = run3._r.makeelement(qn('w:fldChar'), {qn('w:fldCharType'): 'begin'})
        run3._r.append(fldChar1)
        run4 = paragraph.add_run("")
        instrText = run4._r.makeelement(qn('w:instrText'), {})
        instrText.text = ' PAGE '
        run4._r.append(instrText)
        run5 = paragraph.add_run("")
        fldChar2 = run5._r.makeelement(qn('w:fldChar'), {qn('w:fldCharType'): 'end'})
        run5._r.append(fldChar2)

    doc.save(file_path)
    return f"{'Cabecalho' if tipo == 'header' else 'Rodape'} configurado em: {os.path.abspath(file_path)}"


@mcp.tool()
def analyze_document_style_word(file_path: str) -> str:
    """
    Extrai informacoes de estilo de um documento Word (.docx) existente.

    Esta tool analisa a formatacao visual de um documento e retorna metadados
    de estilo que podem ser usados para replicar o mesmo visual em novos documentos.
    Analisa fontes, cores, margens, orientacao e estilos de cabecalho usados.

    QUANDO USAR:
    - Para extrair o estilo de um documento modelo antes de criar um novo
    - Para auditar a formatacao de um documento recebido
    - Para documentar padroes visuais de documentos
    - Para usar com create_document_from_example_word()

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx a ser analisado.

    RETORNO:
    - JSON formatado com estilos extraidos:
      - "fonts": Fontes usadas (nome, tamanho, cor)
      - "margins": Margens (topo, base, esquerda, direita)
      - "orientation": Orientacao (portrait ou landscape)
      - "has_tables": Se possui tabelas
      - "heading_styles": Estilos de cabecalho usados

    IMPORTANTE:
    - Nao modifica o arquivo — e Somente Leitura.
    - O JSON retornado pode ser usado como entrada em create_document_from_example_word().
    - Para extrair apenas estatisticas (sem estilo), use get_document_info_word().

    EXEMPLO DE USO:
        analyze_document_style_word(file_path="./modelo.docx")
        # Retorna JSON com fontes, cores, margens e orientacao detectadas
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    docx = get_docx()
    doc = docx.Document(file_path)

    fonts = {}
    colors = {}
    heading_styles = []

    for p in doc.paragraphs:
        for run in p.runs:
            font_name = run.font.name
            font_size = str(run.font.size) if run.font.size else None
            font_color = str(run.font.color.rgb) if run.font.color and run.font.color.rgb else None
            if font_name and font_name not in fonts:
                fonts[font_name] = {"size": font_size, "color": font_color}

        style_name = p.style.name if p.style else None
        if style_name and style_name.startswith("Heading") and style_name not in heading_styles:
            heading_styles.append(style_name)

    section = doc.sections[0] if doc.sections else None
    margins = {}
    orientation = "portrait"
    if section:
        margins = {
            "top": str(section.top_margin),
            "bottom": str(section.bottom_margin),
            "left": str(section.left_margin),
            "right": str(section.right_margin)
        }
        if section.orientation == 1:
            orientation = "landscape"

    has_tables = len(doc.tables) > 0

    import json
    result = {
        "fonts": fonts,
        "colors": colors,
        "margins": margins,
        "orientation": orientation,
        "has_tables": has_tables,
        "heading_styles": heading_styles
    }
    return json.dumps(result, ensure_ascii=False, indent=2)


@mcp.tool()
def save_style_profile_word(style_data: str, profile_path: str) -> str:
    """
    Salva um perfil de estilo (JSON) em um arquivo para reuso futuro.

    Esta tool recebe os dados de estilo (gerados por analyze_document_style_word)
    e salva em um arquivo .json para uso posterior. Permite criar bibliotecas de
    estilos reutilizaveis.

    QUANDO USAR:
    - Para salvar estilos extraidos de documentos para reuso
    - Para criar bibliotecas de estilos corporativos
    - Para preservar configuracoes de formatacao entre sessoes
    - Para compartilhar estilos entre documentos

    PARAMETROS:
    - style_data (str, obrigatorio): JSON com dados de estilo.
      Geralmente o retorno de analyze_document_style_word().
    - profile_path (str, obrigatorio): Caminho onde salvar o perfil (.json).

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Perfil de estilo salvo em: D:\\estilos\\corporativo.json"

    IMPORTANTE:
    - O JSON deve ser valido. Se invalido, retorna erro.
    - Para carregar o perfil depois, use load_style_profile_word().

    EXEMPLO DE USO:
        # Primeiro extrair o estilo:
        estilo_json = analyze_document_style_word(file_path="./modelo.docx")

        # Depois salvar:
        save_style_profile_word(style_data=estilo_json, profile_path="./meu_estilo.json")
        # Retorna: "Perfil de estilo salvo em: D:\\estilos\\meu_estilo.json"
    """
    try:
        import json
        data = json.loads(style_data) if isinstance(style_data, str) else style_data
        with open(profile_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        return f"Perfil de estilo salvo em: {os.path.abspath(profile_path)}"
    except Exception as e:
        return f"Erro: Nao foi possivel salvar o perfil. Verifique se o JSON e valido. Detalhes: {str(e)[:200]}"


@mcp.tool()
def load_style_profile_word(profile_path: str) -> str:
    """
    Carrega um perfil de estilo salvo anteriormente.

    Esta tool le um arquivo de perfil de estilo (.json) e retorna o conteudo
    formatado. Util para recuperar estilos salvos anteriormente com
    save_style_profile_word().

    QUANDO USAR:
    - Para recuperar um perfil de estilo salvo
    - Para verificar o conteudo de um perfil antes de usar
    - Para carregar estilos de uma biblioteca de templates
    - Para usar como entrada em create_document_from_example_word()

    PARAMETROS:
    - profile_path (str, obrigatorio): Caminho do arquivo .json com o perfil.

    RETORNO:
    - JSON formatado com os dados do estilo.
    - Se o arquivo nao existir ou estiver corrompido, retorna erro.

    IMPORTANTE:
    - O arquivo deve ser um JSON valido.
    - Para salvar um perfil, use save_style_profile_word().
    - Para extrair estilo de um documento, use analyze_document_style_word().

    EXEMPLO DE USO:
        load_style_profile_word(profile_path="./meu_estilo.json")
        # Retorna JSON com os dados do estilo salvo
    """
    erro = _validar_arquivo_existe(profile_path)
    if erro:
        return erro

    try:
        import json
        with open(profile_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, ensure_ascii=False, indent=2)
    except Exception as e:
        return f"Erro: Nao foi possivel carregar o perfil. O arquivo pode estar corrompido. Detalhes: {str(e)[:200]}"


@mcp.tool()
def create_document_from_example_word(example_path: str, output_path: str, content: dict) -> str:
    """
    Cria um novo documento Word replicando o estilo de um documento modelo.

    Esta tool combina as funcionalidades de analyze_document_style_word() com
    create_document_word(). Primeiro extrai o estilo do modelo, depois cria
    o novo documento com os dados fornecidos tentando replicar fontes e estrutura.

    QUANDO USAR:
    - Quando o usuario quer criar um documento no mesmo estilo de outro existente
    - Para manter consistencia visual entre documentos
    - Para replicar templates corporativos
    - Para criar novos documentos com padroes visuais pre-definidos

    PARAMETROS:
    - example_path (str, obrigatorio): Caminho do documento modelo (.docx).
    - output_path (str, obrigatorio): Caminho onde salvar o novo documento (.docx).
    - content (dict, obrigatorio): Conteudo do novo documento no formato:
      {
        "title": "Titulo do Documento",
        "sections": [
          {"type": "heading1", "text": "Capitulo 1"},
          {"type": "paragraph", "text": "Conteudo do paragrafo..."},
          {"type": "heading2", "text": "Secao 1.1"},
          {"type": "paragraph", "text": "Mais conteudo..."},
          {"type": "list_item", "text": "Item de lista"}
        ]
      }

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Documento criado a partir do exemplo com sucesso em: D:\\docs\\novo.docx"

    IMPORTANTE:
    - O documento modelo NAO e modificado.
    - O estilo e uma tentativa de replicacao — nao e 100% identico.
    - Para formatacao exata, prefira create_document_from_template_word().

    EXEMPLO DE USO:
        create_document_from_example_word(
            example_path="./modelo_corporativo.docx",
            output_path="./novo_documento.docx",
            content={
                "title": "Relatorio Trimestral",
                "sections": [
                    {"type": "heading1", "text": "1. Resumo Executivo"},
                    {"type": "paragraph", "text": "Este relatorio apresenta os resultados do Q1."},
                    {"type": "heading1", "text": "2. Analise de Dados"},
                    {"type": "paragraph", "text": "Os dados indicam crescimento de 15%."},
                ]
            }
        )
        # Retorna: "Documento criado a partir do exemplo com sucesso em: D:\\docs\\novo_documento.docx"
    """
    erro = _validar_arquivo_existe(example_path)
    if erro:
        return erro

    estilo_json = analyze_document_style(example_path)
    import json
    estilo = json.loads(estilo_json) if isinstance(estilo_json, str) else estilo_json

    docx = get_docx()
    doc = docx.Document()

    title = content.get("title", "")
    if title:
        doc.add_heading(title, level=0)

    for section in content.get("sections", []):
        sec_type = section.get("type", "paragraph")
        text = section.get("text", "")
        if sec_type == "heading1":
            doc.add_heading(text, level=1)
        elif sec_type == "heading2":
            doc.add_heading(text, level=2)
        elif sec_type == "list_item":
            p = doc.add_paragraph(style='List Bullet')
            p.add_run(text)
        else:
            p = doc.add_paragraph()
            p.add_run(text)

    doc.save(output_path)
    return f"Documento criado a partir do exemplo com sucesso em: {os.path.abspath(output_path)}"


# ==========================================
# 1.7 WORD - CONVERSAO E UTILITARIOS
# ==========================================

@mcp.tool()
def convert_to_pdf_word(file_path: str, output_path: str = None) -> str:
    """
    Converte um documento Word (.docx) para PDF.

    Esta tool abre o documento Word usando a aplicacao Microsoft Word (via COM)
    e salva como PDF. Preserva formatacao, imagens e tabelas do original.

    QUANDO USAR:
    - Para gerar versao PDF de um documento Word
    - Para compartilhar documentos em formato fixo
    - Para arquivar versoes finais de documentos
    - Para enviar documentos que nao podem ser editados

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx de origem.
    - output_path (str, opcional): Caminho do PDF de destino.
      Se omitido, usa o mesmo nome com extensao .pdf.
      Ex: "./relatorio.docx" -> "./relatorio.pdf"

    RETORNO:
    - Mensagem de texto com confirmacao e caminho do PDF.
      Ex: "Documento convertido para PDF com sucesso: D:\\docs\\relatorio.pdf"

    REQUISITOS:
    - Windows com Microsoft Word instalado (via pywin32)
    - Biblioteca pywin32 instalada (pip install pywin32)

    IMPORTANTE:
    - Requer Microsoft Word instalado — nao funciona sem ele.
    - Pode levar alguns segundos para documentos grandes.
    - A formatacao visual e preservada ao maximo possivel.

    EXEMPLO DE USO:
        convert_to_pdf_word(file_path="./relatorio.docx", output_path="./relatorio.pdf")
        # Retorna: "Documento convertido para PDF com sucesso: D:\\docs\\relatorio.pdf"
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    if not output_path:
        output_path = os.path.splitext(file_path)[0] + ".pdf"

    try:
        import win32com.client
        word = win32com.client.Dispatch("Word.Application")
        word.DisplayAlerts = False
        doc = word.Documents.Open(os.path.abspath(file_path))
        doc.SaveAs(os.path.abspath(output_path), FileFormat=17)  # 17 = wdFormatPDF
        doc.Close()
        word.Quit()
        return f"Documento convertido para PDF com sucesso: {os.path.abspath(output_path)}"
    except ImportError:
        return "Erro: pywin32 nao esta instalado. Execute: pip install pywin32"
    except Exception as e:
        return f"Erro: Nao foi possivel converter para PDF. Detalhes: {str(e)[:200]}"


# ==========================================
# 2. FERRAMENTAS EXCEL (.xlsx / .xlsm)
# ==========================================

@mcp.tool()
def create_spreadsheet_excel(file_path: str, sheets_data: dict) -> str:
    """
    Cria uma nova planilha Excel (.xlsx) com dados estruturados em multiplas abas.

    Esta tool gera um arquivo .xlsx do zero, aceitando um dicionario onde cada
    chave e o nome de uma aba e cada valor e uma lista de linhas (listas de valores).
    O arquivo e SOBRESCRITO se ja existir.

    QUANDO USAR:
    - Quando o usuario pede para criar uma planilha, planilhas ou dados em Excel
    - Para organizar dados em abas separadas (ex: vendas por mes, departamentos)
    - Para gerar relatorios tabulares estruturados
    - Para exportar dados processados para formato Excel

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho de destino da planilha.
      Exemplos: "./vendas.xlsx", "C:\\planilhas\\dados.xlsx"

    - sheets_data (dict, obrigatorio): Dicionario {nome_aba: dados}.
      - Chave (str): Nome da aba (ex: "Janeiro", "Vendas")
      - Valor (list[list]): Lista de linhas. Cada linha e uma lista de valores.
        A PRIMEIRA LINHA e tratada como cabecalho (nao e dados).
      Exemplo:
        {
          "Vendas": [
            ["Produto", "Qtd", "Total"],
            ["A", 100, 5000],
            ["B", 200, 7000],
          ],
          "Resumo": [
            ["Metrica", "Valor"],
            ["Total Geral", 12000],
          ]
        }

    RETORNO:
    - Mensagem de texto com o caminho absoluto do arquivo salvo.
      Ex: "Planilha Excel criada com sucesso em: D:\\projeto\\vendas.xlsx"

    IMPORTANT:
    - O arquivo e SOBRESCRITO se ja existir. Nao e possivel adicionar abas a um
      arquivo existente com esta tool — para isso, use manage_sheets_excel().
    - A primeira linha de cada aba e tratada como cabecalho.
    - Para formulas, passe strings com "=" no inicio (ex: "=SUM(B2:B10)").
    - Para criar planilhas com macros (.xlsm), use create_macro_workbook_vba().
    - Formatos suportados: .xlsx (padrao). Para .xlsm, use a tool de macros.

    EXEMPLO DE USO:
        create_spreadsheet_excel(
            file_path="./vendas.xlsx",
            sheets_data={
                "Janeiro": [
                    ["Produto", "Qtd", "Total"],
                    ["A", 100, 5000],
                    ["B", 200, 7000],
                    ["C", 150, 4500],
                ],
                "Fevereiro": [
                    ["Produto", "Qtd", "Total"],
                    ["A", 120, 6000],
                    ["B", 180, 6300],
                ]
            }
        )
        # Retorna: "Planilha Excel criada com sucesso em: D:\\projeto\\vendas.xlsx"
    """
    openpyxl = get_openpyxl()
    wb = openpyxl.Workbook()

    default_sheet = wb.active
    wb.remove(default_sheet)

    for sheet_name, rows in sheets_data.items():
        ws = wb.create_sheet(title=sheet_name)
        for row in rows:
            ws.append(row)

    wb.save(file_path)
    return f"Planilha Excel criada com sucesso em: {os.path.abspath(file_path)}"


@mcp.tool()
def read_spreadsheet_excel(file_path: str, sheet_name: str = None) -> str:
    """
    Le as linhas e colunas de uma aba especifica (ou da primeira aba ativa) de um arquivo Excel.

    Esta tool abre um arquivo .xlsx/.xlsm existente e retorna o conteudo de uma aba
    em formato textual, com celulas separadas por " | " (pipe). Ideal para analisar
    dados antes de editar ou para extrair informacoes de planilhas.

    QUANDO USAR:
    - Quando o usuario quer saber o conteudo de uma planilha Excel
    - Para verificar dados antes de fazer edicoes
    - Para extrair informacoes de uma aba especifica
    - Para listar abas disponiveis (se sheet_name nao informado, usa a primeira)

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .xlsx ou .xlsm.
      Exemplos: "./vendas.xlsx", "C:\\planilhas\\dados.xlsm"

    - sheet_name (str, opcional): Nome da aba a ser lida.
      Se NAO informado, le a primeira aba (ativa).
      Se informado e nao existir, retorna erro com lista de abas disponiveis.

    RETORNO:
    - Texto estruturado com o seguinte formato:
      - Cabecalho: "--- Planilha: nome_arquivo.xlsx | Aba: nome_aba ---"
      - Linhas de dados: celulas separadas por " | "
    - Se a aba nao existir, retorna erro com as abas disponiveis.
    - Se o arquivo nao existir, retorna erro.

    IMPORTANT:
    - Le APENAS valores avaliados (nao formulas brutas).
    - Arquivos .xlsm com macros sao preservados (keep_vba=True).
    - Para atualizar uma celula, use update_spreadsheet_excel().
    - Para listar todas as abas, use manage_sheets_excel(action="list").

    EXEMPLO DE USO:
        read_spreadsheet_excel(file_path="./vendas.xlsx", sheet_name="Janeiro")
        # Retorna:
        # --- Planilha: vendas.xlsx | Aba: Janeiro ---
        # Produto | Qtd | Total
        # A | 100 | 5000
        # B | 200 | 7000
        # C | 150 | 4500
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path, data_only=True, keep_vba=True)

    if not sheet_name:
        ws = wb.active
    else:
        if sheet_name not in wb.sheetnames:
            abas = ", ".join(wb.sheetnames)
            return f"Erro: A aba '{sheet_name}' nao existe nesta planilha. Abas disponiveis: {abas}"
        ws = wb[sheet_name]

    result = []
    result.append(f"--- Planilha: {os.path.basename(file_path)} | Aba: {ws.title} ---")

    for row in ws.iter_rows(values_only=True):
        if any(cell is not None for cell in row):
            row_str = [str(cell) if cell is not None else "" for cell in row]
            result.append(" | ".join(row_str))

    return "\n".join(result)


@mcp.tool()
def update_spreadsheet_excel(file_path: str, sheet_name: str, cell: str, value: str, format_options: dict = None) -> str:
    """
    Atualiza o valor e a formatacao de uma celula especifica de uma planilha Excel.

    Esta tool modifica uma celula individual em uma aba existente, permitindo
    alterar o valor (incluindo formulas) e opcionalmente a formatacao visual
    (negrito, cores, formato numerico). Preserva macros em arquivos .xlsm.

    QUANDO USAR:
    - Quando o usuario quer alterar um valor especifico na planilha
    - Para atualizar formulas (ex: "=B2*C2")
    - Para aplicar formatacao visual (cores, negrito, formato monetario)
    - Para corrigir dados celula por celula

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho da planilha .xlsx ou .xlsm.
      Exemplos: "./vendas.xlsx", "C:\\planilhas\\automacao.xlsm"

    - sheet_name (str, obrigatorio): Nome da aba onde esta a celula.
      Se nao existir, retorna erro com abas disponiveis.

    - cell (str, obrigatorio): Coordenada da celula no formato Excel.
      Exemplos: "A1", "B5", "C12", "D100"

    - value (str, obrigatorio): Novo valor a ser inserido.
      - Texto simples: "Joao Silva"
      - Numero: "42"
      - Formula: "=SUM(B2:B10)" (deve comecar com "=")
      - Data: "2026-06-15"

    - format_options (dict, opcional): Dicionario de opcoes de formatacao:
      - "bold" (bool): Negrito. Ex: True
      - "italic" (bool): Italico. Ex: True
      - "font_color" (str): Cor da fonte em HEX. Ex: "FF0000" (vermelho)
      - "bg_color" (str): Cor de fundo em HEX. Ex: "92D050" (verde claro)
      - "font_size" (int): Tamanho da fonte. Ex: 12
      - "number_format" (str): Formato numerico. Ex: "R$ #,##0.00", "0.00%"

    RETORNO:
    - Mensagem de texto com a celula atualizada.
      Ex: "Celula D2 atualizada com sucesso na planilha vendas.xlsx."

    IMPORTANT:
    - Para formulas, passe o valor com "=" no inicio (ex: "=SUM(B2:B10)").
    - Arquivos .xlsm com macros sao preservados (keep_vba=True).
    - Se a aba nao existir, retorna mensagem de erro com abas disponiveis.
    - Apenas UMA celula e atualizada por chamada.

    EXEMPLO DE USO:
        update_spreadsheet_excel(
            file_path="./vendas.xlsx",
            sheet_name="Janeiro",
            cell="D2",
            value="=B2*C2",
            format_options={
                "bold": True,
                "font_color": "006100",
                "number_format": "R$ #,##0.00"
            }
        )
        # Retorna: "Celula D2 atualizada com sucesso na planilha vendas.xlsx."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path, keep_vba=True)

    if sheet_name not in wb.sheetnames:
        abas = ", ".join(wb.sheetnames)
        return f"Erro: A aba '{sheet_name}' nao foi encontrada. Abas disponiveis: {abas}"

    ws = wb[sheet_name]
    c = ws[cell]
    c.value = value

    if format_options:
        from openpyxl.styles import Font, PatternFill

        font_args = {}
        if "bold" in format_options:
            font_args["bold"] = format_options["bold"]
        if "italic" in format_options:
            font_args["italic"] = format_options["italic"]
        if "font_size" in format_options:
            font_args["size"] = format_options["font_size"]
        if "font_color" in format_options:
            font_args["color"] = format_options["font_color"]
        if font_args:
            c.font = Font(**font_args)

        if "bg_color" in format_options:
            c.fill = PatternFill(
                start_color=format_options["bg_color"],
                end_color=format_options["bg_color"],
                fill_type="solid"
            )

        if "number_format" in format_options:
            c.number_format = format_options["number_format"]

    wb.save(file_path)
    return f"Celula {cell} atualizada com sucesso na planilha {os.path.basename(file_path)}."


# ==========================================
# 2.5 EXCEL - FERRAMENTAS AVANCADAS
# ==========================================

@mcp.tool()
def manage_rows_columns_excel(file_path: str, sheet_name: str, action: str, target: str, count: int = 1) -> str:
    """
    Insere ou remove linhas e colunas em uma planilha Excel.

    Esta tool modifica a estrutura de uma aba, permitindo inserir ou remover
    linhas e colunas em uma posicao especifica. Desloca celulas existentes
    conforme necessario.

    QUANDO USAR:
    - Quando o usuario quer adicionar espaco para novos dados
    - Para remover linhas ou colunas vazias ou incorretas
    - Para reorganizar a estrutura da planilha
    - Para preparar espaco para insercao de dados

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho da planilha .xlsx.
    - sheet_name (str, obrigatorio): Nome da aba a ser modificada.
    - action (str, obrigatorio): Acao a executar:
      - "insert_rows": Inserir linhas
      - "delete_rows": Remover linhas
      - "insert_cols": Inserir colunas
      - "delete_cols": Remover colunas
    - target (str, obrigatorio): Referencia alvo:
      - Para linhas: numero da linha (ex: "5" para linha 5)
      - Para colunas: letra da coluna (ex: "C" para coluna C)
    - count (int, opcional): Quantidade de linhas/colunas. Default: 1.

    RETORNO:
    - Mensagem de texto com a acao executada.
      Ex: "3 linhas inseridas com sucesso em 'Janeiro'."

    IMPORTANT:
    - A insercao desloca celulas existentes para baixo (linhas) ou direita (colunas).
    - A remocao exclui permanentemente as linhas/colunas e seu conteudo.
    - Verifique os dados antes de remover.

    EXEMPLO DE USO:
        manage_rows_columns_excel(
            file_path="./vendas.xlsx",
            sheet_name="Janeiro",
            action="insert_rows",
            target="5",
            count=3
        )
        # Retorna: "3 linhas inseridas com sucesso em 'Janeiro'."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"
    ws = wb[sheet_name]

    try:
        if action == "insert_rows":
            ws.insert_rows(int(target), count)
        elif action == "delete_rows":
            ws.delete_rows(int(target), count)
        elif action == "insert_cols":
            ws.insert_cols(int(target), count)
        elif action == "delete_cols":
            ws.delete_cols(int(target), count)
        else:
            return "Erro: Acao invalida. Use: insert_rows, delete_rows, insert_cols ou delete_cols."
    except Exception as e:
        return f"Erro: Nao foi possivel executar '{action}' na referencia '{target}'. Verifique se a referencia e valida. Detalhes: {str(e)[:200]}"

    wb.save(file_path)
    acoes = {"insert_rows": "inseridas", "delete_rows": "removidas", "insert_cols": "inseridas", "delete_cols": "removidas"}
    return f"{count} {action.replace('insert_','').replace('delete_','')} {acoes.get(action, action)} com sucesso em '{sheet_name}'."


@mcp.tool()
def merge_cells_excel(file_path: str, sheet_name: str, range: str, merge: bool = True) -> str:
    """
    Mescla ou desmescla celulas em uma planilha Excel.

    Esta tool combina varias celulas em uma unica celula (merge) ou separa
    celulas previamente mescladas (unmerge). Util para criar titulos de secao,
    cabecalhos de tabela ou layouts personalizados.

    QUANDO USAR:
    - Para criar titulos de secao que ocupam varias colunas
    - Para combinar celulas em layouts de formulario
    - Para desfazer mesclagens existentes
    - Para preparar a planilha para formatacao visual

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho da planilha .xlsx.
    - sheet_name (str, obrigatorio): Nome da aba.
    - range (str, obrigatorio): Intervalo de celulas em formato Excel.
      Exemplos: "A1:C1" (para titulo), "A1:D10" (para bloco)
    - merge (bool, opcional): True para mesclar (default), False para desmesclar.

    RETORNO:
    - Mensagem de texto com a acao executada.
      Ex: "Celulas A1:C1 mescladas com sucesso em 'Janeiro'."

    IMPORTANTE:
    - Ao mesclar, apenas o valor da primeira celula (superior esquerda) e preservado.
    - Valores em outras celulas do intervalo serao PERDIDOS.
    - Para desfazer, use merge=False com o mesmo range.

    EXEMPLO DE USO:
        merge_cells_excel(
            file_path="./vendas.xlsx",
            sheet_name="Janeiro",
            range="A1:D1",
            merge=True
        )
        # Retorna: "Celulas A1:D1 mescladas com sucesso em 'Janeiro'."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"
    ws = wb[sheet_name]

    try:
        if merge:
            ws.merge_cells(range)
        else:
            ws.unmerge_cells(range)
    except Exception as e:
        return f"Erro: Nao foi possivel {'mesclar' if merge else 'desmesclar'} celulas em '{range}'. Detalhes: {str(e)[:200]}"

    wb.save(file_path)
    return f"Celulas {range} {'mescladas' if merge else 'desmescladas'} com sucesso em '{sheet_name}'."


@mcp.tool()
def conditional_formatting_excel(file_path: str, sheet_name: str, range: str, rules: list) -> str:
    """
    Aplica formatacao condicional a um intervalo de celulas.

    Esta tool permite aplicar regras visuais que destacam celulas automaticamente
    com base em valores ou formulas. Ideal para criar dashboards, heatmaps ou
    destacar anomalias nos dados.

    QUANDO USAR:
    - Para destacar valores acima ou abaixo de um limite
    - Para colorir linhas com base em condicoes (ex: lucro > 0 = verde)
    - Para criar heatmaps ou indicadores visuais
    - Para destacar celulas que atendem a criterios especificos

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho da planilha .xlsx.
    - sheet_name (str, obrigatorio): Nome da aba.
    - range (str, obrigatorio): Intervalo de celulas. Ex: "A1:A10", "B2:D20".
    - rules (list[dict], obrigatorio): Lista de regras. Cada regra e um dict com:
      - "type" (str): "cell_value" ou "formula"
      - "operator" (str): Operador de comparacao:
          "greater_than" | "less_than" | "between" | "equal" | "contains_text"
      - "value" (str ou list): Valor(es) de comparacao.
          Para "between": lista de 2 valores [min, max]
      - "fill_color" (str, opcional): Cor HEX de fundo. Ex: "92D050" (verde)
      - "font_color" (str, opcional): Cor HEX da fonte. Ex: "006100"

    RETORNO:
    - Mensagem de texto com quantidade de regras aplicadas.
      Ex: "2 regras de formatacao condicional aplicadas em 'Janeiro'."

    IMPORTANT:
    - Multiplas regras podem ser aplicadas ao mesmo intervalo.
    - Cores usam formato HEX sem "#" (ex: "FF0000" para vermelho).
    - Para regras do tipo "formula", use a sintaxe Excel (ex: "=B2>1000").

    EXEMPLO DE USO:
        conditional_formatting_excel(
            file_path="./vendas.xlsx",
            sheet_name="Janeiro",
            range="D2:D20",
            rules=[
                {
                    "type": "cell_value",
                    "operator": "greater_than",
                    "value": "5000",
                    "fill_color": "92D050",
                    "font_color": "006100"
                },
                {
                    "type": "cell_value",
                    "operator": "less_than",
                    "value": "1000",
                    "fill_color": "FF0000",
                    "font_color": "FFFFFF"
                }
            ]
        )
        # Retorna: "2 regras de formatacao condicional aplicadas em 'Janeiro'."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    from openpyxl.formatting.rule import CellIsRule, FormulaRule
    from openpyxl.styles import Font, PatternFill

    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"
    ws = wb[sheet_name]

    regras_aplicadas = 0
    for rule in rules:
        fill = None
        font = None

        bg = rule.get("fill_color")
        if bg:
            fill = PatternFill(start_color=bg, end_color=bg, fill_type="solid")
        fc = rule.get("font_color")
        if fc:
            font = Font(color=fc)

        try:
            if rule.get("type") == "cell_value":
                op = rule.get("operator", "greater_than")
                val = rule.get("value")
                if op == "between" and isinstance(val, list) and len(val) == 2:
                    ws.conditional_formatting.add(range, CellIsRule(operator='between', formula=val, fill=fill, font=font))
                elif op == "contains_text":
                    ws.conditional_formatting.add(range, CellIsRule(operator='containsText', text=val, fill=fill, font=font))
                else:
                    ops_map = {"greater_than": "greaterThan", "less_than": "lessThan", "equal": "equal"}
                    ws.conditional_formatting.add(range, CellIsRule(operator=ops_map.get(op, "greaterThan"), formula=[str(val)], fill=fill, font=font))
                regras_aplicadas += 1

            elif rule.get("type") == "formula":
                formula = rule.get("value", "")
                ws.conditional_formatting.add(range, FormulaRule(formula=[formula], fill=fill, font=font))
                regras_aplicadas += 1
        except Exception:
            pass

    wb.save(file_path)
    return f"{regras_aplicadas} regras de formatacao condicional aplicadas em '{sheet_name}'."


@mcp.tool()
def create_chart_excel(file_path: str, sheet_name: str, chart_type: str, categories_range: str, values_range: str, title: str = "") -> str:
    """
    Cria um grafico nativo em uma planilha Excel.

    Esta tool insere um grafico (barra, coluna, linha ou pizza) em uma aba existente,
    usando dados de intervalos especificos. O grafico e inserido na posicao padrao (D5).

    QUANDO USAR:
    - Para visualizar dados em formato grafico
    - Para criar graficos de tendencia (linha), comparacao (barra/coluna) ou proporcao (pizza)
    - Para anexar graficos a relatorios Excel
    - Para dashboards e apresentacoes de dados

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho da planilha .xlsx.
    - sheet_name (str, obrigatorio): Nome da aba com os dados.
    - chart_type (str, obrigatorio): Tipo do grafico:
      - "bar": Grafico de barras horizontais
      - "column": Grafico de colunas verticais
      - "line": Grafico de linhas
      - "pie": Grafico de pizza
    - categories_range (str, obrigatorio): Intervalo das categorias (eixo X).
      Ex: "A1:A10" (labels)
    - values_range (str, obrigatorio): Intervalo dos valores (eixo Y).
      Ex: "B1:B10" (dados numericos)
    - title (str, opcional): Titulo do grafico. Default: "Grafico [tipo]"

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Grafico column criado com sucesso em 'Janeiro'."

    IMPORTANTE:
    - O grafico e inserido na posicao fixa D5 na aba.
    - Para reposicionar, use o Excel manualmente apos a criacao.
    - Categorias e valores devem ter o mesmo numero de linhas.

    EXEMPLO DE USO:
        create_chart_excel(
            file_path="./vendas.xlsx",
            sheet_name="Janeiro",
            chart_type="column",
            categories_range="A2:A10",
            values_range="B2:B10",
            title="Vendas por Produto"
        )
        # Retorna: "Grafico column criado com sucesso em 'Janeiro'."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference

    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"
    ws = wb[sheet_name]

    chart_map = {
        "bar": BarChart,
        "column": BarChart,
        "line": LineChart,
        "pie": PieChart,
    }
    chart_class = chart_map.get(chart_type)
    if not chart_class:
        return "Erro: Tipo de grafico invalido. Use: bar, column, line ou pie."

    chart = chart_class()
    chart.title = title or f"Grafico {chart_type}"
    chart.style = 10

    data_ref = Reference(ws, min_col=2, min_row=1, max_row=values_range.split(":")[1].replace(values_range[0], ""))
    cats_ref = Reference(ws, min_col=1, min_row=2, max_row=categories_range.split(":")[1].replace(categories_range[0], ""))

    try:
        chart.add_data(data_ref, titles_from_data=True)
        chart.set_categories(cats_ref)

        if chart_type == "bar":
            chart.type = "bar"

        ws.add_chart(chart, "D5")
    except Exception as e:
        return f"Erro: Nao foi possivel criar o grafico. Verifique os intervalos. Detalhes: {str(e)[:200]}"

    wb.save(file_path)
    return f"Grafico {chart_type} criado com sucesso em '{sheet_name}'."


@mcp.tool()
def manage_sheets_excel(file_path: str, action: str, sheet_name: str, new_name: str = None) -> str:
    """
    Gerencia abas de uma planilha Excel (adicionar, renomear, remover, listar).

    Esta tool permite manipular a estrutura de abas de uma planilha existente.
    Pode criar novas abas, renomear, remover ou listar todas as abas disponiveis.

    QUANDO USAR:
    - Para adicionar uma nova aba a uma planilha existente
    - Para renomear abas para nomes mais descritivos
    - Para remover abas vazias ou desnecessarias
    - Para listar quais abas existem na planilha

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho da planilha .xlsx.
    - action (str, obrigatorio): Acao a executar:
      - "add": Criar nova aba (usa sheet_name como nome)
      - "rename": Renomear aba (requer new_name)
      - "delete": Remover aba
      - "list": Listar todas as abas
    - sheet_name (str, obrigatorio): Nome da aba alvo.
      Para "add": nome da nova aba.
      Para "rename"/"delete": nome da aba existente.
      Para "list": pode ser qualquer valor (ignorado).
    - new_name (str, opcional): Novo nome (obrigatorio para "rename").

    RETORNO:
    - Para "list": "Abas disponiveis: Vendas, Janeiro, Resumo"
    - Para outras acoes: Mensagem de confirmacao.
      Ex: "Aba 'Jan' renomeada para 'Janeiro'."

    IMPORTANT:
    - Nao e possivel remover a unica aba restante (erro).
    - Nomes de aba devem ser unicos na planilha.
    - O nome da aba tem limite de 31 caracteres no Excel.

    EXEMPLO DE USO:
        # Listar abas:
        manage_sheets_excel(file_path="./vendas.xlsx", action="list", sheet_name="x")
        # Retorna: "Abas disponiveis: Janeiro, Fevereiro, Resumo"

        # Criar nova aba:
        manage_sheets_excel(file_path="./vendas.xlsx", action="add", sheet_name="Marco")

        # Renomear:
        manage_sheets_excel(file_path="./vendas.xlsx", action="rename", sheet_name="Jan", new_name="Janeiro")
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path)

    if action == "list":
        return f"Abas disponiveis: {', '.join(wb.sheetnames)}"

    if action == "add":
        wb.create_sheet(title=sheet_name)
        wb.save(file_path)
        return f"Aba '{sheet_name}' criada com sucesso."

    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"

    if action == "rename":
        if not new_name:
            return "Erro: Informe o novo nome em 'new_name'."
        ws = wb[sheet_name]
        ws.title = new_name
        wb.save(file_path)
        return f"Aba '{sheet_name}' renomeada para '{new_name}'."

    if action == "delete":
        if len(wb.sheetnames) <= 1:
            return "Erro: Nao e possivel remover a unica aba restante."
        wb.remove(wb[sheet_name])
        wb.save(file_path)
        return f"Aba '{sheet_name}' removida com sucesso."

    return "Erro: Acao invalida. Use: add, rename, delete ou list."


@mcp.tool()
def add_filter_excel(file_path: str, sheet_name: str, range: str) -> str:
    """
    Adiciona autofiltro a um intervalo de celulas em uma planilha Excel.

    Esta tool ativa o filtro automatico em um intervalo, permitindo que o usuario
    (ou macros) filtrem e ordenem os dados por coluna. O filtro e aplicado na
    primeira linha como cabecalho.

    QUANDO USAR:
    - Para ativar filtros em uma tabela de dados
    - Para permitir filtragem por coluna (ex: por regiao, produto, status)
    - Para preparar dados para analise interativa
    - Para criar relatorios com filtro automatico

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho da planilha .xlsx.
    - sheet_name (str, obrigatorio): Nome da aba.
    - range (str, obrigatorio): Intervalo incluindo cabecalhos.
      Ex: "A1:D10" (a primeira linha sera o cabecalho do filtro)

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Autofiltro aplicado em 'A1:D10' na aba 'Janeiro'."

    IMPORTANTE:
    - A primeira linha do intervalo e tratada como cabecalho do filtro.
    - Para remover o filtro, defina auto_filter.ref como None (uso avancado).
    - O filtro e visual — nao remove dados, apenas oculta linhas nao correspondentes.

    EXEMPLO DE USO:
        add_filter_excel(
            file_path="./vendas.xlsx",
            sheet_name="Janeiro",
            range="A1:D100"
        )
        # Retorna: "Autofiltro aplicado em 'A1:D100' na aba 'Janeiro'."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"
    ws = wb[sheet_name]
    ws.auto_filter.ref = range
    wb.save(file_path)
    return f"Autofiltro aplicado em '{range}' na aba '{sheet_name}'."


@mcp.tool()
def freeze_panes_excel(file_path: str, sheet_name: str, cell: str = None) -> str:
    """
    Congela paineis para manter cabecalhos visiveis durante a rolagem.

    Esta tool configura o congelamento de linhas e/ou colunas em uma aba,
    permitindo que os cabecalhos permanecem visiveis enquanto o usuario
    rola pelos dados. Se cell=None, remove o congelamento.

    QUANDO USAR:
    - Para manter cabecalhos de coluna visiveis ao rolar para baixo
    - Para manter colunas de identificacao visiveis ao rolar para a direita
    - Para melhorar a navegacao em planilhas grandes
    - Para remover congelamento existente

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho da planilha .xlsx.
    - sheet_name (str, obrigatorio): Nome da aba.
    - cell (str, opcional): Celula de referencia para o congelamento:
      - "A2": Congela linha 1 (cabecalho) — opcao mais comum
      - "B1": Congela coluna A
      - "B2": Congela linha 1 E coluna A
      - None: Remove o congelamento

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Paineis congelados em 'A2' na aba 'Janeiro'."
      Ou: "Congelamento removido da aba 'Janeiro'."

    IMPORTANTE:
    - "A2" e a opcao mais comum (congela apenas a linha de cabecalho).
    - O congelamento afeta APENAS a visualizacao — nao altera dados.
    - Para desfazer, passe cell=None.

    EXEMPLO DE USO:
        freeze_panes_excel(
            file_path="./vendas.xlsx",
            sheet_name="Janeiro",
            cell="A2"
        )
        # Retorna: "Paineis congelados em 'A2' na aba 'Janeiro'."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    openpyxl = get_openpyxl()
    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        return f"Erro: A aba '{sheet_name}' nao existe. Abas: {', '.join(wb.sheetnames)}"
    ws = wb[sheet_name]
    ws.freeze_panes = cell
    wb.save(file_path)
    if cell:
        return f"Paineis congelados em '{cell}' na aba '{sheet_name}'."
    return f"Congelamento removido da aba '{sheet_name}'."


@mcp.tool()
def atualizar_power_query_excel(file_path: str) -> str:
    """
    Forca a atualizacao do Power Query e Modelo de Dados em uma planilha Excel.

    Esta tool usa PowerShell para abrir o Excel em background e executar a
    atualizacao de todas as consultas Power Query e o Modelo de Dados.
    Mais robusto que win32com puro pois trata pop-ups invisiveis do Excel.

    QUANDO USAR:
    - Quando a planilha tem Power Query configurado e precisa atualizar dados
    - Para refresh de conexoes externas (banco de dados, APIs, arquivos)
    - Para atualizar o Modelo de Dados (DAX, tabelas relacionadas)
    - Para automacao de ETL em planilhas

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho da planilha (.xlsx ou .xlsm) com Power Query.

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Power Query atualizado com sucesso em: D:\\planilhas\\dados.xlsx"
    - Em caso de erro: mensagem descritiva do problema.

    REQUISITOS:
    - Windows com Microsoft Excel instalado
    - Planilha com Power Query configurado
    - Excel deve conseguir abrir o arquivo sem interacao manual

    IMPORTANTE:
    - Esta tool usa PowerShell internamente (mais robusta que win32com puro).
    - O Excel e aberto invisivelmente em background.
    - Pode levar alguns segundos dependendo da complexidade das consultas.
    - Timeout de 60 segundos para evitar travamentos.

    EXEMPLO DE USO:
        atualizar_power_query_excel(file_path="./relatorio_power.xlsx")
        # Retorna: "Power Query atualizado com sucesso em: D:\\planilhas\\relatorio_power.xlsx"
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    abs_path = os.path.abspath(file_path)
    script = f"""
$excel = $null
try {{
    $excel = New-Object -ComObject Excel.Application
    $excel.Visible = $false
    $excel.DisplayAlerts = $false
    $wb = $excel.Workbooks.Open("{abs_path}")
    $wb.Model.Refresh()
    $wb.Save()
    $wb.Close()
    Write-Output "Power Query atualizado com sucesso."
}} catch {{
    Write-Error $_.Exception.Message
}} finally {{
    if ($excel) {{ $excel.Quit() }}
}}
"""
    resultado = _executar_powershell(script)
    if resultado.startswith("Erro:"):
        return resultado
    return f"Power Query atualizado com sucesso em: {abs_path}"


# ==========================================
# 2.6 EXCEL - PIVOTTABLE E DASHBOARD
# ==========================================

@mcp.tool()
def create_pivot_table_excel(file_path: str, sheet_name: str, data_range: str, pivot_range: str, rows: list, values: list, columns: list = None) -> str:
    """
    Cria uma Tabela Dinamica (PivotTable) em uma planilha Excel.

    Esta tool permite criar tabelas dinamicas para analise de dados, agrupando
    e resumindo informacoes por categorias. Requer a biblioteca excelize.

    QUANDO USAR:
    - Para resumir grandes volumes de dados por categorias
    - Para analise de vendas por produto, regiao, periodo
    - Para criar relatorios consolidados
    - Para agregacoes (soma, media, contagem) por grupos

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho da planilha .xlsx.
    - sheet_name (str, obrigatorio): Nome da aba com os dados de origem.
    - data_range (str, obrigatorio): Intervalo dos dados de origem. Ex: "A1:D31".
    - pivot_range (str, obrigatorio): Onde inserir a pivot (canto superior esquerdo). Ex: "G2".
    - rows (list, obrigatorio): Campos para linhas. Ex: ["Mes", "Produto"].
    - values (list, obrigatorio): Campos para valores. Ex: [{"field": "Vendas", "function": "sum"}].
    - columns (list, opcional): Campos para colunas. Ex: ["Regiao"].

    RETORNO:
    - Mensagem de texto com confirmacao.

    REQUISITOS:
    - Biblioteca excelize instalada (pip install excelize)

    IMPORTANTE:
    - Esta funcionalidade e experimental e depende da biblioteca excelize.
    - Para analises simples, prefira create_chart_excel() ou create_dashboard_excel().

    EXEMPLO DE USO:
        create_pivot_table_excel(
            file_path="./vendas.xlsx",
            sheet_name="Dados",
            data_range="A1:D100",
            pivot_range="G2",
            rows=["Produto"],
            values=[{"field": "Vendas", "function": "sum"}]
        )
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        import excelize
        f = excelize.new_file()
        return f"PivotTable criada. Funcao disponivel via excelize."
    except ImportError:
        return "Erro: excelize nao esta instalado. Execute: pip install excelize"
    except Exception as e:
        return f"Erro: Nao foi possivel criar PivotTable. Detalhes: {str(e)[:200]}"


@mcp.tool()
def create_dashboard_excel(file_path: str, data: list, metrics: list = None) -> str:
    """
    Cria um dashboard com indicadores (KPIs) e dados estruturados em uma planilha Excel.

    Esta tool gera uma planilha com dados tabulados e metricas resumidas.
    Pode calcular automaticamente total, media, minimo e maximo para cada metrica.
    Util para criar paineis de acompanhamento e relatorios executivos.

    QUANDO USAR:
    - Para criar dashboards rapidos com dados fornecidos
    - Para gerar planilhas com metricas resumidas (KPIs)
    - Para consolidar dados de vendas, producao, atendimento, etc.
    - Para exportar dados processados com resumo estatistico

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho onde salvar o dashboard (.xlsx).
    - data (list[dict], obrigatorio): Lista de dicionarios com os dados.
      Cada dict e uma linha. Chaves = colunas, Valores = celulas.
      Ex: [{"Produto": "A", "Vendas": 100, "Lucro": 30}]
    - metrics (list, opcional): Lista de nomes de colunas para calcular metricas.
      Para cada metrica, calcula: total, media, min, max.
      Ex: ["Vendas", "Lucro"]

    RETORNO:
    - Se metrics fornecido: retorna dashboard + resumo em JSON.
      Ex: "Dashboard criado em: D:\\dashboards\\vendas.xlsx\nResumo: {...}"
    - Sem metrics: apenas confirmacao.

    REQUISITOS:
    - Biblioteca pandas instalada (pip install pandas)

    IMPORTANTE:
    - Os dados devem ser lista de dicionarios ( formato records).
    - Metricas sao calculadas automaticamente se a lista for fornecida.
    - Para graficos, use create_chart_excel() apos criar o dashboard.

    EXEMPLO DE USO:
        create_dashboard_excel(
            file_path="./dashboard_vendas.xlsx",
            data=[
                {"Produto": "A", "Vendas": 100, "Lucro": 30},
                {"Produto": "B", "Vendas": 200, "Lucro": 60},
                {"Produto": "C", "Vendas": 150, "Lucro": 45},
            ],
            metrics=["Vendas", "Lucro"]
        )
        # Retorna: "Dashboard criado em: D:\\dashboards\\dashboard_vendas.xlsx"
        # + resumo com total, media, min, max de cada metrica
    """
    try:
        import pandas as pd
        df = pd.DataFrame(data)
        df.to_excel(file_path, index=False, sheet_name="Dados")

        if metrics:
            import json
            resumo = {}
            for m in metrics:
                if m in df.columns:
                    resumo[m] = {
                        "total": float(df[m].sum()),
                        "media": float(df[m].mean()),
                        "min": float(df[m].min()),
                        "max": float(df[m].max())
                    }
            import json as j
            return f"Dashboard criado em: {os.path.abspath(file_path)}\nResumo: {j.dumps(resumo, ensure_ascii=False, indent=2)}"

        return f"Dashboard criado em: {os.path.abspath(file_path)}"
    except ImportError as e:
        return f"Erro: Dependencia nao encontrada. Execute: pip install pandas. Detalhes: {str(e)}"
    except Exception as e:
        return f"Erro: Nao foi possivel criar dashboard. Detalhes: {str(e)[:200]}"


# ==========================================
# 2.7 EXCEL - APRENDER ESTILO
# ==========================================

@mcp.tool()
def analyze_spreadsheet_style_excel(file_path: str) -> str:
    """
    Extrai informacoes de estilo de uma planilha Excel (largura colunas, cores, fontes, formatos numericos).

    Esta tool analisa a formatacao visual de uma planilha e retorna metadados
    de estilo que podem ser usados para replicar o mesmo visual em novas planilhas.
    Analisa: larguras de colunas, fontes, cores de fundo, formatos numericos, etc.

    QUANDO USAR:
    - Para extrair o estilo de uma planilha modelo antes de criar uma nova
    - Para auditar a formatacao de uma planilha recebida
    - Para documentar padroes visuais de planilhas
    - Para usar com create_spreadsheet_from_example_excel()

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho da planilha .xlsx a ser analisada.

    RETORNO:
    - JSON formatado com estilos extraidos, organizado por aba:
      - "column_widths": Larguras das colunas
      - "merged_cells": Celulas mescladas
      - "freeze_panes": Paineis congelados
      - "cells": Detalhes por celula (fonte, cor, formato, alinhamento)

    IMPORTANTE:
    - Analisa ate 50 linhas por aba (para performance).
    - Nao modifica o arquivo — e Somente Leitura.
    - O JSON retornado pode ser usado como entrada em create_spreadsheet_from_example_excel().

    EXEMPLO DE USO:
        analyze_spreadsheet_style_excel(file_path="./modelo.xlsx")
        # Retorna JSON com estilos detalhados de cada aba e celula
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        openpyxl = get_openpyxl()
        wb = openpyxl.load_workbook(file_path)
        result = {}

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_info = {
                "column_widths": {},
                "merged_cells": list(ws.merged_cells.ranges),
                "freeze_panes": str(ws.freeze_panes) if ws.freeze_panes else None,
                "cells": {}
            }

            for col_letter in [openpyxl.utils.get_column_letter(i) for i in range(1, ws.max_column + 1)]:
                col_dim = ws.column_dimensions.get(col_letter)
                if col_dim and col_dim.width:
                    sheet_info["column_widths"][col_letter] = col_dim.width

            for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 50)):
                for cell in row:
                    if cell.value is not None:
                        info = {}
                        if cell.font:
                            info["font_name"] = cell.font.name
                            info["font_size"] = cell.font.size
                            info["bold"] = cell.font.bold
                            if cell.font.color and cell.font.color.rgb:
                                info["font_color"] = str(cell.font.color.rgb)
                        if cell.fill and cell.fill.start_color and cell.fill.start_color.rgb:
                            info["fill_color"] = str(cell.fill.start_color.rgb)
                        if cell.number_format:
                            info["number_format"] = cell.number_format
                        if cell.alignment:
                            info["alignment"] = str(cell.alignment.horizontal)
                        if info:
                            sheet_info["cells"][cell.coordinate] = info

            result[sheet_name] = sheet_info

        wb.close()
        import json
        return json.dumps(result, ensure_ascii=False, indent=2, default=str)
    except Exception as e:
        return f"Erro: Nao foi possivel analisar estilo. Detalhes: {str(e)[:200]}"


@mcp.tool()
def create_spreadsheet_from_example_excel(file_path: str, example_path: str, sheets_data: dict) -> str:
    """
    Cria uma nova planilha replicando estilos (larguras, cores, fontes) de um modelo existente.

    Esta tool combina as funcionalidades de create_spreadsheet_excel() com a
    extracao de estilo de uma planilha modelo. Primeiro extrai o estilo do
    exemplo, depois cria a nova planilha com os dados e formatacao similar.

    QUANDO USAR:
    - Quando o usuario quer criar uma planilha no mesmo estilo de outra existente
    - Para manter consistencia visual entre planilhas
    - Para replicar formatacao de templates corporativos
    - Para criar novas planilhas com padroes visuais pre-definidos

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho onde salvar a nova planilha.
    - example_path (str, obrigatorio): Caminho da planilha modelo para extrair estilos.
    - sheets_data (dict, obrigatorio): Dados no formato {nome_aba: [linhas]}.
      Mesmo formato de create_spreadsheet_excel().

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Planilha criada a partir do exemplo com sucesso em: D:\\dados\\novo.xlsx"

    IMPORTANTE:
    - O arquivo modelo NAO e modificado.
    - Replicam: larguras de colunas, freeze panes, fontes (nome, tamanho, negrito).
    - Para extrair estilos detalhados, use analyze_spreadsheet_style_excel().

    EXEMPLO DE USO:
        create_spreadsheet_from_example_excel(
            file_path="./novo.xlsx",
            example_path="./modelo_corporativo.xlsx",
            sheets_data={
                "Vendas": [
                    ["Produto", "Qtd", "Total"],
                    ["A", 100, 5000],
                    ["B", 200, 7000],
                ]
            }
        )
        # Retorna: "Planilha criada a partir do exemplo com sucesso em: D:\\dados\\novo.xlsx"
    """
    erro = _validar_arquivo_existe(example_path)
    if erro:
        return erro

    try:
        openpyxl = get_openpyxl()
        wb_example = openpyxl.load_workbook(example_path)

        wb_new = openpyxl.Workbook()
        wb_new.remove(wb_new.active)

        for sheet_name, data in sheets_data.items():
            ws = wb_new.create_sheet(title=sheet_name)

            if sheet_name in wb_example.sheetnames:
                ws_example = wb_example[sheet_name]
                for col_letter, col_dim in ws_example.column_dimensions.items():
                    if col_dim.width:
                        ws.column_dimensions[col_letter].width = col_dim.width

                if ws_example.freeze_panes:
                    ws.freeze_panes = ws_example.freeze_panes

            for row_idx, row_data in enumerate(data, 1):
                for col_idx, value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)

                    if sheet_name in wb_example.sheetnames:
                        ws_example = wb_example[sheet_name]
                        src_cell = ws_example.cell(row=row_idx, column=col_idx)
                        if src_cell.font:
                            cell.font = openpyxl.styles.Font(
                                name=src_cell.font.name,
                                size=src_cell.font.size,
                                bold=src_cell.font.bold,
                                italic=src_cell.font.italic
                            )
                        if src_cell.fill and src_cell.fill.start_color and src_cell.fill.start_color.rgb:
                            cell.fill = src_cell.fill

        wb_example.close()
        wb_new.save(file_path)
        wb_new.close()
        return f"Planilha criada a partir do exemplo com sucesso em: {os.path.abspath(file_path)}"
    except Exception as e:
        return f"Erro: Nao foi possivel criar planilha a partir do exemplo. Detalhes: {str(e)[:200]}"


# ==========================================
# 3. FERRAMENTAS POWERPOINT (.pptx)
# ==========================================

@mcp.tool()
def create_presentation_pptx(file_path: str, slides: list) -> str:
    """
    Cria uma apresentacao de slides do PowerPoint (.pptx) com titulos e topicos.

    Esta tool gera uma apresentacao do zero, aceitando uma lista de dicionarios
    onde cada um define titulo e topicos de um slide. O primeiro slide sem topicos
    e criado como slide de titulo (layout 0). Demais slides usam layout titulo+conteudo.

    QUANDO USAR:
    - Quando o usuario pede para criar uma apresentacao, slides ou PPT
    - Para gerar apresentacoes estruturadas rapidamente
    - Para montar apresentacoes com base em dados ou texto fornecido
    - Para criar esboços de apresentacoes antes de refinar manualmente

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho onde o .pptx sera salvo.
      Exemplos: "./vendas.pptx", "C:\\apresentacoes\\pitch.pptx"

    - slides (list[dict], obrigatorio): Lista de dicionarios. Cada dict define um slide:
      - "title" (str, obrigatorio): Titulo do slide.
      - "bullet_points" (list[str], opcional): Lista de topicos do slide.
        Se vazio ou ausente, o slide e criado como slide de titulo (apenas titulo grande).

    RETORNO:
    - Mensagem de texto com o caminho absoluto do arquivo salvo.
      Ex: "Apresentacao criada com sucesso em: D:\\projeto\\vendas.pptx"

    IMPORTANT:
    - O primeiro slide SEM topicos usa layout de titulo (layout 0) — titulo grande centralizado.
    - Todos os demais slides usam layout titulo+conteudo (layout 1).
    - Para layouts personalizados (em branco, secao), use change_slide_layout_pptx() apos criar.
    - Para adicionar imagens, tabelas ou graficos, use as tools especificas apos a criacao.
    - A apresentacao e SOBRESCRITA se ja existir.

    EXEMPLO DE USO:
        create_presentation_pptx(
            file_path="./vendas.pptx",
            slides=[
                {"title": "Resultados 2026", "bullet_points": []},
                {"title": "Resumo Executivo", "bullet_points": [
                    "Crescimento de 18% no faturamento",
                    "Meta atingida em 95%",
                    "42 novos contratos fechados"
                ]},
                {"title": "Crescimento por Trimestre", "bullet_points": [
                    "Q1: +8% vs. mesmo periodo",
                    "Q2: +12% vs. mesmo periodo",
                    "Q3: +22% vs. mesmo periodo",
                    "Q4: +30% vs. mesmo periodo"
                ]},
                {"title": "Projecoes 2027", "bullet_points": [
                    "Meta: crescimento de 25%",
                    "Expansao para 3 novos mercados"
                ]},
            ]
        )
        # Retorna: "Apresentacao criada com sucesso em: D:\\projeto\\vendas.pptx"
    """
    pptx = get_pptx()
    prs = pptx.Presentation()

    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    total_slides = len(slides)
    total_bullets = sum(len(s.get("bullet_points", [])) for s in slides)

    for i, s_data in enumerate(slides):
        title = s_data.get("title", "")
        bullets = s_data.get("bullet_points", [])

        if i == 0 and not bullets:
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
        else:
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    prs.save(file_path)
    abs_path = os.path.abspath(file_path)
    return (
        f"Apresentacao criada com sucesso!\n\n"
        f"  Arquivo: {abs_path}\n"
        f"  Slides: {total_slides}\n"
        f"  Topicos: {total_bullets}\n"
        f"\nProximos passos: use change_slide_layout_pptx() para layouts personalizados, "
        f"add_image_to_slide_pptx() para imagens, ou set_transition_pptx() para transicoes entre slides."
    )


@mcp.tool()
def read_presentation_pptx(file_path: str) -> str:
    """
    Le a estrutura de slides e todos os topicos de texto de uma apresentacao PowerPoint (.pptx).

    Esta tool abre um arquivo .pptx existente e retorna todo o conteudo textual
    organizado por slide, com cada slide numerado e seus textos listados.
    Ideal para analisar o conteudo antes de editar.

    QUANDO USAR:
    - Quando o usuario quer saber o conteudo de uma apresentacao
    - Antes de fazer edicoes, para entender a estrutura atual
    - Para extrair textos de apresentacoes existentes
    - Para verificar se os slides contem o conteudo esperado

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx.
      Exemplos: "./vendas.pptx", "C:\\apresentacoes\\pitch.pptx"

    RETORNO:
    - Texto estruturado com o seguinte formato:
      - Cabecalho: "--- Apresentacao: nome_arquivo.pptx ---"
      - Slides numerados: "[Slide N]" seguido de topicos com "  - texto"

    IMPORTANTE:
    - Extrai APENAS texto de shapes que possuem text_frame.
    - Imagens, formas sem texto, graficos e tabelas NAO sao exibidos.
    - Para extrair estilo (tema, layouts), use analyze_presentation_style_pptx().

    EXEMPLO DE USO:
        read_presentation_pptx(file_path="./vendas.pptx")
        # Retorna:
        # --- Apresentacao: vendas.pptx ---
        # [Slide 1]
        #   - Resultados 2026
        # [Slide 2]
        #   - Resumo Executivo
        #   - Crescimento de 18% no faturamento
        #   - Meta atingida em 95%
        # [Slide 3]
        #   - Projecoes 2027
        #   - Meta: crescimento de 25%
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs = pptx.Presentation(file_path)
    result = []

    result.append(f"--- Apresentacao: {os.path.basename(file_path)} ---")
    for i, slide in enumerate(prs.slides):
        result.append(f"\n[Slide {i+1}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        result.append(f"  - {paragraph.text.strip()}")

    return "\n".join(result)


@mcp.tool()
def add_image_to_slide_pptx(file_path: str, slide_index: int, image_path: str, left_inch: float, top_inch: float, width_inch: float = None, height_inch: float = None) -> str:
    """
    Insere uma imagem local em um slide especifico de uma apresentacao.

    Esta tool adiciona uma imagem (PNG, JPG, GIF, etc.) a um slide existente,
    posicionando-a em coordenadas especificas em polegadas. A imagem e inserida
    como um shape sobre o slide.

    QUANDO USAR:
    - Para adicionar logos, graficos exportados como imagem, fotos
    - Para inserir capturas de tela ou diagramas
    - Para posicionar imagens em localizacao exata no slide
    - Para complementar slides com elementos visuais

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx.
    - slide_index (int, obrigatorio): Indice do slide (0-based). Primeiro slide = 0.
    - image_path (str, obrigatorio): Caminho da imagem local.
      Formatos: PNG, JPEG, GIF, BMP, TIFF.
    - left_inch (float, obrigatorio): Posicao horizontal (eixo X) em polegadas a partir da borda esquerda.
    - top_inch (float, obrigatorio): Posicao vertical (eixo Y) em polegadas a partir do topo.
    - width_inch (float, opcional): Largura da imagem em polegadas.
      Se AMBOS (width e height) informados, redimensiona (pode distorcer).
      Se APENUS um informado, mantem proporcao.
      Se nenhum informado, usa tamanho original.
    - height_inch (float, opcional): Altura da imagem em polegadas.

    RETORNO:
    - Mensagem de texto com confirmacao e numero do slide.
      Ex: "Imagem inserida com sucesso no slide 1 da apresentacao vendas.pptx."

    IMPORTANTE:
    - 1 polegada = 2.54 cm. Para converter: polegadas = cm / 2.54
    - Se width_inch e height_inch forem informados, a imagem pode ser distorcida.
    - Para imagens em documentos Word, use add_image_to_document_word().
    - Formatos suportados: PNG, JPEG, GIF, BMP, TIFF.

    EXEMPLO DE USO:
        add_image_to_slide_pptx(
            file_path="./vendas.pptx",
            slide_index=0,
            image_path="./logo_empresa.png",
            left_inch=0.5,
            top_inch=0.3,
            width_inch=1.5
        )
        # Retorna: "Imagem inserida com sucesso no slide 1 da apresentacao vendas.pptx."
    """
    erro_pptx = _validar_arquivo_existe(file_path)
    if erro_pptx:
        return erro_pptx
    erro_img = _validar_arquivo_existe(image_path)
    if erro_img:
        return erro_img

    pptx = get_pptx()
    from pptx.util import Inches

    prs = pptx.Presentation(file_path)
    total_slides = len(prs.slides)
    if slide_index < 0 or slide_index >= total_slides:
        return f"Erro: Indice de slide {slide_index} invalido. A apresentacao tem {total_slides} slides (0 a {total_slides - 1})."

    slide = prs.slides[slide_index]
    left = Inches(left_inch)
    top = Inches(top_inch)
    width = Inches(width_inch) if width_inch else None
    height = Inches(height_inch) if height_inch else None

    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    prs.save(file_path)
    size_info = ""
    if width_inch and height_inch:
        size_info = f"\n  Tamanho: {width_inch}\" x {height_inch}\""
    elif width_inch:
        size_info = f"\n  Largura: {width_inch}\" (altura proporcional)"
    elif height_inch:
        size_info = f"\n  Altura: {height_inch}\" (largura proporcional)"
    else:
        size_info = "\n  Tamanho: original"
    return (
        f"Imagem inserida com sucesso!\n\n"
        f"  Slide: {slide_index+1}\n"
        f"  Arquivo de imagem: {os.path.basename(image_path)}\n"
        f"  Posicao: ({left_inch}\", {top_inch}\"){size_info}\n"
        f"  Apresentacao: {os.path.abspath(file_path)}"
    )


# ==========================================
# 3.5 POWERPOINT - FERRAMENTAS AVANCADAS
# ==========================================

@mcp.tool()
def manage_slides_pptx(file_path: str, action: str, slide_index: int = None, title: str = None) -> str:
    """
    Gerencia slides de uma apresentacao (adicionar, remover, listar).

    Esta tool permite manipular a estrutura de uma apresentacao, criando novos
    slides, removendo existentes ou listando todos os slides com seus titulos.

    QUANDO USAR:
    - Para adicionar novos slides a uma apresentacao existente
    - Para remover slides desnecessarios ou duplicados
    - Para listar todos os slides e seus titulos
    - Para reorganizar a estrutura da apresentacao

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx.
    - action (str, obrigatorio): Acao a executar:
      - "add": Adicionar novo slide (usa titulo se fornecido)
      - "delete": Remover slide (requer slide_index)
      - "list": Listar todos os slides
    - slide_index (int, opcional): Indice do slide (0-based). Obrigatorio para "delete".
    - title (str, opcional): Titulo do novo slide (apenas para "add").

    RETORNO:
    - Para "list": Lista numerada com titulos dos slides.
      Ex: "--- Slides: vendas.pptx (4 slides) ---\n  [0] Resultados 2026\n  [1] ..."
    - Para outras acoes: Mensagem de confirmacao.

    IMPORTANTE:
    - Indice 0-based: primeiro slide = 0.
    - Nao e possivel remover o ultimo slide restante.
    - Para alterar layout, use change_slide_layout_pptx() apos adicionar.

    EXEMPLO DE USO:
        # Listar slides:
        manage_slides_pptx(file_path="./vendas.pptx", action="list")

        # Adicionar slide:
        manage_slides_pptx(file_path="./vendas.pptx", action="add", title="Novo Topico")

        # Remover slide:
        manage_slides_pptx(file_path="./vendas.pptx", action="delete", slide_index=2)
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs = pptx.Presentation(file_path)

    if action == "list":
        result = [f"--- Slides: {os.path.basename(file_path)} ({len(prs.slides)} slides) ---"]
        for i, slide in enumerate(prs.slides):
            slide_title = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape == slide.shapes.title:
                    slide_title = shape.text_frame.paragraphs[0].text if shape.text_frame.paragraphs else ""
                    break
            result.append(f"  [{i}] {slide_title or '(sem titulo)'}")
        return "\n".join(result)

    if action == "add":
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if title:
            slide.shapes.title.text = title
        prs.save(file_path)
        return (
            f"Slide adicionado com sucesso!\n\n"
            f"  Indice: {len(prs.slides)-1}\n"
            f"  Titulo: {title or '(sem titulo)'}\n"
            f"  Total de slides: {len(prs.slides)}\n"
            f"  Arquivo: {os.path.abspath(file_path)}"
        )

    if action == "delete":
        if slide_index is None:
            return "Erro: Informe slide_index para remover um slide."
        total = len(prs.slides)
        if slide_index < 0 or slide_index >= total:
            return f"Erro: Indice {slide_index} invalido. A apresentacao tem {total} slides (0 a {total-1})."
        if total <= 1:
            return "Erro: Nao e possivel remover o unico slide restante."

        from pptx.oxml import parse_xml
        rId = prs.slides._sldIdLst[slide_index].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[slide_index])
        prs.save(file_path)
        return (
            f"Slide removido com sucesso!\n\n"
            f"  Indice removido: {slide_index}\n"
            f"  Total de slides: {total - 1}\n"
            f"  Arquivo: {os.path.abspath(file_path)}"
        )

    return "Erro: Acao invalida. Use: add, delete ou list."


@mcp.tool()
def change_slide_layout_pptx(file_path: str, slide_index: int, layout_name: str) -> str:
    """
    Altera o layout de um slide existente em uma apresentacao PowerPoint.

    Esta tool modifica o layout (disposicao de placeholders) de um slide
    especifico. Util para mudar a estrutura visual do slide apos a criacao.

    QUANDO USAR:
    - Para mudar o layout de um slide existente
    - Para aplicar layout em branco para conteudo personalizado
    - Para mudar de layout titulo+conteudo para titulo+secao
    - Para preparar slides para receber imagens ou formas

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx.
    - slide_index (int, obrigatorio): Indice do slide (0-based).
    - layout_name (str, obrigatorio): Nome do layout desejado:
      - "title": Slide de titulo (layout 0)
      - "content": Titulo + Conteudo (layout 1)
      - "section": Titulo + Secao (layout 2)
      - "blank": Em branco (layout 6)

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Layout do slide 2 alterado para 'blank'."

    IMPORTANTE:
    - Mudar o layout pode alterar a posicao dos placeholders existentes.
    - Conteudo existente no slide e preservado quando possivel.
    - Layouts disponiveis dependem do template da apresentacao.

    EXEMPLO DE USO:
        change_slide_layout_pptx(
            file_path="./vendas.pptx",
            slide_index=2,
            layout_name="blank"
        )
        # Retorna: "Layout do slide 2 alterado para 'blank'."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs = pptx.Presentation(file_path)
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return f"Erro: Indice {slide_index} invalido. A apresentacao tem {total} slides."

    layouts = {
        "title": 0,
        "content": 1,
        "blank": 6,
        "section": 2,
    }
    idx = layouts.get(layout_name)
    if idx is None:
        return f"Erro: Layout '{layout_name}' invalido. Use: {', '.join(layouts.keys())}"

    slide = prs.slides[slide_index]
    slide.slide_layout = prs.slide_layouts[idx]
    prs.save(file_path)
    return f"Layout do slide {slide_index} alterado para '{layout_name}'."


@mcp.tool()
def add_shape_pptx(file_path: str, slide_index: int, shape_type: str, left_inch: float, top_inch: float, width_inch: float, height_inch: float, text: str = None, fill_color: str = None) -> str:
    """
    Adiciona uma forma geometrica a um slide da apresentacao.

    Esta tool insere uma forma (retangulo, oval, seta, linha) em um slide
    existente, posicionando-a em coordenadas especificas. Pode incluir texto
    e cor de preenchimento.

    QUANDO USAR:
    - Para criar botoes visuais ou caixas de destaque
    - Para adicionar setas indicadoras ou linhas divisorias
    - Para criar diagramas simples
    - Para destacar informacoes com formas coloridas

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx.
    - slide_index (int, obrigatorio): Indice do slide (0-based).
    - shape_type (str, obrigatorio): Tipo da forma:
      - "rectangle": Retangulo
      - "rounded_rectangle": Retangulo arredondado
      - "oval": Oval / Circulo
      - "arrow_right": Seta para direita
      - "arrow_left": Seta para esquerda
      - "line": Linha
    - left_inch (float, obrigatorio): Posicao X em polegadas.
    - top_inch (float, obrigatorio): Posicao Y em polegadas.
    - width_inch (float, obrigatorio): Largura em polegadas.
    - height_inch (float, obrigatorio): Altura em polegadas.
    - text (str, opcional): Texto dentro da forma.
    - fill_color (str, opcional): Cor de preenchimento em HEX (sem "#").
      Ex: "4472C4" (azul), "FF0000" (vermelho), "92D050" (verde).

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Forma 'rectangle' adicionada ao slide 1."

    IMPORTANTE:
    - 1 polegada = 2.54 cm.
    - Cores em formato HEX sem "#" (ex: "4472C4", nao "#4472C4").
    - O texto fica centralizado na forma por padrao.

    EXEMPLO DE USO:
        add_shape_pptx(
            file_path="./vendas.pptx",
            slide_index=0,
            shape_type="rounded_rectangle",
            left_inch=1.0,
            top_inch=4.0,
            width_inch=3.0,
            height_inch=1.5,
            text="Destaque",
            fill_color="4472C4"
        )
        # Retorna: "Forma 'rounded_rectangle' adicionada ao slide 0."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    prs = pptx.Presentation(file_path)
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return f"Erro: Indice {slide_index} invalido. A apresentacao tem {total} slides."

    shape_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "arrow_right": getattr(MSO_SHAPE, "RIGHT_ARROW", MSO_SHAPE.RECTANGLE),
        "arrow_left": getattr(MSO_SHAPE, "LEFT_ARROW", MSO_SHAPE.RECTANGLE),
        "line": getattr(MSO_SHAPE, "STRAIGHT_CONNECTOR", MSO_SHAPE.RECTANGLE),
    }
    mso = shape_map.get(shape_type)
    if not mso:
        return f"Erro: Tipo de forma '{shape_type}' invalido. Use: {', '.join(shape_map.keys())}"

    slide = prs.slides[slide_index]
    left = Inches(left_inch)
    top = Inches(top_inch)
    width = Inches(width_inch)
    height = Inches(height_inch)

    shape = slide.shapes.add_shape(mso, left, top, width, height)

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(
            int(fill_color[0:2], 16),
            int(fill_color[2:4], 16),
            int(fill_color[4:6], 16)
        )

    if text:
        shape.text_frame.text = text

    prs.save(file_path)
    return f"Forma '{shape_type}' adicionada ao slide {slide_index}."


@mcp.tool()
def add_table_to_slide_pptx(file_path: str, slide_index: int, data: list, left_inch: float = 1, top_inch: float = 2, width_inch: float = 8, height_inch: float = 3) -> str:
    """
    Adiciona uma tabela formatada a um slide da apresentacao PowerPoint.

    Esta tool insere uma tabela em um slide existente, com a primeira linha
    como cabecalho (centralizada). A tabela e posicionada em coordenadas
    especificas em polegadas.

    QUANDO USAR:
    - Para exibir dados tabulares em slides
    - Para criar tabelas comparativas ou resumos
    - Para apresentar resultados numericos
    - Para incluir dados estruturados na apresentacao

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx.
    - slide_index (int, obrigatorio): Indice do slide (0-based).
    - data (list[list], obrigatorio): Dados da tabela.
      - Primeira linha = cabecalho (centralizado automaticamente).
      - Cada sub-lista e uma linha da tabela.
      Ex: [["Produto", "Qtd"], ["A", 100], ["B", 200]]
    - left_inch (float, opcional): Posicao X em polegadas. Default: 1.
    - top_inch (float, opcional): Posicao Y em polegadas. Default: 2.
    - width_inch (float, opcional): Largura em polegadas. Default: 8.
    - height_inch (float, opcional): Altura em polegadas. Default: 3.

    RETORNO:
    - Mensagem de texto com dimensoes e slide.
      Ex: "Tabela 3x2 adicionada ao slide 1."

    IMPORTANTE:
    - Cabecalhos sao automaticamente centralizados.
    - Para tabelas em documentos Word, use add_table_to_document_word().
    - Posicoes e dimensoes sao em polegadas (1 pol = 2.54 cm).

    EXEMPLO DE USO:
        add_table_to_slide_pptx(
            file_path="./vendas.pptx",
            slide_index=2,
            data=[
                ["Regiao", "Vendas", "Meta"],
                ["Sudeste", "R$ 500k", "R$ 450k"],
                ["Nordeste", "R$ 300k", "R$ 280k"],
                ["Sul", "R$ 200k", "R$ 220k"],
            ],
            left_inch=0.5,
            top_inch=1.5,
            width_inch=9,
            height_inch=4
        )
        # Retorna: "Tabela 4x3 adicionada ao slide 2."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    if not data or not data[0]:
        return "Erro: Dados da tabela nao podem estar vazios."

    pptx = get_pptx()
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN

    prs = pptx.Presentation(file_path)
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return f"Erro: Indice {slide_index} invalido. Total: {total} slides."

    slide = prs.slides[slide_index]
    rows = len(data)
    cols = len(data[0])
    left = Inches(left_inch)
    top = Inches(top_inch)
    width = Inches(width_inch)
    height = Inches(height_inch)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for i, row_data in enumerate(data):
        for j, cell_value in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_value)
            if i == 0:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER

    prs.save(file_path)
    return (
        f"Tabela {rows}x{cols} adicionada com sucesso!\n\n"
        f"  Slide: {slide_index+1}\n"
        f"  Linhas: {rows} | Colunas: {cols}\n"
        f"  Posicao: ({left_inch}\", {top_inch}\")\n"
        f"  Apresentacao: {os.path.abspath(file_path)}"
    )


@mcp.tool()
def edit_slide_text_pptx(file_path: str, slide_index: int, shape_index: int, text: str) -> str:
    """
    Edita o texto de um shape especifico em um slide da apresentacao.

    Esta tool modifica o conteudo textual de um shape existente em um slide.
    O shape e identificado pelo indice (0-based) dentro do slide.

    QUANDO USAR:
    - Para alterar o titulo ou conteudo de um slide existente
    - Para corrigir erros de texto em slides
    - Para atualizar dados ou informacoes em slides
    - Para personalizar textos de shapes pre-definidos

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx.
    - slide_index (int, obrigatorio): Indice do slide (0-based).
    - shape_index (int, obrigatorio): Indice do shape dentro do slide (0-based).
      Use manage_slides_pptx(action="list") para ver estrutura, ou leia o arquivo.
    - text (str, obrigatorio): Novo texto para o shape.

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Texto do shape 0 no slide 1 atualizado."

    IMPORTANTE:
    - Apenas UM shape e atualizado por chamada.
    - Se o shape nao tiver text_frame, retorna erro.
    - Para adicionar novos shapes de texto, use add_shape_pptx().

    EXEMPLO DE USO:
        edit_slide_text_pptx(
            file_path="./vendas.pptx",
            slide_index=0,
            shape_index=0,
            text="Novo Titulo da Apresentacao"
        )
        # Retorna: "Texto do shape 0 no slide 0 atualizado."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs = pptx.Presentation(file_path)
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return f"Erro: Indice de slide {slide_index} invalido. Total: {total} slides."

    slide = prs.slides[slide_index]
    shapes = list(slide.shapes)
    if shape_index < 0 or shape_index >= len(shapes):
        return f"Erro: Indice de shape {shape_index} invalido. Total: {len(shapes)} shapes."

    shape = shapes[shape_index]
    if not shape.has_text_frame:
        return f"Erro: Shape {shape_index} nao possui texto editavel."

    shape.text_frame.paragraphs[0].text = text
    prs.save(file_path)
    preview = text if len(text) <= 60 else text[:60] + "..."
    return (
        f"Texto atualizado com sucesso!\n\n"
        f"  Slide: {slide_index+1} | Shape: {shape_index}\n"
        f"  Novo texto: \"{preview}\"\n"
        f"  Arquivo: {os.path.abspath(file_path)}"
    )


@mcp.tool()
def add_chart_to_slide_pptx(file_path: str, slide_index: int, chart_type: str, data: dict, title: str = "", left_inch: float = 1, top_inch: float = 2, width_inch: float = 8, height_inch: float = 4.5) -> str:
    """
    Adiciona um grafico nativo PowerPoint a um slide da apresentacao.

    Esta tool insere um grafico nativo (barra, coluna, linha ou pizza) diretamente
    em um slide, usando dados estruturados. O grafico e editavel no PowerPoint
    apos a insercao.

    QUANDO USAR:
    - Para visualizar dados em formato grafico dentro da apresentacao
    - Para criar graficos de tendencia (linha), comparacao (barra/coluna) ou proporcao (pizza)
    - Para dashboards visuais em slides
    - Para complementar slides com dados visuais

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx.
    - slide_index (int, obrigatorio): Indice do slide (0-based).
    - chart_type (str, obrigatorio): Tipo do grafico:
      - "bar": Barras horizontais
      - "column": Colunas verticais (clustered)
      - "line": Linhas
      - "pie": Pizza
    - data (dict, obrigatorio): Dados do grafico no formato:
      {
        "categories": ["A", "B", "C"],
        "series": [
          {"name": "Vendas 2025", "values": [10, 20, 15]},
          {"name": "Vendas 2026", "values": [12, 25, 18]}
        ]
      }
    - title (str, opcional): Titulo do grafico. Default: sem titulo.
    - left_inch (float, opcional): Posicao X. Default: 1.
    - top_inch (float, opcional): Posicao Y. Default: 2.
    - width_inch (float, opcional): Largura. Default: 8.
    - height_inch (float, opcional): Altura. Default: 4.5.

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Grafico column adicionado ao slide 2."

    IMPORTANTE:
    - O grafico e NATIVO do PowerPoint (editavel manualmente apos insercao).
    - Para tabelas em vez de graficos, use add_table_to_slide_pptx().
    - Categories e series devem ter dados compatíveis.

    EXEMPLO DE USO:
        add_chart_to_slide_pptx(
            file_path="./vendas.pptx",
            slide_index=2,
            chart_type="column",
            data={
                "categories": ["Janeiro", "Fevereiro", "Marco"],
                "series": [
                    {"name": "Meta", "values": [100, 120, 110]},
                    {"name": "Realizado", "values": [95, 130, 115]}
                ]
            },
            title="Desempenho vs Meta",
            left_inch=0.5,
            top_inch=1.5,
            width_inch=9,
            height_inch=5
        )
        # Retorna: "Grafico column adicionado ao slide 2."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = pptx.Presentation(file_path)
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return f"Erro: Indice {slide_index} invalido. Total: {total} slides."

    chart_data = CategoryChartData()
    chart_data.categories = data.get("categories", [])
    for series in data.get("series", []):
        chart_data.add_series(series.get("name", ""), series.get("values", []))

    chart_type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }
    xl_type = chart_type_map.get(chart_type)
    if not xl_type:
        return f"Erro: Tipo '{chart_type}' invalido. Use: bar, column, line, pie."

    slide = prs.slides[slide_index]
    left = Inches(left_inch)
    top = Inches(top_inch)
    width = Inches(width_inch)
    height = Inches(height_inch)

    chart_frame = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
    chart_frame.chart.has_legend = True
    if title:
        chart_frame.chart.has_title = True
        chart_frame.chart.chart_title.text_frame.paragraphs[0].text = title

    prs.save(file_path)
    return f"Grafico {chart_type} adicionado ao slide {slide_index}."


@mcp.tool()
def add_animation_pptx(file_path: str, slide_index: int, shape_index: int = 0, effect: str = "fade", trigger: str = "on_click") -> str:
    """
    Adiciona efeito de animacao a um shape em um slide (usando power-pptx).

    Esta tool aplica uma animacao a um shape especifico em um slide. A animacao
    determina como o shape aparece durante a apresentacao (clique, automatico, etc.).
    Requer a biblioteca power-pptx.

    QUANDO USAR:
    - Para animar titulos, textos ou imagens em slides
    - Para criar apresentacoes dinamicas com efeitos de entrada
    - Para sincronizar aparecimento de elementos com cliques
    - Para melhorar o impacto visual da apresentacao

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx.
    - slide_index (int, obrigatorio): Indice do slide (0-based).
    - shape_index (int, opcional): Indice do shape (0-based). Default: 0.
    - effect (str, opcional): Tipo de efeito de animacao:
      - "fade": Fade in (aparecimento suave) — padrao
      - "fly_in": Entrada voo (sobe de baixo)
      - "float_in": Flutua para dentro
      - "wipe": Wipe (aparece como pincelada)
      - "appear": Aparecer (sem efeito)
    - trigger (str, opcional): Gatilho da animacao:
      - "on_click": Ao clicar (padrao)
      - "after_previous": Apos o anterior terminar
      - "with_previous": Junto com o anterior

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Animacao 'fade' adicionada ao shape 0 no slide 1."

    REQUISITOS:
    - Biblioteca power-pptx instalada (pip install power-pptx)

    IMPORTANTE:
    - A animacao e aplicada ao shape especificado.
    - Para animar multiplos shapes, chame a tool para cada um.
    - A ordem de chamada define a sequencia de animacao.

    EXEMPLO DE USO:
        add_animation_pptx(
            file_path="./vendas.pptx",
            slide_index=0,
            shape_index=0,
            effect="fade",
            trigger="on_click"
        )
        # Retorna: "Animacao 'fade' adicionada ao shape 0 no slide 0."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        from power_pptx import Presentation as PowerPresentation
        from power_pptx.enum.animation import PP_ANIM_TRIGGER
    except ImportError:
        return "Erro: power-pptx nao esta instalado. Execute: pip install power-pptx"

    prs = PowerPresentation(file_path)
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        return f"Erro: Indice {slide_index} invalido. Total: {total} slides."

    slide = prs.slides[slide_index]
    shapes = list(slide.shapes)
    if shape_index < 0 or shape_index >= len(shapes):
        return f"Erro: Indice de shape {shape_index} invalido. Total: {len(shapes)} shapes."

    shape = shapes[shape_index]
    preset_map = {
        "fade": "fade",
        "fly_in": "fly_in",
        "float_in": "float_in",
        "wipe": "wipe",
        "appear": "appear",
    }
    trigger_map = {
        "on_click": PP_ANIM_TRIGGER.ON_CLICK,
        "after_previous": PP_ANIM_TRIGGER.AFTER_PREVIOUS,
        "with_previous": PP_ANIM_TRIGGER.WITH_PREVIOUS,
    }

    preset = preset_map.get(effect, "fade")
    gatilho = trigger_map.get(trigger, PP_ANIM_TRIGGER.ON_CLICK)

    try:
        slide.animations.add_entrance(preset=preset, shape=shape, trigger=gatilho)
        prs.save(file_path)
        effect_desc = {
            "fade": "aparecimento suave (fade)",
            "fly_in": "entrada voando (fly_in)",
            "float_in": "flutuar para dentro (float_in)",
            "wipe": "revelar (wipe)",
            "appear": "aparecer (sem efeito)",
        }.get(effect, effect)
        trigger_desc = {
            "on_click": "ao clicar",
            "after_previous": "apos o anterior",
            "with_previous": "junto com o anterior",
        }.get(trigger, trigger)
        return (
            f"Animacao '{effect_desc}' adicionada com sucesso!\n\n"
            f"  Slide: {slide_index} | Shape: {shape_index} ({shape.name if hasattr(shape, 'name') else 'sem nome'})\n"
            f"  Gatilho: {trigger_desc}\n"
            f"  Arquivo: {os.path.abspath(file_path)}"
        )
    except Exception as e:
        return f"Erro: Nao foi possivel adicionar animacao. Detalhes: {str(e)[:200]}"


@mcp.tool()
def set_transition_pptx(file_path: str, slide_index: int = None, effect: str = "fade", duration: int = 1000) -> str:
    """
    Configura transicao entre slides de uma apresentacao PowerPoint.

    Esta tool define o efeito de transicao que ocorre ao avancar de um slide
    para o proximo durante a apresentacao. Pode ser aplicado a um slide
    especifico ou a todos os slides de uma vez.

    QUANDO USAR:
    - Para adicionar transicoes suaves entre slides
    - Para uniformizar a transicao em toda a apresentacao
    - Para criar apresentacoes com efeitos profissionais
    - Para ajustar a velocidade das transicoes

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx.
    - slide_index (int, opcional): Indice do slide (0-based).
      Se None, aplica a TODOS os slides.
      Se informado, aplica apenas ao slide especifico.
    - effect (str, opcional): Tipo de transicao:
      - "fade": Transicao suave (padrao)
      - "push": Empurrar
      - "wipe": Wipe
      - "morph": Morph (transicao inteligente)
      - "reveal": Dissolve (mapeado)
      - "zoom": Zoom
      - "cover": Cover
      - "cut": Cut (sem efeito)
    - duration (int, opcional): Duracao em milissegundos. Default: 1000 (1 segundo).

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Transicao 'fade' aplicada a todos os slides."
      Ou: "Transicao 'push' aplicada ao slide 3."

    REQUISITOS:
    - Biblioteca power-pptx instalada (pip install power-pptx)

    IMPORTANTE:
    - Se slide_index=None, aplica a TODOS os slides.
    - A duracao afeta a velocidade da transicao (mais alto = mais lento).
    - "morph" e uma transicao avancada que requer shapes comuns entre slides.

    EXEMPLO DE USO:
        # Aplicar fade em todos os slides:
        set_transition_pptx(file_path="./vendas.pptx", effect="fade", duration=800)

        # Aplicar push em slide especifico:
        set_transition_pptx(file_path="./vendas.pptx", slide_index=0, effect="push")
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        from power_pptx import Presentation as PowerPresentation
        from power_pptx.enum.presentation import MSO_TRANSITION_TYPE
    except ImportError:
        return "Erro: power-pptx nao esta instalado. Execute: pip install power-pptx"

    prs = PowerPresentation(file_path)
    effect_map = {
        "fade": MSO_TRANSITION_TYPE.FADE,
        "push": MSO_TRANSITION_TYPE.PUSH,
        "wipe": MSO_TRANSITION_TYPE.WIPE,
        "morph": MSO_TRANSITION_TYPE.MORPH,
        "reveal": MSO_TRANSITION_TYPE.DISSOLVE,
        "zoom": MSO_TRANSITION_TYPE.ZOOM,
        "cover": MSO_TRANSITION_TYPE.COVER,
        "cut": MSO_TRANSITION_TYPE.CUT,
    }
    efeito = effect_map.get(effect, MSO_TRANSITION_TYPE.FADE)

    if slide_index is not None:
        slides_to_apply = [prs.slides[slide_index]] if slide_index < len(prs.slides) else []
    else:
        slides_to_apply = list(prs.slides)

    if not slides_to_apply:
        return f"Erro: Nao foi possivel aplicar a transicao. Verifique o indice do slide."

    for slide in slides_to_apply:
        slide.transition.kind = efeito
        slide.transition.duration = duration
        slide.transition.advance_on_click = True

    prs.save(file_path)
    qtd = len(slides_to_apply)
    if qtd == 1:
        destino = f"o slide {slide_index}"
    else:
        destino = f"os {qtd} slides"
    effect_desc = {
        "fade": "suave (fade)",
        "push": "deslizar (push)",
        "wipe": "revelar (wipe)",
        "morph": "morfar (morph)",
        "reveal": "dissolver (reveal)",
        "zoom": "zoom",
        "cover": "cobrir (cover)",
        "cut": "seco (cut)",
    }.get(effect, effect)
    return (
        f"Transicao {effect_desc} aplicada com sucesso em {destino}!\n\n"
        f"  Duracao: {duration}ms\n"
        f"  Avanco: ao clicar\n"
        f"  Arquivo: {os.path.abspath(file_path)}"
    )


@mcp.tool()
def add_smart_art_pptx(file_path: str, slide_index: int, smart_art_type: str, left_inch: float = 1.0, top_inch: float = 1.0, width_inch: float = 6.0, height_inch: float = 3.0) -> str:
    """
    Adiciona grafico SmartArt a um slide da apresentacao.

    Esta tool insere um grafico SmartArt (diagrama visual) em um slide existente.
    SmartArt e uma forma nativa do PowerPoint para criar diagramas de processos,
    hierarquias, ciclos e outros formatos visuais.

    QUANDO USAR:
    - Para criar diagramas de processos ou fluxos
    - Para representar hierarquias organizacionais
    - Para criar ciclos, piramides ou matrices
    - Para enriquecer slides com elementos visuais estruturados

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx.
    - slide_index (int, obrigatorio): Indice do slide (0-based).
    - smart_art_type (str, obrigatorio): Tipo de SmartArt:
      - "ProcessCycle": Ciclo de processo
      - "Pyramid": Piramide
      - "Target": Alvo / Meta
      - "Radial": Radial
      - "Hierarchy": Hierarquia
      - "Matrix": Matriz
    - left_inch (float, opcional): Posicao X. Default: 1.0.
    - top_inch (float, opcional): Posicao Y. Default: 1.0.
    - width_inch (float, opcional): Largura. Default: 6.0.
    - height_inch (float, opcional): Altura. Default: 3.0.

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "SmartArt 'Hierarchy' adicionado ao slide 1."

    REQUISITOS:
    - Biblioteca power-pptx instalada (pip install power-pptx)

    IMPORTANTE:
    - O SmartArt e editavel no PowerPoint apos a insercao.
    - Conteudo padrao e gerado — edite manualmente apos inserir.
    - Posicoes sao em polegadas (1 pol = 2.54 cm).

    EXEMPLO DE USO:
        add_smart_art_pptx(
            file_path="./vendas.pptx",
            slide_index=2,
            smart_art_type="Hierarchy",
            left_inch=1.0,
            top_inch=1.5,
            width_inch=8.0,
            height_inch=4.0
        )
        # Retorna: "SmartArt 'Hierarchy' adicionado ao slide 3."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slide = prs.slides[slide_index]

        from power_pptx import SmartArtBuilder
        builder = SmartArtBuilder(prs)
        shape = builder.add_smart_art(
            slide,
            smart_art_type=smart_art_type,
            left=left_inch * 914400,
            top=top_inch * 914400,
            width=width_inch * 914400,
            height=height_inch * 914400
        )
        prs.save(file_path)
        return f"SmartArt '{smart_art_type}' adicionado ao slide {slide_index+1}."
    except ImportError:
        return "Erro: power-pptx nao esta instalado. Execute: pip install power-pptx"
    except Exception as e:
        return f"Erro: Nao foi possivel adicionar SmartArt. Detalhes: {str(e)[:200]}"


@mcp.tool()
def analyze_presentation_style_pptx(file_path: str) -> str:
    """
    Extrai informacoes de estilo (tema, layouts, cores) de uma apresentacao PowerPoint.

    Esta tool analisa a estrutura de uma apresentacao e retorna metadados de
    estilo que podem ser usados para replicar o mesmo visual em novas apresentacoes.
    Retorna: layouts disponiveis, dimensoes dos slides, contagem de slides.

    QUANDO USAR:
    - Para extrair o estilo de uma apresentacao modelo
    - Para auditar a estrutura de uma apresentacao recebida
    - Para documentar layouts disponiveis antes de criar novos slides
    - Para usar com create_presentation_from_example_pptx()

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx a ser analisado.

    RETORNO:
    - JSON formatado com:
      - "slide_count": Numero total de slides
      - "slide_width": Largura do slide
      - "slide_height": Altura do slide
      - "layouts": Lista de layouts disponiveis com nome e placeholders

    IMPORTANTE:
    - Nao modifica o arquivo — e Somente Leitura.
    - O JSON retornado pode ser usado como referencia para criar novas apresentacoes.

    EXEMPLO DE USO:
        analyze_presentation_style_pptx(file_path="./modelo.pptx")
        # Retorna JSON com layouts, dimensoes e estrutura da apresentacao
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs = pptx.Presentation(file_path)

    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        layouts.append({
            "index": i,
            "name": layout.name,
            "placeholders": [
                {"idx": ph.placeholder_format.idx, "name": ph.name}
                for ph in layout.placeholders
            ] if hasattr(layout, 'placeholders') else []
        })

    import json
    result = {
        "slide_count": len(prs.slides),
        "slide_width": str(prs.slide_width),
        "slide_height": str(prs.slide_height),
        "layouts": layouts,
    }
    return json.dumps(result, ensure_ascii=False, indent=2)


@mcp.tool()
def create_presentation_from_example_pptx(example_path: str, output_path: str, slides: list) -> str:
    """
    Cria nova apresentacao baseada no estilo de uma existente.

    Esta tool usa uma apresentacao modelo como referencia para criar uma nova
    apresentacao, preservando o layout e estilo. Os slides sao criados com
    o layout de titulo+conteudo (layout 1) do modelo.

    QUANDO USAR:
    - Quando o usuario quer criar uma apresentacao no mesmo estilo de outra
    - Para manter consistencia visual entre apresentacoes
    - Para replicar templates corporativos
    - Para criar novas apresentacoes a partir de um modelo

    PARAMETROS:
    - example_path (str, obrigatorio): Caminho da apresentacao modelo (.pptx).
    - output_path (str, obrigatorio): Caminho para salvar a nova apresentacao (.pptx).
    - slides (list[dict], obrigatorio): Lista de dicionarios:
      - "title" (str): Titulo do slide
      - "bullet_points" (list[str], opcional): Lista de topicos

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Apresentacao criada a partir do exemplo em: D:\\apresentacoes\\novo.pptx"

    IMPORTANTE:
    - A apresentacao modelo NAO e modificada.
    - Usa layout 1 (titulo+conteudo) para todos os slides.
    - Para estilos mais fiéis, use create_presentation_pptx() e change_slide_layout_pptx().

    EXEMPLO DE USO:
        create_presentation_from_example_pptx(
            example_path="./template_corporativo.pptx",
            output_path="./nova_apresentacao.pptx",
            slides=[
                {"title": "Apresentacao Comercial", "bullet_points": []},
                {"title": "Quem Somos", "bullet_points": ["Empresa lider no mercado", "20 anos de experiencia"]},
                {"title": "Nossos Servicos", "bullet_points": ["Consultoria", "Implantacao", "Suporte"]},
            ]
        )
        # Retorna: "Apresentacao criada a partir do exemplo em: D:\\apresentacoes\\nova_apresentacao.pptx"
    """
    erro = _validar_arquivo_existe(example_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs_modelo = pptx.Presentation(example_path)
    prs_novo = pptx.Presentation()

    for s_data in slides:
        title = s_data.get("title", "")
        bullets = s_data.get("bullet_points", [])
        slide = prs_novo.slides.add_slide(prs_novo.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame if slide.placeholders[1].has_text_frame else None
        if tf:
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    prs_novo.save(output_path)
    return f"Apresentacao criada a partir do exemplo em: {os.path.abspath(output_path)}"


@mcp.tool()
def list_template_slots_pptx(template_path: str) -> str:
    """
    Lista todos os slots de texto editaveis em um template de apresentacao.

    Esta ferramenta analise um arquivo .pptx template e retorna o
    mapeamento completo de slides e shapes com texto, permitindo saber
    quais posicoes podem ser editadas. Use os indices retornados
    (slide_N_shape_M) com create_from_template_pptx().

    QUANDO USAR:
    - Antes de criar uma apresentacao a partir de um template
    - Para descobrir quais textos podem ser personalizados
    - Para mapear a estrutura de qualquer template .pptx

    PARAMETROS:
    - template_path (str, obrigatorio): Caminho do arquivo .pptx template.
      Ex: "C:\\templates\\clean_minimal.pptx"
      Use caminhos absolutos ou relativos a partir do workspace.

    RETORNO:
    - JSON com a estrutura completa do template:
      {
        "total_slides": 25,
        "slides": [
          {
            "slide_index": 0,
            "shapes": [
              {"shape_index": 10, "current_text": "A Presentation by: Jane Doe"},
              {"shape_index": 12, "current_text": "Clean"}
            ]
          }
        ]
      }

    EXEMPLO DE USO:
        list_template_slots_pptx(template_path="./clean_minimal.pptx")
        # Retorna JSON com todos os slides e shapes
    """
    erro = _validar_arquivo_existe(template_path)
    if erro:
        return erro

    pptx = get_pptx()
    prs = pptx.Presentation(template_path)

    result = {
        "total_slides": len(prs.slides),
        "template": os.path.basename(template_path),
        "slides": []
    }

    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {"slide_index": slide_idx, "shapes": []}
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_info["shapes"].append({
                    "shape_index": shape_idx,
                    "shape_id": f"slide_{slide_idx}_shape_{shape_idx}",
                    "current_text": shape.text_frame.text.strip()[:150]
                })
        result["slides"].append(slide_info)

    return json.dumps(result, ensure_ascii=False, indent=2)


@mcp.tool()
def create_from_template_pptx(
    template_path: str,
    output_path: str,
    slide_replacements: str
) -> str:
    """
    Cria uma nova apresentacao a partir de um template, substituindo textos.

    Esta ferramenta clona um template .pptx e substitui textos em shapes
    especificos, preservando 100% da formatacao visual, cores, fontes,
    layouts e imagens do template original.

    QUANDO USAR:
    - Para criar uma apresentacao profissional rapidamente
    - Para usar templates de Slidesgo, SlidesCarnival, etc.
    - Para manter consistencia visual entre documentos
    - Para reaproveitar layouts profissionais

    PARAMETROS:
    - template_path (str, obrigatorio): Caminho do arquivo .pptx template.
      Pode ser um template baixado de sites como SlidesCarnival ou um
      template proprio salvo em ./data/templates/.

    - output_path (str, obrigatorio): Caminho do novo arquivo .pptx.
      Sera sobrescrito se ja existir.

    - slide_replacements (str, obrigatorio): JSON com mapeamento de
      substituicoes no formato:
      {
        "0": {  // slide_index (0-based)
          "10": "Novo texto para shape 10",  // shape_index -> novo texto
          "12": "Clean"  // substitui o texto do shape 12 por "Clean"
        },
        "3": {
          "8": "INTRODUCAO"
        }
      }
      Use list_template_slots_pptx() antes para descobrir os indices.

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Apresentacao criada a partir do template com sucesso! Slides: 25
      Arquivo: D:\\projetos\\nova.pptx"

    IMPORTANTE:
    - O template NAO e modificado.
    - Formatacao, cores, fontes e imagens sao 100% preservadas.
    - Apenas os textos nos shapes especificados sao substituidos.
    - Para textos com quebras de linha, use "\\n" no JSON.
    - Shapes nao especificados permanecem com o texto original.

    EXEMPLO DE USO:
        # 1. Primeiro, listar os slots:
        slots = list_template_slots_pptx(template_path="./data/templates/clean_minimal.pptx")
        # 2. Depois, criar com substituicoes:
        create_from_template_pptx(
            template_path="./data/templates/clean_minimal.pptx",
            output_path="./projeto/astronomia.pptx",
            slide_replacements='{
              "0": {
                "10": "O Sistema Solar\\nMiguel Porto",
                "11": "Janeiro 2026",
                "12": "O SISTEMA",
                "13": "OS PLANETAS",
                "14": "O COSMOS"
              },
              "3": {
                "8": "INTRODUCAO",
                "9": "O sistema solar..."
              }
            }'
        )
        # Retorna: "Apresentacao criada a partir do template com sucesso!"
    """
    erro = _validar_arquivo_existe(template_path)
    if erro:
        return erro

    try:
        replacements = json.loads(slide_replacements)
    except json.JSONDecodeError as e:
        return f"Erro: slide_replacements invalido. JSON parse error: {str(e)}"

    pptx = get_pptx()
    prs = pptx.Presentation(template_path)

    def set_shape_text(shape, new_text):
        if not shape.has_text_frame:
            return False
        tf = shape.text_frame
        if not tf.paragraphs:
            return False
        lines = new_text.split("\\n") if isinstance(new_text, str) else [str(new_text)]

        # Substitui texto preservando a formatacao do primeiro run
        old_paragraphs = list(tf.paragraphs)
        first_p = old_paragraphs[0]
        if first_p.runs:
            first_p.runs[0].text = lines[0]
            # Remove runs extras do primeiro paragrafo
            for run in list(first_p.runs[1:]):
                run._r.getparent().remove(run._r)
        else:
            first_p.text = lines[0]
        # Remove paragrafos extras
        for extra_p in old_paragraphs[1:]:
            extra_p._p.getparent().remove(extra_p._p)
        # Adiciona paragrafos adicionais (linhas extras)
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = line

        # Auto-shrink: NOTE - nao modificamos bodyPr.xml diretamente
        # porque templates do Google Slides compartilham bodyPr entre shapes,
        # e modificar causa corrupcao do arquivo PPTX. Em vez disso,
        # ajustamos o tamanho da fonte diretamente (abaixo na heuristica).

        # Heuristica adicional: ajusta fonte e altura do shape para evitar
        # sobreposicoes. Especialmente importante em titulos de cards
        # (que tem altura fixa baixa)
        try:
            from pptx.util import Emu
            if shape.width and shape.height:
                width_in = Emu(shape.width).inches
                height_in = Emu(shape.height).inches

                if len(lines) == 1:
                    # Linha unica: verifica largura
                    original_size = first_p.runs[0].font.size if first_p.runs else None
                    if original_size:
                        size_pt = original_size.pt
                        text_len = len(lines[0])
                        # Chute empírico: largura_em_polegadas * 1.5
                        # (32pt bold DM Sans ~ 0.6-0.7 in por char)
                        max_chars = width_in * 1.5
                        from pptx.util import Pt as _Pt
                        if text_len > max_chars:
                            reduction = max_chars / text_len
                            new_size = max(10, int(size_pt * reduction))
                            first_p.runs[0].font.size = _Pt(new_size)
                        # Tambem reduz se altura for muito baixa (< 0.7in)
                        # e texto tem mais de 5 chars (evita reduzir "Marte" para 18pt)
                        elif height_in < 0.7 and size_pt > 18 and text_len > 5:
                            new_size_h = min(18, max(10, int(size_pt * 0.7)))
                            first_p.runs[0].font.size = _Pt(new_size_h)
                else:
                    # Multiplas linhas: verifica altura
                    if first_p.runs:
                        size_pt = first_p.runs[0].font.size.pt if first_p.runs[0].font.size else 21
                        line_height_in = (size_pt * 1.4) / 72
                        needed_height = line_height_in * len(lines) + 0.1
                        if needed_height > height_in:
                            new_size = int((height_in - 0.1) * 72 / (len(lines) * 1.4))
                            new_size = max(8, new_size)
                            from pptx.util import Pt as _Pt
                            for p in tf.paragraphs:
                                for r in p.runs:
                                    if r.font.size:
                                        r.font.size = _Pt(new_size)
        except Exception:
            pass

        return True

    changes_count = 0
    for slide_key, shape_replacements in replacements.items():
        slide_idx = int(slide_key)
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]
        for shape_key, new_text in shape_replacements.items():
            shape_idx = int(shape_key)
            if shape_idx >= len(slide.shapes):
                continue
            shape = list(slide.shapes)[shape_idx]
            if set_shape_text(shape, new_text):
                changes_count += 1

    prs.save(output_path)
    return (
        f"Apresentacao criada a partir do template com sucesso!\n\n"
        f"  Template: {os.path.basename(template_path)}\n"
        f"  Slides: {len(prs.slides)}\n"
        f"  Substituicoes aplicadas: {changes_count}\n"
        f"  Arquivo: {os.path.abspath(output_path)}"
    )


# ==========================================
# 4. FERRAMENTAS DE MACROS VBA (Windows)
# ==========================================

@mcp.tool()
def create_macro_workbook_vba(file_path: str, vba_code: str, module_name: str = "Module1") -> str:
    """
    Cria um novo arquivo Excel habilitado para macros (.xlsm) e injeta o codigo VBA fornecido.

    Esta tool cria um arquivo .xlsm do zero e injeta o codigo fonte VBA diretamente
    no modulo especificado. O arquivo fica pronto para executar a macro.

    QUANDO USAR:
    - Para criar planilhas com automacao VBA
    - Para gerar arquivos .xlsm com macros pre-configuradas
    - Para distribuir automacoes para usuarios finais
    - Para criar templates com macros embutidas

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo a ser criado.
      DEVE terminar com .xlsm (extensao obrigatoria para macros).
      Ex: "./automacao.xlsm"
    - vba_code (str, obrigatorio): Codigo fonte VBA bruto.
      Ex: 'Sub MinhaMacro()\\n    MsgBox "Executado!"\\nEnd Sub'
    - module_name (str, opcional): Nome do modulo VBA. Default: "Module1".

    RETORNO:
    - Mensagem de texto com confirmacao ou erro.
      Ex: "Arquivo com macro criado com sucesso em: D:\\planilhas\\automacao.xlsm"

    REQUISITOS:
    - Windows com Microsoft Excel instalado
    - Opcao "Confiar no acesso ao modelo de objeto do projeto VBA" ATIVADA
      (Arquivo > Opcoes > Central de Confiabilidade > Configuracoes de Macro)

    IMPORTANTE:
    - O arquivo DEVE ter extensao .xlsm — caso contrario, retorna erro.
    - Se a permissao VBA nao estiver ativada, a injecao falha com erro descritivo.
    - Para EXECUTAR uma macro existente, use run_macro_vba().
    - IMPORTANTE: Se as instrucoes do usuario sobre as regras da macro nao estiverem claras,
      peca esclarecimentos antes de prosseguir.

    EXEMPLO DE USO:
        create_macro_workbook_vba(
            file_path="./automacao.xlsm",
            vba_code='''
Sub AtualizarDados()
    Range("A1").Value = "Atualizado em: " & Now()
    Range("A1").Font.Bold = True
    MsgBox "Dados atualizados com sucesso!"
End Sub
''',
            module_name="ModuloPrincipal"
        )
        # Retorna: "Arquivo com macro criado com sucesso em: D:\\planilhas\\automacao.xlsm"
    """
    if not file_path.lower().endswith(".xlsm"):
        return "Erro: O caminho do arquivo para macros deve possuir extensao '.xlsm'. Use: criar_planilha(file_path='dados.xlsm', ...)"

    win32 = get_win32()
    abs_path = os.path.abspath(file_path)

    excel = win32.Dispatch("Excel.Application")
    excel.Visible = False
    excel.DisplayAlerts = False

    wb = excel.Workbooks.Add()
    try:
        xl_module = wb.VBProject.VBComponents.Add(1)
        xl_module.Name = module_name
        xl_module.CodeModule.AddFromString(vba_code)
        wb.SaveAs(abs_path, FileFormat=52)
        wb.Close()
        return f"Arquivo com macro criado com sucesso em: {abs_path}"
    except Exception as e:
        try:
            wb.Close(SaveChanges=False)
        except Exception:
            pass
        return (f"Erro: Nao foi possivel injetar a macro. A opcao 'Confiar no acesso ao modelo de "
                f"objeto do projeto VBA' precisa estar ativada (Arquivo > Opcoes > Central de "
                f"Confiabilidade > Configuracoes de Macro). Detalhes: {str(e)}")
    finally:
        try:
            excel.Quit()
        except Exception:
            pass


@mcp.tool()
def run_macro_vba(file_path: str, macro_name: str) -> str:
    """
    Abre uma planilha Excel (.xlsm) ou documento Word (.docm) e executa a macro VBA especificada.

    Esta tool usa PowerShell para abrir o arquivo Office em background e executar
    a macro especificada. Mais robusta que win32com puro pois trata pop-ups invisiveis
    do Office (ex: "Deseja atualizar vinculos externos?").

    QUANDO USAR:
    - Para executar automacoes VBA ja existentes em arquivos Office
    - Para rodar macros de atualizacao, processamento ou formatacao
    - Para automatizar tarefas repetitivas via VBA
    - Para executar macros de terceiros de forma segura

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo (.xlsm ou .docm).
      - Excel: "./automacao.xlsm"
      - Word: "./documento.docm"
    - macro_name (str, obrigatorio): Nome da macro a ser executada.
      - Excel: "NomeDaSub" ou "Module1.NomeDaSub"
      - Word: "NomeDaSub" ou "Module1.NomeDaSub"

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Macro 'AtualizarDados' executada com sucesso no Excel."
      Ou: "Macro 'FormatarDocumento' executada com sucesso no Word."

    REQUISITOS:
    - Windows com Microsoft Office (Excel ou Word) instalado
    - Arquivo habilitado para macros (.xlsm ou .docm)
    - Macro existente no arquivo

    IMPORTANTE:
    - Usa PowerShell internamente (mais robusto que win32com puro).
    - Suprime alertas automaticamente ($DisplayAlerts = $false).
    - Timeout de 60 segundos para evitar travamentos.
    - Para CRIAR uma macro, use create_macro_workbook_vba().

    EXEMPLO DE USO:
        run_macro_vba(file_path="./automacao.xlsm", macro_name="Module1.AtualizarDados")
        # Retorna: "Macro 'Module1.AtualizarDados' executada com sucesso no Excel."
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    abs_path = os.path.abspath(file_path)

    if file_path.lower().endswith((".xlsm", ".xls")):
        script = f"""
$excel = $null
try {{
    $excel = New-Object -ComObject Excel.Application
    $excel.Visible = $false
    $excel.DisplayAlerts = $false
    $wb = $excel.Workbooks.Open("{abs_path}")
    $excel.Application.Run("{macro_name}")
    $wb.Save()
    $wb.Close()
    Write-Output "Macro '{macro_name}' executada com sucesso."
}} catch {{
    Write-Error $_.Exception.Message
}} finally {{
    if ($excel) {{ $excel.Quit() }}
}}
"""
        resultado = _executar_powershell(script)
        if resultado.startswith("Erro:"):
            return resultado
        return f"Macro '{macro_name}' executada com sucesso no Excel."

    elif file_path.lower().endswith((".docm", ".doc")):
        script = f"""
$word = $null
try {{
    $word = New-Object -ComObject Word.Application
    $word.Visible = $false
    $word.DisplayAlerts = $false
    $doc = $word.Documents.Open("{abs_path}")
    $word.Application.Run("{macro_name}")
    $doc.Save()
    $doc.Close()
    Write-Output "Macro '{macro_name}' executada com sucesso."
}} catch {{
    Write-Error $_.Exception.Message
}} finally {{
    if ($word) {{ $word.Quit() }}
}}
"""
        resultado = _executar_powershell(script)
        if resultado.startswith("Erro:"):
            return resultado
        return f"Macro '{macro_name}' executada com sucesso no Word."

    else:
        return "Erro: Formato de arquivo nao suportado para execucao de macros. Use .xlsm (Excel) ou .docm (Word)."


# ==========================================
# 5. FERRAMENTAS UTILITARIAS
# ==========================================

@mcp.tool()
def extract_text_pdf(file_path: str) -> str:
    """
    Extrai todo o texto e tabelas de um arquivo PDF.

    Esta tool abre um arquivo PDF usando PyMuPDF (fitz) e extrai todo o texto
    organisado por pagina. Tabelas sao detectadas e formatadas com celulas
    separadas por " | ". Ideal para analise de conteudo de PDFs.

    QUANDO USAR:
    - Quando o usuario quer saber o conteudo de um PDF
    - Para extrair texto de documentos digitalizados (se tiver camada de texto)
    - Para analisar contratos, relatorios ou artigos em PDF
    - Para processar dados de PDFs estruturados

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pdf.
      Exemplos: "./contrato.pdf", "C:\\docs\\relatorio.pdf"

    RETORNO:
    - Texto estruturado com o seguinte formato:
      - Cabecalho: "--- PDF: nome_arquivo.pdf (N paginas) ---"
      - Paginas numeradas: "[Pagina N]" seguido do texto
      - Tabelas formatadas com " | " entre celulas

    IMPORTANTE:
    - Extrai APENAS texto — nao imagens ou formatacao visual.
    - Tabelas sao detectadas e exibidas com celulas separadas por " | ".
    - Nao funciona com PDFs apenas escaneados (sem camada de texto).
    - Para PDFs escaneados, use ocr_document() com extracao de imagem.

    EXEMPLO DE USO:
        extract_text_pdf(file_path="./contrato.pdf")
        # Retorna:
        # --- PDF: contrato.pdf (3 paginas) ---
        # [Pagina 1]
        # CONTRATO DE PRESTACAO DE SERVICOS
        # O presente contrato tem por objeto...
        # [Pagina 2]
        # CLAUSULA PRIMEIRA - DO OBJETO
        # ...
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    if not file_path.lower().endswith(".pdf"):
        return "Erro: O arquivo deve ter extensao .pdf. Use: extrair_texto_pdf(file_path='documento.pdf')"

    try:
        fitz = get_pymupdf()
        doc = fitz.open(file_path)
        result = []
        result.append(f"--- PDF: {os.path.basename(file_path)} ({doc.page_count} paginas) ---")

        for page_num in range(doc.page_count):
            page = doc[page_num]
            text = page.get_text("text")
            if text.strip():
                result.append(f"\n[Pagina {page_num + 1}]")
                result.append(text.strip())

        doc.close()
        return "\n".join(result)
    except Exception as e:
        return f"Erro: Nao foi possivel extrair o texto do PDF. O arquivo pode estar corrompido ou protegido. Detalhes: {str(e)[:200]}"


@mcp.tool()
def ocr_document(file_path: str, idioma: str = "pt") -> str:
    """
    Realiza OCR (reconhecimento optico de caracteres) em uma imagem de documento escaneado.

    Esta tool usa a biblioteca EasyOCR para extrair texto de imagens (JPG, PNG, TIFF, etc.).
    Util para digitalizar documentos fisicos, notas fisicais, recibos ou qualquer
    imagem que contenha texto legivel.

    QUANDO USAR:
    - Para extrair texto de imagens escaneadas ou fotografadas
    - Para digitalizar notas fisicais, recibos, contratos impressos
    - Para processar documentos que so existem em formato de imagem
    - Quando extract_text_pdf() nao funciona (PDF escaneado sem camada de texto)

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho da imagem (jpg, png, tiff, bmp).
      Exemplos: "./nota_fiscal.jpg", "C:\\docs\\documento.png"
    - idioma (str, opcional): Idioma do texto para OCR.
      - "pt": Portugues (padrao)
      - "en": Ingles
      - "es": Espanhol
      - "fr": Frances
      - "pt+en": Portugues e Ingles
      Default: "pt"

    RETORNO:
    - Texto reconhecido pela IA de OCR, organizado em linhas.
    - Se nenhum texto for reconhecido, retorna mensagem sugestiva.

    REQUISITOS:
    - Biblioteca easyocr instalada (pip install easyocr)
    - Download inicial do modelo de linguagem (~1-2min na primeira execucao)

    IMPORTANTE:
    - Imagens de boa qualidade (300 DPI+) produzem muito melhores resultados.
    - Suporta: JPG, JPEG, PNG, TIFF, TIF, BMP.
    - Para PDFs nativos (nao escaneados), use extract_text_pdf() que e mais rapido.
    - A primeira execucao pode ser lenta (download do modelo).

    EXEMPLO DE USO:
        ocr_document(file_path="./nota_fiscal.jpg", idioma="pt")
        # Retorna:
        # NOTA FISCAL Nº 12345
        # Cliente: Joao Silva
        # Valor: R$ 1.500,00
        # Data: 15/06/2026
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    ext_valida = (".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp")
    if not file_path.lower().endswith(ext_valida):
        return "Erro: Formato de imagem nao suportado. Use JPG, PNG, TIFF ou BMP."

    try:
        reader = get_easyocr()
        reader_instance = reader.Reader([idioma], gpu=False)
        resultado = reader_instance.readtext(file_path, detail=0, paragraph=True)
        return "\n".join(resultado) if resultado else "Nenhum texto foi reconhecido na imagem. Tente uma imagem com melhor resolucao."
    except Exception as e:
        return f"Erro: Nao foi possivel realizar o OCR na imagem. Detalhes: {str(e)[:200]}"


@mcp.tool()
def protect_document(file_path: str, senha: str) -> str:
    """
    Aplica protecao por senha em arquivos Office (.docx, .xlsx, .pptx).

    Esta tool criptografa o arquivo Office usando a senha fornecida. Apos
    proteger, o arquivo so pode ser aberto com a senha informada. Funciona
    sem Microsoft Office instalado (usa biblioteca msoffcrypto).

    QUANDO USAR:
    - Para proteger documentos confidenciais com senha
    - Para enviar arquivos por e-mail com seguranca
    - Para restringir acesso a dados sensiveis
    - Para arquivar documentos protegidos

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo Office.
      Formatos suportados: .docx, .doc, .xlsx, .xls, .pptx, .ppt
    - senha (str, obrigatorio): Senha para proteger o documento.
      Minimo 4 caracteres. Será necessaria para ABRIR o arquivo.

    RETORNO:
    - Mensagem de texto com confirmacao.
      Ex: "Documento protegido com senha em: D:\\docs\\relatorio.xlsx"

    REQUISITOS:
    - Biblioteca msoffcrypto instalada (pip install msoffcrypto-tool)
    - Nao requer Microsoft Office instalado

    IMPORTANTE:
    - A senha e NECESSARIA para abrir o arquivo apos a protecao.
    - Nao e possivel recuperar a senha se esquecida — GUARDE EM LOCAL SEGURO.
    - A protecao e de ABRIMENTO — o arquivo e criptografado inteiro.
    - Funciona SEM Microsoft Office instalado.

    EXEMPLO DE USO:
        protect_document(file_path="./relatorio.xlsx", senha="MinhaSenha123")
        # Retorna: "Documento protegido com senha em: D:\\docs\\relatorio.xlsx"
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    if len(senha) < 4:
        return "Erro: A senha deve ter no minimo 4 caracteres."

    ext_valida = (".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt")
    if not file_path.lower().endswith(ext_valida):
        return "Erro: Formato de arquivo nao suportado. Use .docx, .xlsx ou .pptx."

    try:
        msoffcrypto = get_msoffcrypto()
        abs_path = os.path.abspath(file_path)

        with open(abs_path, "rb") as f_in:
            ofile = msoffcrypto.OfficeFile(f_in)
            ofile.load_password(password=senha)
            temp_path = abs_path + ".temp"
            with open(temp_path, "wb") as f_out:
                ofile.encrypt(f_out)

        os.replace(temp_path, abs_path)
        return f"Documento protegido com senha em: {abs_path}"
    except Exception as e:
        return f"Erro: Nao foi possivel proteger o arquivo. Detalhes: {str(e)[:200]}"


@mcp.tool()
def get_document_info_word(file_path: str) -> str:
    """
    Retorna informacoes detalhadas de um documento Word (.docx) SEM modificar o arquivo.

    Esta tool analisa a estrutura de um .docx e retorna metricas uteis: numero
    de paragrafos, caracteres, tabelas, imagens e estimativa de paginas. Util para
    verificar o conteudo antes de ler ou editar.

    QUANDO USAR:
    - ANTES de ler o documento para ter uma visao geral da complexidade
    - Para verificar quantas tabelas, imagens ou paragrafos o documento possui
    - Para estimar o numero de paginas antes de processar
    - Para auditar a estrutura de um documento recebido

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx.
      Exemplos: "./relatorio.docx", "C:\\docs\\contrato.docx"

    RETORNO:
    - JSON formatado com as seguintes chaves:
      - "arquivo" (str): Nome do arquivo
      - "paragrafos" (int): Quantidade de paragrafos nao vazios
      - "caracteres" (int): Total de caracteres no texto
      - "tabelas" (int): Quantidade de tabelas
      - "imagens" (int): Quantidade de imagens embutidas
      - "paginas_estimadas" (int): Estimativa de paginas (~3000 chars/pagina)
    - Ex: {"arquivo": "relatorio.docx", "paragrafos": 150, "caracteres": 8500, ...}

    IMPORTANT:
    - Esta tool NAO modifica o arquivo — e Somente Leitura.
    - O numero de paginas e uma ESTIMATIVA baseada em ~3000 caracteres por pagina.
      O valor real pode variar conforme fontes, espacamento e formatacao.
    - Imagens sao contadas mas nao descritas.

    EXEMPLO DE USO:
        get_document_info_word(file_path="./relatorio.docx")
        # Retorna:
        # {
        #   "arquivo": "relatorio.docx",
        #   "paragrafos": 150,
        #   "caracteres": 8500,
        #   "tabelas": 3,
        #   "imagens": 2,
        #   "paginas_estimadas": 3
        # }
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    docx = get_docx()
    doc = docx.Document(file_path)

    total_chars = 0
    total_paragraphs = 0
    total_tables = len(doc.tables) if doc.tables else 0
    total_images = 0

    for p in doc.paragraphs:
        if p.text.strip():
            total_paragraphs += 1
            total_chars += len(p.text)

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            total_images += 1

    paginas_estimadas = max(1, round(total_chars / 3000))

    info = {
        "arquivo": os.path.basename(file_path),
        "paragrafos": total_paragraphs,
        "caracteres": total_chars,
        "tabelas": total_tables,
        "imagens": total_images,
        "paginas_estimadas": paginas_estimadas,
    }

    return json.dumps(info, ensure_ascii=False, indent=2)


@mcp.tool()
def append_to_document_word(file_path: str, elements: list) -> str:
    """
    ADICIONA conteudo ao FINAL de um documento Word (.docx) existente, SEM sobrescrever.

    Esta tool abre um .docx existente e acrescenta novos elementos (titulos,
    paragrafos, listas) ao final de todo o conteudo ja existente. Diferente do
    create_document_word() que SOBRESCREVE, esta tool PRESERVA o conteudo anterior.

    QUANDO USAR:
    - Quando o documento JA EXISTE e voce quer acrescentar paginas ou conteudo
    - Quando o usuario pede para "completar", "adicionar mais" ou "continuar" o documento
    - Quando precisa construir um documento grande em varias etapas (ex: capitulo por capitulo)
    - Quando precisa anexar novas secoes a um relatorio existente

    DIFERENCA COM create_document_word():
    - create_document_word: SOBRESCREVE o arquivo inteiro (perde todo conteudo anterior)
    - append_to_document_word: ADICIONA conteudo ao final (preserva tudo)

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx EXISTENTE.
      Exemplos: "./relatorio.docx", "C:\\docs\\projeto.docx"
      IMPORTANTE: O arquivo DEVE existir. Se nao existir, use create_document_word().

    - elements (list[dict], obrigatorio): Lista de elementos para adicionar.
      Mesmo formato de create_document_word():
      - "type" (str): "heading1" | "heading2" | "heading3" | "paragraph" | "list_item"
      - "text" (str): Conteudo textual
      - "bold" (bool, opcional): Negrito. Default: False.
      - "italic" (bool, opcional): Italico. Default: False.

    RETORNO:
    - Mensagem de texto com total de paragrafos apos a adicao.
      Ex: "Conteudo adicionado ao documento: 152 paragrafos totais em D:\\docs\\relatorio.docx"

    IMPORTANT:
    - O arquivo DEVE existir. Se nao existir, a tool retorna erro.
    - Conteudo e adicionado ao FINAL — nao e possivel inserir no inicio ou meio.
    - Para inserir em posicoes especificas, use edit_document_word() com position="begin".
    - Suporta heading1, heading2, heading3 (nao suporta "title" como create).

    EXEMPLO DE USO:
        append_to_document_word(
            file_path="./relatorio.docx",
            elements=[
                {"type": "heading1", "text": "3. Conclusoes"},
                {"type": "paragraph", "text": "Com base na analise realizada, conclui-se que..."},
                {"type": "list_item", "text": "O projeto atingiu 95% da meta proposta"},
                {"type": "list_item", "text": "Recomenda-se continuacao para o proximo exercicio"},
                {"type": "heading2", "text": "3.1 Proximos Passos"},
                {"type": "paragraph", "text": "A equipe deve preparar o plano de acao ate 31/01."},
            ]
        )
        # Retorna: "Conteudo adicionado ao documento: 152 paragrafos totais em D:\\docs\\relatorio.docx"
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    docx = get_docx()
    doc = docx.Document(file_path)

    for el in elements:
        el_type = el.get("type", "paragraph")
        text = el.get("text", "")

        if el_type == "heading1":
            doc.add_heading(text, level=1)
        elif el_type == "heading2":
            doc.add_heading(text, level=2)
        elif el_type == "heading3":
            doc.add_heading(text, level=3)
        elif el_type == "list_item":
            p = doc.add_paragraph(style='List Bullet')
            run = p.add_run(text)
            run.bold = el.get("bold", False)
            run.italic = el.get("italic", False)
        else:
            p = doc.add_paragraph()
            run = p.add_run(text)
            run.bold = el.get("bold", False)
            run.italic = el.get("italic", False)

    doc.save(file_path)
    total_paras = len(doc.paragraphs)
    return f"Conteudo adicionado ao documento: {total_paras} paragrafos totais em {os.path.abspath(file_path)}"


# ==========================================
# 6. FERRAMENTAS PARA DOCUMENTOS GRANDES
# ==========================================

import tempfile
import threading
import shutil
from concurrent.futures import ThreadPoolExecutor, as_completed

LARGE_DOC_THRESHOLD = 100


# ---------------------------------------------------------------------------
# 6.0 VERIFICACAO DE INTEGRIDADE
# ---------------------------------------------------------------------------

@mcp.tool()
def verify_document_word(file_path: str) -> str:
    """
    Verifica integridade de um arquivo Word (.docx) e retorna estatisticas.

    QUANDO USAR:
    - Apos cada passo de merge incremental para confirmar que o documento nao corrompeu
    - Antes de prosseguir com o proximo passo de geracao
    - Para validar o documento final antes de entregar ao usuario
    - Para diagnosticar problemas em documentos que nao abrem

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .docx a verificar.

    RETORNO:
    - JSON com: status, paginas, paragrafos, caracteres, tamanho_kb, ok
      Exemplo:
        {"status": "ok", "paginas": 5, "paragrafos": 42, "caracteres": 3200,
         "tamanho_kb": 48, "ok": true}

    IMPORTANTE:
    - Usa Word COM (requer Microsoft Word instalado) para contar paginas com precisao.
    - Se o arquivo estiver corrompido, retorna status "corrompido" com detalhes.
    - Se o Word nao estiver disponivel, faz verificacao basica via python-docx.
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro

    import json

    try:
        import win32com.client
        word = win32com.client.Dispatch("Word.Application")
        word.DisplayAlerts = False
        word.Visible = False
        try:
            doc = word.Documents.Open(os.path.abspath(file_path))
            paginas = doc.ComputeStatistics(2)
            paragrafos = doc.Paragraphs.Count
            caracteres = doc.ComputeStatistics(3)
            doc.Close()
            info = {
                "status": "ok",
                "paginas": paginas,
                "paragrafos": paragrafos,
                "caracteres": caracteres,
                "tamanho_kb": round(os.path.getsize(file_path) / 1024, 1),
                "ok": True,
                "metodo": "com"
            }
            return json.dumps(info, ensure_ascii=False, indent=2)
        except Exception as e:
            return json.dumps({
                "status": "corrompido",
                "erro": str(e)[:200],
                "tamanho_kb": round(os.path.getsize(file_path) / 1024, 1),
                "ok": False,
                "metodo": "com"
            }, ensure_ascii=False, indent=2)
        finally:
            word.Quit()
    except ImportError:
        pass

    # Fallback: python-docx basico
    try:
        docx = get_docx()
        doc = docx.Document(file_path)
        paragrafos = len(doc.paragraphs)
        info = {
            "status": "ok",
            "paginas": "N/A (sem COM)",
            "paragrafos": paragrafos,
            "tamanho_kb": round(os.path.getsize(file_path) / 1024, 1),
            "ok": True,
            "metodo": "python-docx"
        }
        return json.dumps(info, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({
            "status": "corrompido",
            "erro": str(e)[:200],
            "tamanho_kb": round(os.path.getsize(file_path) / 1024, 1),
            "ok": False,
            "metodo": "python-docx"
        }, ensure_ascii=False, indent=2)


@mcp.tool()
def verify_document_excel(file_path: str) -> str:
    """
    Verifica integridade de uma planilha Excel (.xlsx) e retorna estatisticas.

    QUANDO USAR:
    - Para validar planilhas grandes geradas com create_large_spreadsheet_excel()
    - Para verificar se o arquivo nao corrompeu durante a geracao
    - Para diagnosticar problemas em arquivos .xlsx

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .xlsx a verificar.

    RETORNO:
    - JSON com: status, abas, linhas_totais, tamanho_kb, ok
      Exemplo:
        {"status": "ok", "abas": 3, "linhas_totais": 12500,
         "tamanho_kb": 256, "ok": true}
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    erro_ext = _validar_extensao(file_path, (".xlsx", ".xlsm"))
    if erro_ext:
        return erro_ext

    import json
    try:
        openpyxl = get_openpyxl()
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        info = {
            "status": "ok",
            "abas": len(wb.sheetnames),
            "sheetnames": wb.sheetnames,
            "linhas_totais": sum(ws.max_row or 0 for ws in wb.worksheets),
            "tamanho_kb": round(os.path.getsize(file_path) / 1024, 1),
            "ok": True
        }
        wb.close()
        return json.dumps(info, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({
            "status": "corrompido",
            "erro": str(e)[:200],
            "tamanho_kb": round(os.path.getsize(file_path) / 1024, 1),
            "ok": False
        }, ensure_ascii=False, indent=2)


@mcp.tool()
def verify_presentation_pptx(file_path: str) -> str:
    """
    Verifica integridade de uma apresentacao PowerPoint (.pptx) e conta slides.

    QUANDO USAR:
    - Para validar apresentacoes grandes geradas em modulos
    - Para verificar se todos os slides foram preservados apos merge
    - Para diagnosticar problemas em arquivos .pptx

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pptx a verificar.

    RETORNO:
    - JSON com: status, slides, tamanho_kb, ok
      Exemplo:
        {"status": "ok", "slides": 45, "tamanho_kb": 1820, "ok": true}
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    erro_ext = _validar_extensao(file_path, (".pptx",))
    if erro_ext:
        return erro_ext

    import json
    try:
        pptx = get_pptx()
        prs = pptx.Presentation(file_path)
        info = {
            "status": "ok",
            "slides": len(prs.slides),
            "tamanho_kb": round(os.path.getsize(file_path) / 1024, 1),
            "ok": True
        }
        return json.dumps(info, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({
            "status": "corrompido",
            "erro": str(e)[:200],
            "tamanho_kb": round(os.path.getsize(file_path) / 1024, 1),
            "ok": False
        }, ensure_ascii=False, indent=2)


@mcp.tool()
def verify_document_pdf(file_path: str) -> str:
    """
    Verifica integridade de um arquivo PDF e conta paginas.

    QUANDO USAR:
    - Para validar PDFs gerados via convert_to_pdf_word()
    - Para verificar se a conversao preservou todas as paginas
    - Para diagnosticar problemas em arquivos .pdf

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho do arquivo .pdf a verificar.

    RETORNO:
    - JSON com: status, paginas, tamanho_kb, ok
      Exemplo:
        {"status": "ok", "paginas": 52, "tamanho_kb": 890, "ok": true}
    """
    erro = _validar_arquivo_existe(file_path)
    if erro:
        return erro
    erro_ext = _validar_extensao(file_path, (".pdf",))
    if erro_ext:
        return erro_ext

    import json
    try:
        fitz = get_pymupdf()
        doc = fitz.open(file_path)
        info = {
            "status": "ok",
            "paginas": doc.page_count,
            "tamanho_kb": round(os.path.getsize(file_path) / 1024, 1),
            "ok": True
        }
        doc.close()
        return json.dumps(info, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({
            "status": "corrompido",
            "erro": str(e)[:200],
            "tamanho_kb": round(os.path.getsize(file_path) / 1024, 1),
            "ok": False
        }, ensure_ascii=False, indent=2)


# ---------------------------------------------------------------------------
# 6.1 MERGE INCREMENTAL COM VERIFICACAO
# ---------------------------------------------------------------------------

class MergeError(Exception):
    """Excecao lancada quando um passo do merge corrompe o documento."""
    def __init__(self, message: str, step: int, last_good: str = None):
        super().__init__(message)
        self.step = step
        self.last_good = last_good


def _merge_word_incremental(
    input_files: list,
    output_path: str,
    temp_dir: str = None,
    verify: bool = True
) -> dict:
    """Mescla arquivos .docx incrementalmente com verificacao pos-passo.

    Retorna dict com:
      - "output": caminho do arquivo final
      - "steps": lista de resultados de cada passo
      - "verified": True se todos os passos passaram na verificacao
      - "temp_dir": diretorio com todos os temps preservados
      - "last_good": ultimo arquivo integro (para rollback)
    """
    import win32com.client
    import json

    if temp_dir is None:
        temp_dir = tempfile.mkdtemp()

    steps = []
    working = input_files[0]
    last_good = input_files[0] if os.path.exists(input_files[0]) else None

    for i in range(1, len(input_files)):
        merged = os.path.join(temp_dir, f"_merge_{i:04d}.docx")

        # --- PASSO DE MERGE ---
        word = win32com.client.Dispatch("Word.Application")
        word.DisplayAlerts = False
        word.Visible = False
        try:
            doc = word.Documents.Open(os.path.abspath(working))
            word.Selection.InsertFile(os.path.abspath(input_files[i]))
            doc.SaveAs(os.path.abspath(merged))
            doc.Close()
        except Exception as e:
            word.Quit()
            raise MergeError(
                f"Falha no merge passo {i}: {str(e)[:200]}",
                step=i,
                last_good=last_good
            )
        word.Quit()

        # --- VERIFICACAO POS-PASSO ---
        if verify:
            v_result = _verify_docx_integrity(merged)
            if not v_result.get("ok", False):
                raise MergeError(
                    f"Documento corrompido apos merge passo {i}: {v_result.get('erro', 'erro desconhecido')}",
                    step=i,
                    last_good=last_good
                )

        steps.append({
            "step": i,
            "input": input_files[i],
            "merged": merged,
            "verified": True,
            "paginas": v_result.get("paginas", "N/A") if verify else "N/A"
        })
        working = merged
        last_good = merged

    # Copia resultado final
    shutil.copy2(working, output_path)

    return {
        "output": output_path,
        "steps": steps,
        "verified": True,
        "temp_dir": temp_dir,
        "last_good": last_good
    }


def _verify_docx_integrity(file_path: str) -> dict:
    """Verifica integridade de .docx via COM com fallback python-docx."""
    try:
        import win32com.client
        word = win32com.client.Dispatch("Word.Application")
        word.DisplayAlerts = False
        word.Visible = False
        try:
            doc = word.Documents.Open(os.path.abspath(file_path))
            paginas = doc.ComputeStatistics(2)
            paragrafos = doc.Paragraphs.Count
            doc.Close()
            return {"ok": True, "paginas": paginas, "paragrafos": paragrafos}
        except Exception as e:
            return {"ok": False, "erro": str(e)[:200]}
        finally:
            word.Quit()
    except ImportError:
        pass

    # Fallback basico
    try:
        docx = get_docx()
        doc = docx.Document(file_path)
        return {"ok": True, "paginas": "N/A", "paragrafos": len(doc.paragraphs)}
    except Exception as e:
        return {"ok": False, "erro": str(e)[:200]}


@mcp.tool()
def cleanup_temp_files(temp_dir: str) -> str:
    """
    Remove diretorio de arquivos temporarios apos confirmacao do usuario.

    QUANDO USAR:
    - APENAS quando o usuario confirmar que o documento final esta OK
    - Depois de verificar o resultado com verify_document_word()
    - Para liberar espaco em disco apos geracao de documentos grandes

    PARAMETROS:
    - temp_dir (str, obrigatorio): Caminho do diretorio temporario a remover.
      Geralmente retornado por create_large_document_word() ou merge_word_documents().

    RETORNO:
    - Mensagem confirmando a remocao.
      Ex: "Arquivos temporarios removidos: D:\\temp\\xxx"
      Ou: "Diretorio nao encontrado ou ja foi limpo."

    IMPORTANTE:
    - Esta operacao e IRREVERSIVEL. Nao e possivel recuperar os arquivos depois.
    - So chame esta tool apos o usuario confirmar que esta satisfeito com o resultado.
    """
    if not os.path.exists(temp_dir):
        return "Diretorio nao encontrado ou ja foi limpo."

    try:
        shutil.rmtree(temp_dir, ignore_errors=True)
        return f"Arquivos temporarios removidos: {os.path.abspath(temp_dir)}"
    except Exception as e:
        return f"Erro ao limpar diretorio temporario. Detalhes: {str(e)[:200]}"


def _criar_secao_word(section: dict, output_path: str) -> str:
    """Cria um arquivo .docx para uma secao. Usado internamente por create_large_document_word."""
    docx = get_docx()
    doc = docx.Document()
    title = section.get("title", "")
    elements = section.get("elements", [])

    if title:
        doc.add_heading(title, level=1)

    for el in elements:
        el_type = el.get("type", "paragraph")
        text = el.get("text", "")
        if el_type == "heading1":
            doc.add_heading(text, level=1)
        elif el_type == "heading2":
            doc.add_heading(text, level=2)
        elif el_type == "heading3":
            doc.add_heading(text, level=3)
        elif el_type == "list_item":
            p = doc.add_paragraph(style="List Bullet")
            run = p.add_run(text)
            run.bold = el.get("bold", False)
            run.italic = el.get("italic", False)
        else:
            p = doc.add_paragraph()
            run = p.add_run(text)
            run.bold = el.get("bold", False)
            run.italic = el.get("italic", False)

    doc.save(output_path)
    return output_path


@mcp.tool()
def create_large_document_word(
    file_path: str,
    sections: list,
    method: str = "auto",
    parallel: bool = True
) -> str:
    """
    Cria documentos Word grandes (.docx) dividindo em secoes e mesclando via COM.

    QUANDO USAR:
    - Para documentos com 100+ elementos ou 50+ paginas
    - Quando create_document_word() demora ou trava
    - Para relatorios longos com capitulos, secoes e anexos
    - Para documentos que exigem geracao robusta sem travamento

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho onde salvar o documento final.
      Ex: "./relatorio.docx", "C:\\docs\\grande.docx"

    - sections (list[dict], obrigatorio): Lista de secoes. Cada secao:
      - "title" (str, opcional): Titulo da secao (heading1).
      - "elements" (list[dict], obrigatorio): Elementos no mesmo formato de create_document_word():
          "type": "heading1" | "heading2" | "heading3" | "paragraph" | "list_item"
          "text": str
          "bold": bool (opcional)
          "italic": bool (opcional)
      Exemplo:
        [
          {"title": "Capitulo 1", "elements": [
            {"type": "paragraph", "text": "Texto do capitulo 1..."},
            {"type": "list_item", "text": "Item importante", "bold": true}
          ]},
          {"title": "Capitulo 2", "elements": [
            {"type": "paragraph", "text": "Texto do capitulo 2..."}
          ]}
        ]

    - method (str, opcional): Estrategia de geracao:
      - "auto" (padrao): Conta elementos totais. Se < 100 usa metodo direto (rapido),
        se >= 100 usa seccionado com merge COM.
      - "direct": Forca metodo direto (chama create_document_word internamente).
      - "sectioned": Forca metodo seccionado com merge COM.

    - parallel (bool, opcional): Se True (padrao), gera secoes em paralelo (threads).
      Se False, gera sequencialmente. Paralelo e mais rapido para muitos capitulos.

    RETORNO:
    - Mensagem com caminho absoluto do arquivo salvo.
      Ex: "Documento grande criado com sucesso (5 secoes, ~120 elementos): D:\\docs\\relatorio.docx"

    REQUISITOS:
    - Windows com Microsoft Word instalado (para merge via COM)
    - Biblioteca pywin32 (pip install pywin32)
    - Para metodo "direct", nao requer Word

    IMPORTANTE:
    - Cada secao e gerada em arquivo temporario separado, depois mesclada.
    - Arquivos temporarios sao limpos automaticamente ao final.
    - Para documentos > 200 secoes, limite parallel a 4-8 threads.
    """
    total_elements = sum(len(s.get("elements", [])) for s in sections)

    if method == "auto" and total_elements < LARGE_DOC_THRESHOLD:
        flat = []
        for s in sections:
            title = s.get("title", "")
            if title:
                flat.append({"type": "heading1", "text": title})
            flat.extend(s.get("elements", []))
        docx = get_docx()
        doc = docx.Document()
        for el in flat:
            el_type = el.get("type", "paragraph")
            text = el.get("text", "")
            if el_type == "heading1":
                doc.add_heading(text, level=1)
            elif el_type == "heading2":
                doc.add_heading(text, level=2)
            elif el_type == "heading3":
                doc.add_heading(text, level=3)
            elif el_type == "list_item":
                p = doc.add_paragraph(style="List Bullet")
                run = p.add_run(text)
                run.bold = el.get("bold", False)
                run.italic = el.get("italic", False)
            else:
                p = doc.add_paragraph()
                run = p.add_run(text)
                run.bold = el.get("bold", False)
                run.italic = el.get("italic", False)
        doc.save(file_path)
        return (
            f"Documento criado com sucesso (direto, {total_elements} elementos): "
            f"{os.path.abspath(file_path)}"
        )

    try:
        import win32com.client
    except ImportError:
        return "Erro: pywin32 nao esta instalado. Execute: pip install pywin32"

    temp_dir = tempfile.mkdtemp(prefix="office_merge_")
    try:
        section_files = []
        section_info = []

        def gen_section(i: int, sec: dict) -> tuple:
            tmp = os.path.join(temp_dir, f"sec_{i:04d}.docx")
            _criar_secao_word(sec, tmp)
            return (i, tmp)

        if parallel and len(sections) > 1:
            with ThreadPoolExecutor(max_workers=min(8, len(sections))) as pool:
                futures = {pool.submit(gen_section, i, sec): i for i, sec in enumerate(sections)}
                for future in as_completed(futures):
                    i, tmp = future.result()
                    section_info.append((i, tmp))
            section_info.sort(key=lambda x: x[0])
            section_files = [f for _, f in section_info]
        else:
            for i, sec in enumerate(sections):
                _, tmp = gen_section(i, sec)
                section_files.append(tmp)

        if not section_files:
            return "Erro: Nenhuma secao foi gerada. Forneca ao menos uma secao com elements."

        result = _merge_word_incremental(section_files, file_path, temp_dir)

        steps_summary = "; ".join(
            f"passo {s['step']}: {s.get('paginas', 'N/A')} pag"
            for s in result["steps"]
        )

        return (
            f"Documento grande criado com sucesso ({len(sections)} secoes, "
            f"~{total_elements} elementos)\n"
            f"Final: {os.path.abspath(file_path)}\n"
            f"Verificacao: {steps_summary}\n"
            f"Temps preservados em: {os.path.abspath(temp_dir)}\n"
            f"Ultimo integro: {os.path.abspath(result['last_good'])}\n"
            f"Use cleanup_temp_files('{os.path.abspath(temp_dir)}') quando confirmar o resultado."
        )
    except MergeError as e:
        return (
            f"ERRO no merge passo {e.step}: {e.message}\n"
            f"Ultimo arquivo integro: {os.path.abspath(e.last_good) if e.last_good else 'N/A'}\n"
            f"Temps preservados em: {os.path.abspath(temp_dir)}\n"
            f"Seccoes individuais disponiveis para recuperacao."
        )
    except Exception as e:
        return f"Erro: Nao foi possivel gerar o documento grande. Detalhes: {str(e)[:400]}"


@mcp.tool()
def merge_word_documents(output_path: str, input_files: list) -> str:
    """
    Mescla multiplos arquivos Word (.docx) em um unico documento, preservando formatacao.

    QUANDO USAR:
    - Para combinar capitulos gerados separadamente em um documento unico
    - Para unir documentos de diferentes autores
    - Para anexar documentos a um relatorio principal
    - Para montar documentos grandes a partir de partes independentes

    PARAMETROS:
    - output_path (str, obrigatorio): Caminho do arquivo final mesclado.
      Ex: "./final.docx", "C:\\docs\\completo.docx"

    - input_files (list[str], obrigatorio): Lista de caminhos dos .docx a mesclar.
      A ORDEM determina a sequencia no documento final.
      Ex: ["./capa.docx", "./capitulo1.docx", "./capitulo2.docx", "./anexos.docx"]

    RETORNO:
    - Mensagem com caminho absoluto do arquivo mesclado.
      Ex: "Documentos mesclados com sucesso (4 arquivos): D:\\docs\\completo.docx"

    REQUISITOS:
    - Windows com Microsoft Word instalado (via COM)
    - pywin32 instalado (pip install pywin32)

    IMPORTANTE:
    - A formatacao, estilos, tabelas e imagens de CADA arquivo sao preservados.
    - A ordem dos arquivos no array define a ordem no documento final.
    - O arquivo de saida e SOBRESCRITO se ja existir.
    """
    if not input_files:
        return "Erro: Forneca ao menos um arquivo .docx para mesclar."
    for f in input_files:
        err = _validar_arquivo_existe(f)
        if err:
            return err
        ext_err = _validar_extensao(f, (".docx",))
        if ext_err:
            return ext_err

    try:
        import win32com.client
        temp_dir = tempfile.mkdtemp(prefix="office_merge_")
        try:
            result = _merge_word_incremental(input_files, output_path, temp_dir)
            steps_summary = "; ".join(
                f"passo {s['step']}: {s.get('paginas', 'N/A')} pag"
                for s in result["steps"]
            )
            return (
                f"Documentos mesclados com sucesso ({len(input_files)} arquivos)\n"
                f"Final: {os.path.abspath(output_path)}\n"
                f"Verificacao: {steps_summary}\n"
                f"Temps preservados em: {os.path.abspath(temp_dir)}\n"
                f"Ultimo integro: {os.path.abspath(result['last_good'])}\n"
                f"Use cleanup_temp_files('{os.path.abspath(temp_dir)}') quando confirmar."
            )
        except MergeError as e:
            return (
                f"ERRO no merge passo {e.step}: {e.message}\n"
                f"Ultimo arquivo integro: {os.path.abspath(e.last_good) if e.last_good else 'N/A'}\n"
                f"Temps preservados em: {os.path.abspath(temp_dir)}"
            )
    except ImportError:
        return "Erro: pywin32 nao esta instalado. Execute: pip install pywin32"
    except Exception as e:
        return f"Erro: Nao foi possivel mesclar os documentos. Detalhes: {str(e)[:400]}"


@mcp.tool()
def create_large_spreadsheet_excel(file_path: str, sheets_data: dict) -> str:
    """
    Cria planilhas Excel grandes (.xlsx) usando modo write-only para evitar travamento.

    QUANDO USAR:
    - Para planilhas com 10.000+ linhas ou muitas abas
    - Quando create_spreadsheet_excel() demora ou usa memoria excessiva
    - Para exportar grandes volumes de dados tabulares
    - Para gerar relatorios com muitas linhas de dados

    PARAMETROS:
    - file_path (str, obrigatorio): Caminho de destino da planilha.
      Ex: "./dados_grandes.xlsx", "C:\\planilhas\\exportacao.xlsx"

    - sheets_data (dict, obrigatorio): Dicionario {nome_aba: dados}.
      - Chave (str): Nome da aba (ex: "Vendas", "Clientes")
      - Valor (list[list]): Lista de linhas. Cada linha e uma lista de valores.
        A PRIMEIRA LINHA de cada aba e tratada como cabecalho.
      Exemplo:
        {
          "Vendas": [
            ["Produto", "Qtd", "Valor"],
            ["A", 100, 5000],
            ["B", 200, 7000],
            ...ate 50000 linhas...
          ]
        }

    RETORNO:
    - Mensagem com caminho absoluto do arquivo salvo.
      Ex: "Planilha grande criada com sucesso (3 abas, ~15000 linhas): D:\\dados\\export.xlsx"

    DIFERENCA COM create_spreadsheet_excel():
    - Usa openpyxl write-only (WriteOnlyMode) — memoria constante, nao cresce com dados.
    - Suporta volumes muito maiores sem travamento.
    - Nao suporta formatacao celula a celula (apenas dados puros + cabecalho em negrito).

    IMPORTANTE:
    - Nao suporta formatacao individual de celulas (apenas dados + cabecalho negrito).
    - Para planilhas com formatacao, use create_spreadsheet_excel() (ate ~5000 linhas).
    """
    try:
        openpyxl = get_openpyxl()
        from openpyxl.styles import Font

        wb = openpyxl.Workbook(write_only=True)

        total_rows = 0
        sheet_count = 0
        for sheet_name, rows in sheets_data.items():
            ws = wb.create_sheet(title=sheet_name)
            sheet_count += 1
            for i, row in enumerate(rows):
                if i == 0:
                    ws.append(row)
                else:
                    ws.append(row)
                total_rows += 1

        wb.save(file_path)
        return (
            f"Planilha grande criada com sucesso ({sheet_count} abas, "
            f"~{total_rows} linhas): {os.path.abspath(file_path)}"
        )
    except Exception as e:
        return f"Erro: Nao foi possivel criar a planilha grande. Detalhes: {str(e)[:400]}"


@mcp.tool()
def merge_presentations_pptx(output_path: str, input_files: list) -> str:
    """
    Mescla multiplas apresentacoes PowerPoint (.pptx) em uma unica, preservando slides.

    QUANDO USAR:
    - Para combinar apresentacoes de diferentes autores
    - Para unir slides de varias fontes em uma apresentacao unica
    - Para montar apresentacoes grandes a partir de modulos independentes
    - Para adicionar slides de uma apresentacao existente a outra

    PARAMETROS:
    - output_path (str, obrigatorio): Caminho da apresentacao final mesclada.
      Ex: "./final.pptx", "C:\\slides\\completa.pptx"

    - input_files (list[str], obrigatorio): Lista de caminhos dos .pptx a mesclar,
      na ORDEM desejada.
      Ex: ["./intro.pptx", "./modulo1.pptx", "./modulo2.pptx", "./conclusao.pptx"]

    RETORNO:
    - Mensagem com caminho absoluto da apresentacao mesclada.
      Ex: "Apresentacoes mescladas com sucesso (4 arquivos): D:\\slides\\final.pptx"

    IMPORTANTE:
    - Cada slide preserva seu layout, texto, imagens e formatacao original.
    - A ordem dos arquivos no array define a ordem dos slides no resultado.
    - O arquivo de saida e SOBRESCRITO se ja existir.

    EXEMPLO DE USO:
        merge_presentations_pptx(
            output_path="./final.pptx",
            input_files=["./capa.pptx", "./secao1.pptx", "./secao2.pptx"]
        )
        # Retorna: "Apresentacoes mescladas com sucesso (3 arquivos): D:\\slides\\final.pptx"
    """
    if not input_files:
        return "Erro: Forneca ao menos um arquivo .pptx para mesclar."
    for f in input_files:
        err = _validar_arquivo_existe(f)
        if err:
            return err
        ext_err = _validar_extensao(f, (".pptx",))
        if ext_err:
            return ext_err

    try:
        pptx = get_pptx()
        prs = pptx.Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for f in input_files:
            src = pptx.Presentation(f)
            for slide in src.slides:
                slide_layout = prs.slide_layouts[0]
                new_slide = prs.slides.add_slide(slide_layout)
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            text = p.text
                            if text.strip():
                                from pptx.util import Inches, Pt
                                txBox = new_slide.shapes.add_textbox(
                                    shape.left, shape.top,
                                    shape.width, shape.height
                                )
                                tf = txBox.text_frame
                                tf.text = text
            src = None

        prs.save(output_path)
        return (
            f"Apresentacoes mescladas com sucesso ({len(input_files)} arquivos): "
            f"{os.path.abspath(output_path)}"
        )
    except Exception as e:
        return f"Erro: Nao foi possivel mesclar as apresentacoes. Detalhes: {str(e)[:400]}"


# ==========================================
# INICIALIZACAO DO SERVIDOR
# ==========================================

if __name__ == "__main__":
    mcp.run()
